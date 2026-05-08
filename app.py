"""
APSG (Staging Ground) Report — Unified Application
Combined: Daily Report | Excel Rejection | PPT Alignment | Bulk Bundle Filter
"""

import os, io, uuid, threading, time, json, copy, re, traceback, hashlib, logging
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime, date
from collections import OrderedDict
from functools import wraps

from flask import (Flask, request, jsonify, send_file,
                   session, redirect, url_for, render_template_string)
from werkzeug.utils import secure_filename

# ── Python-PPTX (used by PPT Rejection + Photo Merge) ──────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# ── Photo Merge engine ──────────────────────────────────────────────────────
import ppt_merger

# ── Daily Report engines ──────────────────────────────────────────────────────
from report_engine import (
    load_and_validate_file, generate_report, get_date_range,
    filter_preview, validate_and_flag, validate_net_weights, validate_etoken,
)
from wb_engine import (
    load_wb_file, wb_get_date_range, wb_filter_by_date,
    wb_apply_unified_logic, wb_build_pivot, wb_pivot_to_excel,
    wb_pivot_copy_text, wb_find_incomplete_rows, wb_apply_row_decisions,
    wb_net_weight_validation, wb_validate_etoken_match,
    _normalize_accepted as _wb_norm_acc,
)
from rectification_report import (
    build_rr_docx, fetch_row_by_token, fetch_row_from_series,
    apply_user_updates, fetch_and_generate, generate_table_image,
    _parse_dt, _MONTH_B,
    REASON_A, REASON_B, REASON_C,
)
import pandas as pd
import math

# ── DB / state ──────────────────────────────────────────────────────────────
import sqlite3

# ── JSON-safe numeric helpers ────────────────────────────────────────────────
def safe_float(v, default=0.0):
    """Convert v to float, returning default for NaN/None/inf/NA."""
    try:
        result = float(v)
        if math.isnan(result) or math.isinf(result):
            return default
        return result
    except (TypeError, ValueError):
        return default

def safe_int(v, default=0):
    """Convert v to int, returning default for NaN/None/NA."""
    try:
        f = float(v)
        if math.isnan(f) or math.isinf(f):
            return default
        return int(f)
    except (TypeError, ValueError):
        return default

def safe_col_sum(df, col, default=0.0):
    """Sum a DataFrame column safely, returning default if col missing or all NaN."""
    if col not in df.columns:
        return default
    return safe_float(df[col].dropna().sum(), default)

app = Flask(__name__, static_folder='static', static_url_path='/static')
app.secret_key = os.environ.get("SECRET_KEY", "apsg-report-secret-2024")

# ── NaN-safe JSON serialization ──────────────────────────────────────────────
# Flask's default encoder passes NaN to json.dumps which produces invalid JSON.
# Override to replace NaN/Inf with null so browsers never see a parse error.
import json as _json
class _NanSafeProvider(app.json_provider_class):
    def dumps(self, obj, **kw):
        return _json.dumps(obj, **kw, allow_nan=False,
                           default=lambda o: None if (isinstance(o, float) and (math.isnan(o) or math.isinf(o))) else str(o))
    def loads(self, s, **kw):
        return _json.loads(s, **kw)
app.json_provider_class = _NanSafeProvider
app.json = _NanSafeProvider(app)

_BASE = os.path.dirname(os.path.abspath(__file__))
# Fly.io: use /app/data (persistent volume mount) if available, else local
_DATA = os.environ.get("DATA_DIR", os.path.join(_BASE, "data") if os.path.isdir(os.path.join(_BASE, "data")) else _BASE)
DB_PATH = os.environ.get("DB_PATH", os.path.join(_DATA, "apsg_report.db"))
UPLOAD_DIR = os.environ.get("UPLOAD_DIR", os.path.join(_DATA, "uploads"))
OUTPUT_DIR = os.environ.get("OUTPUT_DIR", os.path.join(_DATA, "outputs"))
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ── PPT Rejection in-memory DB ──────────────────────────────────────────────
PPT_DB: dict = {"records": [], "slide_map": {}, "presentations": []}
jobs: dict = {}
jobs_lock = threading.Lock()

# ═══════════════════════════════════════════════════════════════════════════════
#  DATABASE SETUP
# ═══════════════════════════════════════════════════════════════════════════════

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    """Create tables and seed default admin. Safe to call multiple times."""
    conn = get_db()
    c = conn.cursor()
    c.execute("""CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        name TEXT NOT NULL,
        password_hash TEXT NOT NULL,
        plaintext_pw TEXT DEFAULT '',
        role TEXT DEFAULT 'user',
        created_at TEXT DEFAULT (datetime('now'))
    )""")
    # Add plaintext_pw column if upgrading from older schema
    try:
        c.execute("ALTER TABLE users ADD COLUMN plaintext_pw TEXT DEFAULT ''")
    except Exception:
        pass  # Column already exists
    c.execute("""CREATE TABLE IF NOT EXISTS activity_log (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT NOT NULL,
        action TEXT NOT NULL,
        detail TEXT,
        ip TEXT,
        ts TEXT DEFAULT (datetime('now'))
    )""")
    # Always ensure admin exists (re-seeds after Render restarts)
    c.execute("""INSERT OR IGNORE INTO users (username,name,password_hash,role)
                 VALUES (?,?,?,?)""",
              ("admin", "Administrator", hash_pw("Admin@1234"), "admin"))
    # Also ensure karthi account exists
    c.execute("""INSERT OR IGNORE INTO users (username,name,password_hash,role)
                 VALUES (?,?,?,?)""",
              ("karthi", "Karthi", hash_pw("karthi123"), "admin"))
    conn.commit()
    conn.close()

# (DB init is called after all functions are defined — see bottom of file)

def hash_pw(pw: str) -> str:
    return hashlib.sha256(pw.encode()).hexdigest()

def verify_pw(pw: str, h: str) -> bool:
    return hash_pw(pw) == h

@app.before_request
def ensure_db_and_auth():
    """1. Re-init DB if missing. 2. Enforce login on every route."""
    # DB check
    try:
        conn = get_db()
        conn.execute("SELECT 1 FROM users LIMIT 1")
        conn.close()
    except Exception:
        try:
            init_db()
        except Exception as e:
            print(f"DB re-init error: {e}")
    # Auth enforcement — public endpoints that don't need login
    public = {'/login', '/register', '/api/health', '/static'}
    path = request.path
    if any(path.startswith(p) for p in public):
        return  # allow through
    if 'username' not in session:
        if request.is_json or path.startswith('/api/'):
            return jsonify({'error': 'Not authenticated', 'redirect': '/login'}), 401
        return redirect(url_for('login_page'))

ADMIN_EMAIL = os.environ.get("ADMIN_EMAIL", "")
SMTP_HOST   = os.environ.get("SMTP_HOST", "smtp.gmail.com")
SMTP_PORT   = int(os.environ.get("SMTP_PORT", "587"))
SMTP_USER   = os.environ.get("SMTP_USER", "")
SMTP_PASS   = os.environ.get("SMTP_PASS", "")

_notify_lock = threading.Lock()

def _send_email_notification(username: str, action: str, detail: str, ip: str, ts: str):
    """Send admin email notification in background thread. Silent if not configured."""
    if not all([ADMIN_EMAIL, SMTP_USER, SMTP_PASS]):
        return
    def _send():
        try:
            subject = f"[APSG] Activity: {action} by {username}"
            body = (
                f"<h3>APSG System Activity Notification</h3>"
                f"<table style='border-collapse:collapse;font-family:monospace;font-size:14px;'>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>User</td><td>{username}</td></tr>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>Action</td><td>{action}</td></tr>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>Detail</td><td>{detail or chr(8212)}</td></tr>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>IP</td><td>{ip}</td></tr>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>Timestamp</td><td>{ts}</td></tr>"
                f"</table>"
                f"<p style='color:#888;font-size:12px;'>APSG Report System &middot; Developed by Karthi</p>"
            )
            msg = MIMEMultipart("alternative")
            msg["Subject"] = subject
            msg["From"]    = SMTP_USER
            msg["To"]      = ADMIN_EMAIL
            msg.attach(MIMEText(body, "html"))
            with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=10) as srv:
                srv.ehlo(); srv.starttls(); srv.login(SMTP_USER, SMTP_PASS)
                srv.sendmail(SMTP_USER, ADMIN_EMAIL, msg.as_string())
        except Exception as e:
            print(f"[Email] Notification failed: {e}")
    with _notify_lock:
        threading.Thread(target=_send, daemon=True).start()


def send_activity_email(username: str, action: str, detail: str, ts: str, ip: str):
    """Send email notification for key actions (non-blocking, best-effort)."""
    import smtplib, os
    from email.mime.text import MIMEText
    notify_email = os.environ.get("NOTIFY_EMAIL", "")
    smtp_host    = os.environ.get("SMTP_HOST", "smtp.gmail.com")
    smtp_port    = int(os.environ.get("SMTP_PORT", "587"))
    smtp_user    = os.environ.get("SMTP_USER", "")
    smtp_pass    = os.environ.get("SMTP_PASS", "")
    if not (notify_email and smtp_user and smtp_pass):
        return  # Email not configured — skip silently
    try:
        subject = "[APSG] " + action + " - " + username
        sep = "-" * 40
        body = "\n".join([
            "APSG Activity Notification",
            sep,
            "User      : " + username,
            "Action    : " + action,
            "Detail    : " + (detail or "-"),
            "Timestamp : " + ts,
            "IP        : " + (ip or "unknown"),
        ])
        msg = MIMEText(body)
        msg["Subject"] = subject
        msg["From"] = smtp_user
        msg["To"] = notify_email
        with smtplib.SMTP(smtp_host, smtp_port) as srv:
            srv.ehlo(); srv.starttls(); srv.ehlo()
            srv.login(smtp_user, smtp_pass)
            srv.sendmail(smtp_user, [notify_email], msg.as_string())
    except Exception:
        pass  # Never let email errors break the main flow

_EMAIL_ACTIONS = {"LOGIN", "OPEN_APP", "DAILY_REPORT", "RECTIFICATION_REPORT",
                  "EXCEL_REJECTION", "PPT_GENERATE"}

def log_activity(username: str, action: str, detail: str = ""):
    ts = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S UTC")
    ip = ""
    try:
        ip = request.remote_addr or ""
    except Exception:
        pass
    # Send email notification for key actions (in background thread)
    if action in _EMAIL_ACTIONS:
        import threading
        threading.Thread(
            target=send_activity_email,
            args=(username, action, detail, ts, ip),
            daemon=True
        ).start()
    try:
        conn = get_db()
        conn.execute("INSERT INTO activity_log (username,action,detail,ip) VALUES (?,?,?,?)",
                     (username, action, detail, ip))
        conn.commit()
        conn.close()
    except Exception:
        pass
    _send_email_notification(username, action, detail, ip, ts)

# ─── auth helpers ────────────────────────────────────────────────────────────
def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if "username" not in session:
            return redirect(url_for("login_page"))
        return f(*args, **kwargs)
    return decorated

def admin_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if session.get("role") != "admin":
            return jsonify({"error": "Admin only"}), 403
        return f(*args, **kwargs)
    return decorated

# ═══════════════════════════════════════════════════════════════════════════════
#  PPT REJECTION LOGIC  (unchanged from original)
# ═══════════════════════════════════════════════════════════════════════════════

SLIDE_W  = 12192000; SLIDE_H  = 6858000
TBL_LEFT = 597159;   TBL_TOP  = 1189793; TBL_W = 11306716
COL_WIDTHS = [289286,1425361,739037,826948,1091559,928816,774237,2169617,626565,2435290]
ROW_H_HEADER = 681522; ROW_H_DATA = 734233
TITLE_LEFT=0; TITLE_TOP=288093; TITLE_W=12283440; TITLE_H=901700
HEADER_COLOR='00B050'; BORDER_W=6350

COMPANY_ABBREV = {
    'KKL':'KOH KOCK LEONG ENTERPRISE PTE LTD','KTC':'KTC Civil Engineering & Construction Pte Ltd',
    'UNC':'Unity Contractors Pte Ltd','CCE':'Chan & Chan Engineering Pte Ltd',
    'CLM':'Chuan Lim Construction Pte Ltd','HHT':'Hong Aik Engineering Pte Ltd',
    'JME':'JME E&C PTE LTD','KTE':'KOK TONG EARTHWORKS & ENGINEERING PTE LTD',
    'QQC':'QUEK & QUEK CIVIL ENGINEERING PTE LTD','BKH':'Backho (S) Pte Ltd',
    'OTP':'OKT TRANSPORT PTE LTD','RCE':'RECLAIMS ENTERPRISE PTE LTD',
    'SIN':'Sin Heng Transport Pte Ltd','HSE':'Hanshika Engineering & Construction Pte Ltd',
    'JFT':'JIN FENG TRANSPORT PTE LTD','CJC':'Chye Joo Construction Pte. Ltd.',
    'GGC':'Guan Gi Construction Pte Ltd','HTC':'HUATIONG CONTRACTOR PTE LTD',
    'KHT':'KENG HO TRADING AND TRANSPORT PTE LTD','MSH':'Megastone Holdings Pte Ltd',
    'MTC':'Metrocon Pte Ltd','TGT':'Tengah Transportation & Construction Pte Ltd',
    'SAS':'SASAN CONSTRUCTION PTE LTD','SHC':'SIN HUA CIVIL ENGINEERING & CONSTRUCTION PTE LTD',
    'WCF':'Wang Cheng Foundation Pte Ltd','YEC':'YONGSHENG E & C PTE LTD',
    'EEH':'EE HUP CONSTRUCTION PTE LTD','SLE':'Sam Lain Equipment Services Pte Ltd',
}

_BANNER_BLOB_CACHE = [None]

def _get_banner_blob():
    if _BANNER_BLOB_CACHE[0]: return _BANNER_BLOB_CACHE[0]
    for prs in PPT_DB.get('presentations',[]):
        for slide in prs.slides:
            for sh in slide.shapes:
                if sh.shape_type==13 and hasattr(sh,'width') and sh.width>9000000 and sh.top>5500000:
                    xml_s=etree.tostring(sh._element).decode()
                    embeds=re.findall(r'r:embed="(rId\d+)"',xml_s)
                    for rId in embeds:
                        rel=slide.part.rels.get(rId)
                        if rel and 'image' in rel.reltype:
                            _BANNER_BLOB_CACHE[0]=rel.target_part.blob
                            return _BANNER_BLOB_CACHE[0]
    return None

def _add_footer_banner(dst_slide):
    blob=_get_banner_blob()
    if not blob: return
    BANNER_H=557487; BANNER_Y=SLIDE_H-BANNER_H
    img_part,new_rId=dst_slide.part.get_or_add_image_part(io.BytesIO(blob))
    NS_P='xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    NS_A='xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    NS_R=f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
    pic_xml=(f'<p:pic {NS_P} {NS_A} {NS_R}><p:nvPicPr><p:cNvPr id="999" name="FooterBanner"/>'
             f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
             f'<p:blipFill><a:blip r:embed="{new_rId}" cstate="email"><a:extLst/></a:blip>'
             f'<a:stretch><a:fillRect/></a:stretch></p:blipFill>'
             f'<p:spPr><a:xfrm><a:off x="0" y="{BANNER_Y}"/><a:ext cx="{SLIDE_W}" cy="{BANNER_H}"/>'
             f'</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>')
    dst_slide.shapes._spTree.append(etree.fromstring(pic_xml.encode()))

def _borders():
    fill=('<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
          '<a:prstDash val="solid"/><a:round/>'
          '<a:headEnd type="none" w="med" len="med"/>'
          '<a:tailEnd type="none" w="med" len="med"/>')
    return ''.join(f'<a:{s} w="{BORDER_W}" cap="flat" cmpd="sng" algn="ctr">{fill}</a:{s}>'
                   for s in ['lnL','lnR','lnT','lnB'])

def parse_date_ppt(s):
    try: return datetime.strptime(s.strip(),'%d-%m-%Y')
    except: return None

def cell_xml(text,header,col):
    ns='xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    safe=(str(text or '')).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;').replace('"','&quot;')
    if header:
        algn,fa='ctr','b'; sz,bold,color,face=1200,1,HEADER_COLOR,'Calibri'
        mar='marL="8912" marR="8912" marT="8912" marB="0"'
    else:
        algn,fa='ctr','ctr'; sz,bold,color,face=1100,0,'000000','Calibri'
        mar='marL="9525" marR="9525" marT="9525" marB="0"'
    return (f'<a:tc {ns}><a:txBody><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:pPr algn="{algn}" fontAlgn="{fa}"/>'
            f'<a:r><a:rPr lang="en-SG" sz="{sz}" b="{bold}" i="0" u="none" '
            f'strike="noStrike" dirty="0">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:effectLst/><a:latin typeface="{face}"/>'
            f'</a:rPr><a:t>{safe}</a:t></a:r></a:p>'
            f'</a:txBody>'
            f'<a:tcPr {mar} anchor="ctr">{_borders()}<a:noFill/></a:tcPr>'
            f'</a:tc>')

def is_summary_slide(slide):
    return any(sh.has_table and len(sh.table.columns)==10 for sh in slide.shapes)

def is_detail_slide(slide):
    return any(sh.has_table and len(sh.table.columns)==8 for sh in slide.shapes)

def extract_from_prs(prs, prs_idx):
    records,slide_map=[],{}
    for si,slide in enumerate(prs.slides):
        if is_summary_slide(slide):
            for sh in slide.shapes:
                if sh.has_table and len(sh.table.columns)==10:
                    tbl=sh.table
                    for ri in range(1,len(tbl.rows)):
                        cells=[c.text.strip() for c in tbl.rows[ri].cells]
                        if len(cells)>=10 and cells[1]:
                            records.append({'sn':cells[0],'ticket_no':cells[1],'veh_no':cells[2],
                                'material':cells[3],'source_site':cells[4],'date':cells[5],
                                'time':cells[6],'e_token':cells[7],'accepted':cells[8],
                                'reject_reason':cells[9],'_prs':prs_idx})
        elif is_detail_slide(slide):
            for sh in slide.shapes:
                if sh.has_table and len(sh.table.columns)==8:
                    tbl=sh.table
                    if len(tbl.rows)>1:
                        tn=tbl.rows[1].cells[0].text.strip()
                        if tn: slide_map[tn]=(prs_idx,si)
                    break
    return records,slide_map

def add_summary_slide(out_prs,chunk,slide_num,sn_offset=1):
    slide=out_prs.slides.add_slide(out_prs.slide_layouts[6])
    tx=slide.shapes.add_textbox(Emu(TITLE_LEFT),Emu(TITLE_TOP),Emu(TITLE_W),Emu(TITLE_H))
    tf=tx.text_frame; tf.word_wrap=True
    r0=chunk[0] if chunk else {}
    sites=list(dict.fromkeys(r.get('source_site','') for r in chunk if r.get('source_site')))
    site_lbl=(sites[0] if len(sites)==1 else
              (f'{sites[0]}, {sites[1]}' if len(sites)==2 else f'{sites[0]} et al.'))
    p1=tf.paragraphs[0]; p1.alignment=PP_ALIGN.CENTER
    r1=p1.add_run(); r1.text='Rejected Loads Summary'
    r1.font.name='Calibri'; r1.font.size=Pt(28); r1.font.bold=True
    r1.font.underline=True; r1.font.color.rgb=RGBColor(0,0,0)
    p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
    r2=p2.add_run(); r2.text=site_lbl
    r2.font.name='Calibri'; r2.font.size=Pt(20); r2.font.bold=True
    r2.font.color.rgb=RGBColor(0x22,0x22,0x22)
    n_data=len(chunk)
    tbl_frm=slide.shapes.add_table(1+n_data,10,Emu(TBL_LEFT),Emu(TBL_TOP),
                                   Emu(TBL_W),Emu(ROW_H_HEADER+ROW_H_DATA*n_data))
    tbl=tbl_frm.table
    for ci,w in enumerate(COL_WIDTHS): tbl.columns[ci].width=w
    tbl.rows[0].height=ROW_H_HEADER
    for ri in range(1,1+n_data): tbl.rows[ri].height=ROW_H_DATA
    HEADERS=['S/N','Ticket No','Veh No','Material','Source Site','Date In','Time In','E-Token','Accepted','Reject Reason']
    for ci,h in enumerate(HEADERS):
        tc=tbl.rows[0].cells[ci]._tc
        tc.getparent().replace(tc,etree.fromstring(cell_xml(h,True,ci)))
    for ri,rec in enumerate(chunk):
        vals=[str(sn_offset+ri),rec.get('ticket_no',''),rec.get('veh_no',''),
              rec.get('material',''),rec.get('source_site',''),rec.get('date',''),
              rec.get('time',''),rec.get('e_token',''),rec.get('accepted','NO'),rec.get('reject_reason','')]
        for ci,v in enumerate(vals):
            tc=tbl.rows[ri+1].cells[ci]._tc
            tc.getparent().replace(tc,etree.fromstring(cell_xml(v,False,ci)))
    nb=slide.shapes.add_textbox(Emu(8610600),Emu(6356350),Emu(2743200),Emu(365125))
    np_=nb.text_frame.paragraphs[0]; np_.alignment=PP_ALIGN.RIGHT
    nr=np_.add_run(); nr.text=str(slide_num); nr.font.size=Pt(12); nr.font.color.rgb=RGBColor(0,0,0)
    _add_footer_banner(slide)
    return slide

def _replace_rids_single_pass(xml_str,rId_map):
    pattern=re.compile(r'(r:embed|r:link)="(rId\d+)"')
    def replacer(m): return f'{m.group(1)}="{rId_map.get(m.group(2),m.group(2))}"'
    return pattern.sub(replacer,xml_str)

def clone_slide(src_prs,src_idx,dst_prs,remove_slide_number=True,sn=None):
    src=src_prs.slides[src_idx]
    dst=dst_prs.slides.add_slide(dst_prs.slide_layouts[6])
    src_tree=src.shapes._spTree; dst_tree=dst.shapes._spTree
    for ch in list(dst_tree)[2:]: dst_tree.remove(ch)
    for ch in list(src_tree)[2:]: dst_tree.append(copy.deepcopy(ch))
    NS_P='http://schemas.openxmlformats.org/presentationml/2006/main'
    NS_A='http://schemas.openxmlformats.org/drawingml/2006/main'
    if remove_slide_number:
        for sp in list(dst_tree.iter(f'{{{NS_P}}}sp')):
            ph=sp.find(f'.//{{{NS_P}}}ph')
            if ph is not None and ph.get('type')=='sldNum':
                sp.getparent().remove(sp); break
    rId_map={}
    for rel in src.part.rels.values():
        if 'image' not in rel.reltype: continue
        try:
            _,new_rId=dst.part.get_or_add_image_part(io.BytesIO(rel.target_part.blob))
            rId_map[rel.rId]=new_rId
        except Exception as e:
            app.logger.warning(f'Image copy rId={rel.rId}: {e}')
    if rId_map:
        xml_str=etree.tostring(dst_tree).decode()
        xml_fixed=_replace_rids_single_pass(xml_str,rId_map)
        new_tree=etree.fromstring(xml_fixed.encode())
        dst_tree.getparent().replace(dst_tree,new_tree)
    return dst

def get_group_key(source_id):
    if len(source_id)<6: return source_id.upper()
    substr=source_id[5:8]
    alpha=re.sub(r'[^A-Za-z]','',substr).upper()
    return alpha if alpha else substr.upper()

def build_report(records):
    groups=OrderedDict()
    for rec in sorted(records,key=lambda r:(r.get('source_site',''),r.get('e_token',''))):
        site=rec.get('source_site','Unknown')
        groups.setdefault(site,[]).append(rec)
    out=Presentation(); out.slide_width=Emu(SLIDE_W); out.slide_height=Emu(SLIDE_H)
    slide_num=1
    for site,site_recs in groups.items():
        sn_offset=1
        for chunk in [site_recs[i:i+5] for i in range(0,len(site_recs),5)]:
            add_summary_slide(out,chunk,slide_num,sn_offset=sn_offset)
            sn_offset+=len(chunk); slide_num+=1
        for sn_idx,rec in enumerate(site_recs,start=1):
            tn=rec.get('ticket_no','')
            if tn in PPT_DB['slide_map']:
                pi,si=PPT_DB['slide_map'][tn]
                try: clone_slide(PPT_DB['presentations'][pi],si,out,remove_slide_number=True,sn=sn_idx)
                except Exception as e: app.logger.warning(f'Clone failed {tn}: {e}')
            slide_num+=1
    buf=io.BytesIO(); out.save(buf); buf.seek(0); return buf

def build_excel_ppt(records):
    wb=openpyxl.Workbook(); ws=wb.active; ws.title='Rejection Report'
    hdr_font=Font(name='Calibri',bold=True,color='FFFFFF',size=11)
    hdr_fill=PatternFill('solid',fgColor='00B050')
    hdr_align=Alignment(horizontal='center',vertical='center',wrap_text=True)
    hdr_border=Border(left=Side(style='thin'),right=Side(style='thin'),top=Side(style='thin'),bottom=Side(style='thin'))
    data_font=Font(name='Calibri',size=10)
    data_align=Alignment(horizontal='center',vertical='center',wrap_text=True)
    data_border=Border(left=Side(style='thin'),right=Side(style='thin'),top=Side(style='thin'),bottom=Side(style='thin'))
    HEADERS=['S/N','Ticket No','Veh No','Material','Source Site','Date In','Time In','E-Token','Accepted','Reject Reason']
    COL_W_XLSX=[6,18,12,12,18,14,10,32,10,26]
    for ci,(h,w) in enumerate(zip(HEADERS,COL_W_XLSX),start=1):
        cell=ws.cell(row=1,column=ci,value=h)
        cell.font=hdr_font; cell.fill=hdr_fill; cell.alignment=hdr_align; cell.border=hdr_border
        ws.column_dimensions[get_column_letter(ci)].width=w
    ws.row_dimensions[1].height=28
    sorted_recs=sorted(records,key=lambda r:(r.get('source_site',''),r.get('e_token','')))
    current_site=None; row=2; sn=1
    for rec in sorted_recs:
        site=rec.get('source_site','')
        if site!=current_site:
            if current_site is not None:
                ws.row_dimensions[row].height=6; row+=1
            ws.merge_cells(start_row=row,start_column=1,end_row=row,end_column=10)
            gc=ws.cell(row=row,column=1,value=f'▶  Source Site: {site}')
            gc.font=Font(name='Calibri',bold=True,size=11,color='1B5E20')
            gc.fill=PatternFill('solid',fgColor='C8E6C9')
            gc.alignment=Alignment(horizontal='left',vertical='center')
            gc.border=hdr_border; ws.row_dimensions[row].height=22; row+=1; current_site=site; sn=1
        vals=[sn,rec.get('ticket_no',''),rec.get('veh_no',''),rec.get('material',''),
              rec.get('source_site',''),rec.get('date',''),rec.get('time',''),
              rec.get('e_token',''),rec.get('accepted','NO'),rec.get('reject_reason','')]
        is_alt=(sn%2==0)
        row_fill=PatternFill('solid',fgColor='F0FAF4') if is_alt else None
        for ci,v in enumerate(vals,start=1):
            cell=ws.cell(row=row,column=ci,value=v)
            cell.font=data_font; cell.border=data_border; cell.alignment=data_align
            if row_fill: cell.fill=row_fill
        ws.row_dimensions[row].height=20; row+=1; sn+=1
    ws.freeze_panes='A2'
    ws.auto_filter.ref=f'A1:{get_column_letter(len(HEADERS))}1'
    buf=io.BytesIO(); wb.save(buf); buf.seek(0); return buf

def build_zip_ppt(records):
    import zipfile as _zipfile
    company_groups=OrderedDict()
    for rec in sorted(records,key=lambda r:(get_group_key(r.get('source_site','')),r.get('source_site',''),r.get('e_token',''))):
        gkey=get_group_key(rec.get('source_site',''))
        company_groups.setdefault(gkey,[]).append(rec)
    zip_buf=io.BytesIO()
    with _zipfile.ZipFile(zip_buf,'w',_zipfile.ZIP_DEFLATED) as zf:
        for gkey,grp_recs in company_groups.items():
            pptx_buf=build_report(grp_recs)
            safe_key=re.sub(r'[\/:*?"<>|]','_',gkey)
            zf.writestr(f'{safe_key}_Rejected_Reports/{safe_key}_Rejected_Report.pptx',pptx_buf.read())
    zip_buf.seek(0); return zip_buf

# ═══════════════════════════════════════════════════════════════════════════════
#  PHOTO MERGE WORKER
# ═══════════════════════════════════════════════════════════════════════════════

def _save_upload(file_obj, tag=""):
    name=secure_filename(file_obj.filename or "upload.pptx")
    uid=uuid.uuid4().hex[:8]
    path=os.path.join(UPLOAD_DIR,f"{uid}_{tag}_{name}")
    file_obj.save(path); return path

def _run_merge(job_id,path_a,path_b,top_ph,front_ph,verbose,
               top_h_cm,top_w_cm,top_left_cm,front_h_cm,front_w_cm,front_left_cm,center_gap_cm):
    def upd(pct,msg):
        with jobs_lock:
            jobs[job_id]["progress"]=pct; jobs[job_id]["status"]=msg
    with jobs_lock:
        jobs[job_id]={"status":"Starting…","progress":0,"log_lines":[],"result_file":None,
                      "output_name":"","reference_file":"?","error":None,"stats":None}
    try:
        upd(5,"Loading files…"); time.sleep(0.1)
        upd(25,"Extracting tokens and images…")
        tmp_out=os.path.join(OUTPUT_DIR,f"{job_id}_output.pptx")
        result=ppt_merger.merge(path_a,path_b,tmp_out,
            top_ph=top_ph,front_ph=front_ph,verbose=verbose,
            top_h_cm=top_h_cm,top_w_cm=top_w_cm,top_left_cm=top_left_cm,
            front_h_cm=front_h_cm,front_w_cm=front_w_cm,front_left_cm=front_left_cm,
            center_gap_cm=center_gap_cm)
        upd(90,"Finalising…"); time.sleep(0.15)
        final_path=getattr(result,"_final_path",None)
        if final_path and not os.path.exists(final_path): final_path=tmp_out
        with jobs_lock:
            jobs[job_id]["log_lines"]=result.log_lines
            jobs[job_id]["stats"]=result.to_dict()
            jobs[job_id]["output_name"]=result.output_name
            jobs[job_id]["reference_file"]=result.reference_file
            jobs[job_id]["progress"]=100
            if result.ok and final_path and os.path.exists(final_path):
                jobs[job_id]["status"]="complete"
                jobs[job_id]["result_file"]=final_path
                jobs[job_id]["output_name"]=result.output_name
            else:
                err="; ".join(result.errors) if result.errors else "Merge failed"
                jobs[job_id]["status"]="error"; jobs[job_id]["error"]=err
    except Exception as e:
        with jobs_lock:
            jobs[job_id]["status"]="error"; jobs[job_id]["error"]=str(e); jobs[job_id]["progress"]=100
    finally:
        for p in [path_a,path_b]:
            if p and os.path.exists(p):
                try: os.remove(p)
                except: pass

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Auth
# ═══════════════════════════════════════════════════════════════════════════════


@app.route('/static/bg.jpg')
def serve_bg():
    import os
    path = os.path.join(os.path.dirname(__file__), 'static', 'bg.jpg')
    return send_file(path, mimetype='image/jpeg', max_age=86400)


@app.route("/login", methods=["GET","POST"])
def login_page():
    if request.method=="POST":
        data=request.get_json(force=True) if request.is_json else request.form
        username=data.get("username","").strip().lower()
        password=data.get("password","")
        conn=get_db()
        user=conn.execute("SELECT * FROM users WHERE username=?",(username,)).fetchone()
        conn.close()
        if user and verify_pw(password,user["password_hash"]):
            session["username"]=user["username"]
            session["name"]=user["name"]
            session["role"]=user["role"]
            log_activity(username,"LOGIN","Successful login")
            if request.is_json: return jsonify({"ok":True,"name":user["name"],"role":user["role"]})
            return redirect(url_for("dashboard"))
        if request.is_json: return jsonify({"ok":False,"error":"Invalid username or password"}),401
        return redirect(url_for("login_page"))
    return render_template_string(AUTH_HTML, page="login")

@app.route("/register", methods=["POST"])
def register():
    data=request.get_json(force=True)
    name=data.get("name","").strip()
    username=data.get("username","").strip().lower()
    password=data.get("password","")
    confirm=data.get("confirm","")
    if not all([name,username,password]):
        return jsonify({"ok":False,"error":"All fields required"}),400
    if password!=confirm:
        return jsonify({"ok":False,"error":"Passwords do not match"}),400
    if len(password)<6:
        return jsonify({"ok":False,"error":"Password must be at least 6 characters"}),400
    try:
        conn=get_db()
        conn.execute("INSERT INTO users (username,name,password_hash,plaintext_pw,role) VALUES (?,?,?,?,'user')",
                     (username,name,hash_pw(password),password))
        conn.commit(); conn.close()
        log_activity(username,"REGISTER","New user registered")
        # Auto-login after registration
        session["username"] = username
        session["name"] = name
        session["role"] = "user"
        return jsonify({"ok":True,"name":name,"role":"user"})
    except sqlite3.IntegrityError:
        return jsonify({"ok":False,"error":"Username already exists"}),409

@app.route("/logout")
def logout():
    username=session.get("username","")
    if username: log_activity(username,"LOGOUT","")
    session.clear()
    return redirect(url_for("login_page"))

@app.route("/")
@login_required
def dashboard():
    return render_template_string(DASHBOARD_HTML,
        name=session.get("name",""), username=session.get("username",""),
        role=session.get("role",""))

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — App pages (served as embedded SPAs)
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/app/ppt-rejection")
@login_required
def ppt_rejection_page():
    log_activity(session["username"],"OPEN_APP","PPT Alignment / Rejection Filter")
    return render_template_string(PPT_REJECTION_HTML)

@app.route("/app/photo-merge")
@login_required
def photo_merge_page():
    log_activity(session["username"],"OPEN_APP","Front and Top Photo Merge")
    return render_template_string(PHOTO_MERGE_HTML)

@app.route("/app/excel-rejection")
@login_required
def excel_rejection_page():
    log_activity(session["username"],"OPEN_APP","Excel Rejection Report")
    return render_template_string(EXCEL_REJECTION_HTML)

@app.route("/app/daily-report")
@login_required
def daily_report_page():
    log_activity(session["username"],"OPEN_APP","Daily Report")
    return render_template_string(DAILY_REPORT_HTML)

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — PPT Rejection API
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/ppt/upload", methods=["POST"])
@login_required
def ppt_upload():
    try:
        if request.args.get("reset","true").lower()=="true":
            PPT_DB["records"]=[]; PPT_DB["slide_map"]={}; PPT_DB["presentations"]=[]
        loaded,errors=0,[]
        for f in request.files.getlist("files"):
            if not f.filename: continue
            name=f.filename
            try:
                data=f.read()
                if not data: errors.append(f'{name}: empty'); continue
                prs=Presentation(io.BytesIO(data))
                prs_idx=len(PPT_DB["presentations"])
                PPT_DB["presentations"].append(prs)
                recs,smap=extract_from_prs(prs,prs_idx)
                PPT_DB["records"].extend(recs); PPT_DB["slide_map"].update(smap); loaded+=1
            except Exception as e: errors.append(f'{name}: {e}')
        out_recs=[{k:v for k,v in r.items() if k!='_prs'} for r in PPT_DB["records"]]
        result={"records":out_recs,"files_loaded":loaded,"errors":errors}
        log_activity(session["username"],"PPT_UPLOAD",f"{loaded} files, {len(out_recs)} records")
        if loaded==0 and errors:
            result["error"]="All files failed: "+"; ".join(errors[:3])
            return jsonify(result),400
        return jsonify(result)
    except Exception as e:
        return jsonify({"error":str(e),"records":[],"files_loaded":0,"errors":[]}),500

@app.route("/api/ppt/generate", methods=["POST"])
@login_required
def ppt_generate():
    try:
        records=request.get_json(force=True).get("records",[])
        if not records: return jsonify({"error":"No records"}),400
        buf=build_report(records)
        log_activity(session["username"],"PPT_GENERATE",f"{len(records)} records")
        return send_file(buf,mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                         as_attachment=True,download_name="Rejection_Report.pptx")
    except Exception as e:
        return jsonify({"error":str(e)}),500

@app.route("/api/ppt/export_excel", methods=["POST"])
@login_required
def ppt_export_excel():
    try:
        records=request.get_json(force=True).get("records",[])
        if not records: return jsonify({"error":"No records"}),400
        buf=build_excel_ppt(records)
        fname=f"Rejection_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        log_activity(session["username"],"PPT_EXCEL",f"{len(records)} records")
        return send_file(buf,mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                         as_attachment=True,download_name=fname)
    except Exception as e:
        return jsonify({"error":str(e)}),500

@app.route("/api/ppt/generate_zip", methods=["POST"])
@login_required
def ppt_generate_zip():
    try:
        records=request.get_json(force=True).get("records",[])
        if not records: return jsonify({"error":"No records"}),400
        buf=build_zip_ppt(records)
        fname=f"Rejection_Report_ZIP_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
        log_activity(session["username"],"PPT_ZIP",f"{len(records)} records")
        return send_file(buf,mimetype="application/zip",as_attachment=True,download_name=fname)
    except Exception as e:
        return jsonify({"error":str(e)}),500

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Photo Merge API
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/merge", methods=["POST"])
@login_required
def api_merge():
    file_a=request.files.get("file_a"); file_b=request.files.get("file_b")
    if not file_a: return jsonify({"error":"File A is required"}),400
    if not file_b: return jsonify({"error":"File B is required"}),400
    path_a=_save_upload(file_a,"A"); path_b=_save_upload(file_b,"B")
    job_id=uuid.uuid4().hex
    top_ph=request.form.get("top_placeholder",ppt_merger.TOP_PH_NAME)
    front_ph=request.form.get("front_placeholder",ppt_merger.FRONT_PH_NAME)
    verbose=request.form.get("verbose","0")=="1"
    top_h_cm=float(request.form.get("top_h_cm",9.11)); top_w_cm=float(request.form.get("top_w_cm",15.28))
    top_left_cm=float(request.form.get("top_left_cm",0.95)); front_h_cm=float(request.form.get("front_h_cm",9.11))
    front_w_cm=float(request.form.get("front_w_cm",15.51)); front_left_cm=float(request.form.get("front_left_cm",18.73))
    center_gap_cm=float(request.form.get("center_gap_cm",2.5))
    t=threading.Thread(target=_run_merge,args=(job_id,path_a,path_b,top_ph,front_ph,verbose,
        top_h_cm,top_w_cm,top_left_cm,front_h_cm,front_w_cm,front_left_cm,center_gap_cm),daemon=True)
    t.start()
    log_activity(session["username"],"PHOTO_MERGE_START",f"job={job_id}")
    return jsonify({"job_id":job_id})

@app.route("/api/job/<job_id>")
@login_required
def api_job_status(job_id):
    with jobs_lock:
        job=jobs.get(job_id)
    if not job: return jsonify({"error":"Job not found"}),404
    return jsonify({"status":job["status"],"progress":job["progress"],"log_lines":job["log_lines"],
                    "stats":job["stats"],"error":job["error"],"output_name":job.get("output_name",""),
                    "reference_file":job.get("reference_file","?"),"has_result":job["result_file"] is not None})

@app.route("/api/download/<job_id>")
@login_required
def api_download_merge(job_id):
    with jobs_lock:
        job=jobs.get(job_id)
    if not job or not job.get("result_file"): return jsonify({"error":"No result"}),404
    path=job["result_file"]
    if not os.path.exists(path): return jsonify({"error":"File not found"}),404
    display_name=job.get("output_name") or os.path.basename(path)
    log_activity(session["username"],"PHOTO_MERGE_DOWNLOAD",display_name)
    return send_file(path,as_attachment=True,download_name=display_name,
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Excel Rejection API (Streamlit-like logic via REST)
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/excel/dates", methods=["POST"])
@login_required
def excel_get_dates():
    try:
        import sys; sys.path.insert(0,os.path.dirname(__file__))
        from generate_ppt_excel import get_available_dates, format_date
        f=request.files.get("file")
        if not f: return jsonify({"error":"No file"}),400
        dates=get_available_dates(f)
        return jsonify({"dates":[{"label":format_date(d),"value":str(d)} for d in dates]})
    except Exception as e:
        return jsonify({"error":str(e)}),500

@app.route("/api/excel/preview", methods=["POST"])
@login_required
def excel_preview():
    """Return stats (rejections, sites, est_slides) for the selected date — matches standalone."""
    try:
        import sys; sys.path.insert(0, os.path.dirname(__file__))
        from generate_ppt_excel import load_and_filter, format_date
        import datetime as dt
        f = request.files.get("file")
        date_str = request.form.get("date", "")
        if not f:
            return jsonify({"error": "No file"}), 400
        selected_date = dt.date.fromisoformat(date_str) if date_str else None
        groups, data_date = load_and_filter(f, filter_date=selected_date)
        total_rej   = sum(len(v) for v in groups.values())
        total_sites = len(groups)
        est_slides  = 1 + total_sites * 2  # cover + (summary+detail) per site
        eff_date    = selected_date or data_date
        badges = [{"site": site, "count": len(rows)} for site, rows in groups.items()]
        return jsonify({
            "ok": True,
            "rejections": total_rej,
            "sites": total_sites,
            "est_slides": est_slides,
            "date_label": format_date(eff_date) if eff_date else "",
            "badges": badges,
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/excel/generate", methods=["POST"])
@login_required
def excel_generate():
    try:
        import sys; sys.path.insert(0,os.path.dirname(__file__))
        from generate_ppt_excel import generate_ppt, get_available_dates
        import datetime as dt
        f=request.files.get("file")
        date_str=request.form.get("date","")
        if not f: return jsonify({"error":"No file"}),400
        selected_date=dt.date.fromisoformat(date_str) if date_str else None
        f.seek(0)
        ppt_bytes=generate_ppt(f,report_date_obj=None,template_path=None,
                                photo_folder=None,filter_date=selected_date)
        eff=selected_date or dt.date.today()
        fname=f"APSG-Loads_Rejected-{eff.strftime('%d%m%Y')}.pptx"
        log_activity(session["username"],"EXCEL_REJECTION",fname)
        return send_file(io.BytesIO(ppt_bytes),
                         mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                         as_attachment=True,download_name=fname)
    except Exception as e:
        return jsonify({"error":str(e)}),500

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Admin Panel
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/admin")
@login_required
def admin_panel():
    if session.get("role")!="admin":
        return redirect(url_for("dashboard"))
    return render_template_string(ADMIN_HTML, name=session.get("name",""))

@app.route("/api/admin/users")
@login_required
@admin_required
def admin_users():
    conn=get_db()
    users=conn.execute("SELECT id,username,name,role,created_at FROM users ORDER BY created_at DESC").fetchall()
    conn.close()
    return jsonify([dict(u) for u in users])

@app.route("/api/admin/activity")
@login_required
@admin_required
def admin_activity():
    limit  = int(request.args.get("limit", 200))
    action = request.args.get("action", "").strip().upper()
    conn   = get_db()
    if action:
        rows = conn.execute(
            "SELECT * FROM activity_log WHERE action=? ORDER BY ts DESC LIMIT ?",
            (action, limit)
        ).fetchall()
    else:
        rows = conn.execute(
            "SELECT * FROM activity_log ORDER BY ts DESC LIMIT ?", (limit,)
        ).fetchall()
    conn.close()
    return jsonify([dict(r) for r in rows])

@app.route("/api/admin/activity/download")
@login_required
@admin_required
def admin_activity_download():
    conn=get_db()
    rows=conn.execute("SELECT * FROM activity_log ORDER BY ts DESC").fetchall()
    conn.close()
    wb=openpyxl.Workbook(); ws=wb.active; ws.title="Activity Log"
    ws.append(["ID","Username","Action","Detail","IP","Timestamp"])
    for r in rows: ws.append(list(r))
    buf=io.BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf,mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                     as_attachment=True,download_name=f"activity_log_{datetime.now().strftime('%Y%m%d')}.xlsx")

@app.route("/api/health")
def health():
    return jsonify({"status":"ok","version":"1.0.0","app":"APSG (Staging Ground) Report"})

@app.route("/api/admin/user/<username>/password", methods=["POST"])
@login_required
@admin_required
def admin_change_password(username):
    """Admin changes any user's password."""
    data = request.get_json(force=True)
    new_pw = data.get("password","").strip()
    if len(new_pw) < 6:
        return jsonify({"ok":False,"error":"Password must be at least 6 characters"}), 400
    conn = get_db()
    row = conn.execute("SELECT id FROM users WHERE username=?", (username,)).fetchone()
    if not row:
        conn.close()
        return jsonify({"ok":False,"error":"User not found"}), 404
    conn.execute("UPDATE users SET password_hash=?,plaintext_pw=? WHERE username=?", (hash_pw(new_pw), new_pw, username))
    conn.commit()
    conn.close()
    log_activity(session["username"], "ADMIN_CHANGE_PW", f"Changed password for user: {username}")
    return jsonify({"ok":True})

@app.route("/api/admin/user/<username>/view_password", methods=["GET"])
@login_required
@admin_required
def admin_view_password(username):
    """
    Return the stored plaintext password if available.
    If plaintext_pw is set (stored at registration/change), return it directly.
    Otherwise generate and set a new temp password (old behaviour).
    """
    conn = get_db()
    row = conn.execute(
        "SELECT id,name,plaintext_pw FROM users WHERE username=?", (username,)
    ).fetchone()
    if not row:
        conn.close()
        return jsonify({"ok":False,"error":"User not found"}), 404
    stored_pw = (row["plaintext_pw"] or "").strip()
    if stored_pw:
        # Return the plaintext password stored at registration/last change
        conn.close()
        log_activity(session["username"], "ADMIN_VIEW_PW", f"Viewed password for: {username}")
        return jsonify({"ok":True,"password":stored_pw,"username":username,"name":row["name"],"source":"stored"})
    # Fallback: generate temp password
    import secrets, string
    temp_pw = ''.join(secrets.choice(string.ascii_letters + string.digits) for _ in range(12))
    conn.execute("UPDATE users SET password_hash=?,plaintext_pw=? WHERE username=?",
                 (hash_pw(temp_pw), temp_pw, username))
    conn.commit()
    conn.close()
    log_activity(session["username"], "ADMIN_VIEW_PW", f"Generated temp password for: {username}")
    return jsonify({"ok":True,"password":temp_pw,"username":username,"name":row["name"],"source":"generated"})

@app.route("/api/admin/user/<username>/delete", methods=["DELETE"])
@login_required
@admin_required
def admin_delete_user(username):
    """Admin deletes a user (cannot delete own account or built-in admin)."""
    if username == session.get("username"):
        return jsonify({"ok":False,"error":"Cannot delete your own account"}), 400
    if username == "admin":
        return jsonify({"ok":False,"error":"Cannot delete the built-in admin account"}), 400
    conn = get_db()
    conn.execute("DELETE FROM users WHERE username=?", (username,))
    conn.commit()
    conn.close()
    log_activity(session["username"], "ADMIN_DELETE_USER", f"Deleted user: {username}")
    return jsonify({"ok":True})

@app.route("/api/admin/analytics")
@login_required
@admin_required
def admin_analytics():
    conn = get_db()
    users = conn.execute("SELECT id,username,name,role,created_at FROM users ORDER BY created_at DESC").fetchall()
    # Login attempts
    logins = conn.execute("""SELECT username, COUNT(*) as cnt, MAX(ts) as last_login
        FROM activity_log WHERE action='LOGIN' GROUP BY username ORDER BY last_login DESC""").fetchall()
    # All activity summary per user
    activity_sum = conn.execute("""SELECT username, COUNT(*) as total_actions,
        MIN(ts) as first_seen, MAX(ts) as last_seen,
        GROUP_CONCAT(DISTINCT action) as actions_used
        FROM activity_log GROUP BY username ORDER BY last_seen DESC""").fetchall()
    # Login attempts (failed = no matching login row within same minute)
    attempts = conn.execute("""SELECT username, COUNT(*) as attempts, MAX(ts) as last_ts
        FROM activity_log WHERE action IN ('LOGIN','LOGIN_FAIL')
        GROUP BY username ORDER BY last_ts DESC""").fetchall()
    total_users = len(users)
    total_logins = conn.execute("SELECT COUNT(*) FROM activity_log WHERE action='LOGIN'").fetchone()[0]
    total_actions = conn.execute("SELECT COUNT(*) FROM activity_log").fetchone()[0]
    conn.close()
    return jsonify({
        'total_users': total_users,
        'total_logins': total_logins,
        'total_actions': total_actions,
        'users': [dict(u) for u in users],
        'logins': [dict(l) for l in logins],
        'activity_summary': [dict(a) for a in activity_sum],
    })

@app.route("/api/admin/analytics/download")
@login_required
@admin_required
def admin_analytics_download():
    conn = get_db()
    users = conn.execute("SELECT id,username,name,role,created_at FROM users ORDER BY created_at").fetchall()
    activity = conn.execute("""SELECT a.username, u.name, a.action, a.detail, a.ip, a.ts
        FROM activity_log a LEFT JOIN users u ON a.username=u.username
        ORDER BY a.ts DESC""").fetchall()
    conn.close()
    wb = openpyxl.Workbook()
    # Users sheet
    ws1 = wb.active; ws1.title = 'Users'
    ws1.append(['ID','Username','Name','Role','Registered'])
    hf = Font(bold=True,color='FFFFFF'); hfill = PatternFill('solid',fgColor='1E3A8A')
    for cell in ws1[1]: cell.font=hf; cell.fill=hfill; cell.alignment=Alignment(horizontal='center')
    for u in users: ws1.append(list(u))
    for col in ws1.columns:
        ws1.column_dimensions[col[0].column_letter].width = max(len(str(c.value or '')) for c in col)+4
    # Activity sheet
    ws2 = wb.create_sheet('Activity Log')
    ws2.append(['Username','Full Name','Action','Detail','IP','Timestamp'])
    for cell in ws2[1]: cell.font=hf; cell.fill=hfill; cell.alignment=Alignment(horizontal='center')
    for row in activity: ws2.append(list(row))
    for col in ws2.columns:
        ws2.column_dimensions[col[0].column_letter].width = max(len(str(c.value or '')) for c in col)+4
    # Summary sheet
    ws3 = wb.create_sheet('Usage Summary')
    ws3.append(['Username','Total Actions','First Seen','Last Seen'])
    for cell in ws3[1]: cell.font=hf; cell.fill=hfill; cell.alignment=Alignment(horizontal='center')
    conn2 = get_db()
    rows = conn2.execute("""SELECT username,COUNT(*) as total,MIN(ts),MAX(ts)
        FROM activity_log GROUP BY username ORDER BY total DESC""").fetchall()
    conn2.close()
    for r in rows: ws3.append(list(r))
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf,mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                     as_attachment=True,download_name=f'APSG_Analytics_{datetime.now().strftime("%Y%m%d")}.xlsx')

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Daily Report API (engine-powered, no Streamlit needed)
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/daily/upload", methods=["POST"])
@login_required
def daily_upload():
    """Load and validate the Online export file, return metadata + date range."""
    try:
        f = request.files.get("file")
        if not f:
            return jsonify({"error": "No file uploaded"}), 400
        result = load_and_validate_file(f)
        if result["error"]:
            return jsonify({"error": result["error"]}), 400
        df = result["df"]
        min_date, max_date = get_date_range(df)
        # Store df in session via temp file
        import tempfile, pickle
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pkl",
                                          dir=OUTPUT_DIR)
        pickle.dump(df, tmp); tmp.close()
        return jsonify({
            "ok": True,
            "rows": len(df),
            "filename": f.filename,
            "tmp_path": tmp.name,
            "min_date": str(min_date),
            "max_date": str(max_date),
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/daily/generate", methods=["POST"])
@login_required
def daily_generate():
    """Generate the Online report Excel file."""
    try:
        import pickle
        data = request.get_json(force=True)
        tmp_path = data.get("tmp_path", "")
        filter_date_str = data.get("filter_date", "")
        corrections = data.get("corrections", {})

        if not tmp_path or not os.path.exists(tmp_path):
            return jsonify({"error": "Session expired — please re-upload the file"}), 400

        with open(tmp_path, "rb") as fh:
            df = pickle.load(fh)

        from datetime import date as dt_date
        filter_date = dt_date.fromisoformat(filter_date_str) if filter_date_str else None
        if not filter_date:
            return jsonify({"error": "Filter date required"}), 400

        # Build corrections dict with int keys
        corr_int = {int(k): v for k, v in corrections.items()}

        def _build_flag_map(corr_dict):
            fm = {}
            for idx, corr in corr_dict.items():
                decision = corr.get("Accepted", "Yes")
                red_cols = ({"Accepted", "Out Weight", "Net Weight"}
                            if str(decision).strip().lower() == "no"
                            else {"Out Weight", "Net Weight"})
                fm[idx] = {"flagged": True, "red_cols": red_cols}
            return fm

        flag_map = _build_flag_map(corr_int)
        excel_buf, fname, stats = generate_report(
            df, filter_date, filter_date,
            corrections=corr_int, flag_map=flag_map,
        )
        log_activity(session["username"], "DAILY_REPORT", fname)
        return send_file(
            excel_buf,
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            as_attachment=True,
            download_name=fname,
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/daily/validate", methods=["POST"])
@login_required
def daily_validate():
    """Run all validations on filtered data and return results."""
    try:
        import pickle
        data = request.get_json(force=True)
        tmp_path = data.get("tmp_path", "")
        filter_date_str = data.get("filter_date", "")
        corrections = data.get("corrections", {})

        if not tmp_path or not os.path.exists(tmp_path):
            return jsonify({"error": "Session expired — please re-upload the file"}), 400

        with open(tmp_path, "rb") as fh:
            df = pickle.load(fh)

        from datetime import date as dt_date
        filter_date = dt_date.fromisoformat(filter_date_str) if filter_date_str else None
        if not filter_date:
            return jsonify({"error": "Filter date required"}), 400

        preview = filter_preview(df, filter_date, filter_date)
        if preview.empty:
            return jsonify({"ok": True, "rows": 0, "incomplete": [], "errors": [], "stats": {}})

        corr_int = {int(k): v for k, v in corrections.items()}
        val = validate_and_flag(preview)
        incomplete_list = []
        for idx in val["incomplete_accepted"]:
            row = preview.loc[idx]
            incomplete_list.append({
                "idx": int(idx),
                "token": str(row.get("Token", idx)),
                "in_weight": safe_float(row.get("In Weight")),
                "veh": str(row.get("Vehicle Number", "") or ""),
            })

        # Apply corrections to preview
        preview_corr = preview.copy()
        for idx, corr in corr_int.items():
            if idx in preview_corr.index:
                if "Accepted"   in corr: preview_corr.at[idx, "Accepted"]   = corr["Accepted"]
                if "Out Weight" in corr: preview_corr.at[idx, "Out Weight"] = float(corr["Out Weight"])
                if "Net Weight" in corr: preview_corr.at[idx, "Net Weight"] = float(corr["Net Weight"])

        nw_errors = validate_net_weights(preview_corr)
        errors = [f"E-Token [{e['token']}] – Ledger mismatch" for e in nw_errors]

        # Stats
        acc_mask = preview_corr["Accepted"].astype(str).str.strip().str.lower().isin(
            ("yes","1","true","accepted"))
        rej_mask = preview_corr["Accepted"].astype(str).str.strip().str.lower().isin(
            ("no","0","false","rejected","reject"))
        acc_rows = preview_corr[acc_mask]
        stats = {
            "accepted": int(acc_mask.sum()),
            "rejected": int(rej_mask.sum()),
            "wi": safe_col_sum(acc_rows, "In Weight"),
            "wo": safe_col_sum(acc_rows, "Out Weight"),
            "nw": safe_col_sum(acc_rows, "Net Weight"),
        }

        return jsonify({
            "ok": True,
            "rows": len(preview),
            "incomplete": incomplete_list,
            "errors": errors,
            "stats": stats,
            "all_resolved": len(incomplete_list) == 0 or all(
                item["idx"] in corr_int for item in incomplete_list),
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/daily/wb_upload", methods=["POST"])
@login_required
def daily_wb_upload():
    """Load WB file."""
    try:
        import pickle
        f = request.files.get("file")
        if not f:
            return jsonify({"error": "No file"}), 400
        result = load_wb_file(f)
        if result["error"]:
            return jsonify({"error": result["error"]}), 400
        df = result["df"]
        min_date, max_date = wb_get_date_range(df)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pkl", dir=OUTPUT_DIR)
        pickle.dump(df, tmp); tmp.close()
        return jsonify({
            "ok": True, "rows": len(df), "filename": f.filename,
            "tmp_path": tmp.name,
            "min_date": str(min_date) if min_date else "",
            "max_date": str(max_date) if max_date else "",
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

import tempfile

@app.route("/api/daily/wb_pivot", methods=["POST"])
@login_required
def daily_wb_pivot():
    """
    Build WB pivot using the exact same logic as the standalone Daily Report app.
    Runs:
      1. wb_filter_by_date          — date filter
      2. Normalize Accepted (0/1 → Yes/No)
      3. Filter accepted rows only (remove Accepted=No)
      4. wb_find_incomplete_rows    — detect missing Date Out / Time Out
      5. wb_apply_row_decisions     — apply user decisions for incomplete rows
      6. wb_apply_unified_logic     — sync Online rejections + resolve blank Out Weights via E-Token
      7. wb_net_weight_validation   — per-row In−Out=Net check
      8. wb_validate_etoken_match   — E-Token cross-comparison with Online
      9. wb_build_pivot             — build pivot table
      10. Compute wb_stats from wb_df_for_processing (INTENDED user decisions,
          same as standalone) so comparison table matches Online
    """
    try:
        import pickle
        data = request.get_json(force=True)
        wb_tmp          = data.get("wb_tmp_path", "")
        filter_date_str = data.get("filter_date", "")
        wb_decisions    = data.get("wb_decisions", {})
        online_tmp      = data.get("online_tmp_path", "")
        online_date_str = data.get("online_date", filter_date_str)
        corrections_raw = data.get("corrections", {})

        if not wb_tmp or not os.path.exists(wb_tmp):
            return jsonify({"error": "WB session expired — re-upload WB file"}), 400

        with open(wb_tmp, "rb") as fh:
            wb_df_raw = pickle.load(fh)

        from datetime import date as dt_date
        filter_date = dt_date.fromisoformat(filter_date_str) if filter_date_str else None
        if not filter_date:
            return jsonify({"error": "Filter date required"}), 400

        # ── Step 1: Date filter ──────────────────────────────────────────────
        wb_df_filtered = wb_filter_by_date(wb_df_raw, filter_date)
        if wb_df_filtered.empty:
            return jsonify({"error": f"No WB records for {filter_date_str}"}), 400

        # ── Step 2: Normalize Accepted (0/1 → Yes/No) ───────────────────────
        wb_df_filtered = wb_df_filtered.copy()
        wb_df_filtered["Accepted"] = wb_df_filtered["Accepted"].apply(_wb_norm_acc)

        wb_total  = len(wb_df_filtered)
        _acc_mask = wb_df_filtered["Accepted"].str.strip().str.lower() == "yes"
        wb_rej_raw = int((~_acc_mask).sum())

        # ── Step 3: Filter accepted rows only ───────────────────────────────
        wb_df_accepted = wb_df_filtered[_acc_mask].copy()
        if wb_df_accepted.empty:
            return jsonify({"error": "No accepted records remain after filtering Accepted=No rows."}), 400

        # ── Step 4: Find incomplete rows (missing Date Out / Time Out) ───────
        incomplete_rows = wb_find_incomplete_rows(wb_df_accepted)

        # ── Step 5: Apply user row decisions ────────────────────────────────
        dec_int = {int(k): v for k, v in wb_decisions.items()}

        # Check if all incomplete rows have decisions — if not, return them for UI
        if not incomplete_rows.empty:
            undecided = [
                {
                    "idx": int(idx),
                    "etoken": str(row.get("E-Token", f"Row {idx}")),
                    "in_weight": safe_float(row.get("In Weight", 0)),
                }
                for idx, row in incomplete_rows.iterrows()
                if int(idx) not in dec_int
            ]
            if undecided:
                return jsonify({"ok": True, "incomplete": undecided, "rows": [], "wb_stats": {}, "errors": []})

        _user_rejected_count = sum(1 for d in dec_int.values() if d.get("decision") == 0)
        wb_rej_raw_total = wb_rej_raw + _user_rejected_count

        wb_df_for_processing, _ = wb_apply_row_decisions(wb_df_accepted, dec_int)

        # ── Step 6: Load Online data for cross-comparison ───────────────────
        online_filtered = pd.DataFrame()
        corrections = {}

        if online_tmp and os.path.exists(online_tmp):
            try:
                with open(online_tmp, "rb") as fh:
                    online_df_raw = pickle.load(fh)
                corrections = {int(k): v for k, v in corrections_raw.items()}
                online_filtered = filter_preview(online_df_raw, filter_date, filter_date)
            except Exception:
                pass  # graceful degradation — proceed without Online comparison

        # ── Step 7: wb_apply_unified_logic ──────────────────────────────────
        # Matches standalone: syncs Online rejections + resolves blank Out Weights via E-Token
        wb_proc = wb_apply_unified_logic(wb_df_for_processing, online_filtered, corrections)
        wb_accepted_df = wb_proc["wb_accepted_df"]

        # ── Step 8: Collect validation errors (same as standalone) ──────────
        wb_all_errors = []

        # Net weight validation on wb_df_for_processing (pre-sync, matches standalone)
        nw_val_errors = wb_net_weight_validation(wb_df_for_processing)
        for e in nw_val_errors:
            wb_all_errors.append(f"E-Token [{e['etoken']}] – Ledger mismatch")

        # Also include errors from unified logic pass
        for e in wb_proc.get("nw_errors", []):
            msg = f"E-Token [{e['etoken']}] – Ledger mismatch"
            if not any(f"E-Token [{e['etoken']}]" in m and "Ledger mismatch" in m for m in wb_all_errors):
                wb_all_errors.append(msg)

        # E-Token cross-comparison with Online (only if Online data available)
        if not online_filtered.empty:
            etoken_result = wb_validate_etoken_match(
                wb_accepted_df, online_filtered, corrections,
                online_rejected_tokens=wb_proc.get("synced_tokens", set()),
            )
            for tok in etoken_result.get("wb_only", []):
                wb_all_errors.append(f"E-Token [{tok}] – Mismatch between Online & Weighbridge")
            for tok in etoken_result.get("online_only", []):
                wb_all_errors.append(f"E-Token [{tok}] – Mismatch between Online & Weighbridge")
            for tok in etoken_result.get("wb_dupes", []):
                wb_all_errors.append(f"E-Token [{tok}] – Duplicate E-Token in Weighbridge data")

        # ── Step 9: Build pivot ──────────────────────────────────────────────
        pivot_df = wb_build_pivot(wb_accepted_df)

        rows = []
        for _, r in pivot_df.iterrows():
            rows.append({
                "type":  r.get("_row_type", "data"),
                "label": str(r.get("Row Labels", "")),
                "loads": safe_int(r.get("Sum of Loads", 0)) if r.get("_row_type") != "mat_header" else 0,
                "wi":    safe_float(r.get("Sum of Weight In (T)", 0)),
                "wo":    safe_float(r.get("Sum of Weight Out (T)", 0)),
                "nw":    safe_float(r.get("Sum of Net Weight (T)", 0)),
            })

        # ── Step 10: Stats — use wb_df_for_processing (INTENDED accepted count) ─
        # Matches standalone: uses pre-sync df so stats align with what user decided,
        # not what was removed by Online-sync (which would cause count mismatch display).
        _wb_proc_acc_mask = (
            wb_df_for_processing["Accepted"].astype(str).str.strip().str.lower().isin(
                ("yes", "1", "true", "accepted")
            )
            if not wb_df_for_processing.empty and "Accepted" in wb_df_for_processing.columns
            else pd.Series(dtype=bool)
        )
        _wb_proc_accepted = (
            wb_df_for_processing[_wb_proc_acc_mask]
            if not wb_df_for_processing.empty and not _wb_proc_acc_mask.empty
            else pd.DataFrame()
        )
        wb_stats = {
            "accepted": int(_wb_proc_acc_mask.sum()) if not _wb_proc_acc_mask.empty else 0,
            "rejected": wb_rej_raw_total,
            "wi": safe_col_sum(_wb_proc_accepted, "In Weight"),
            "wo": safe_col_sum(_wb_proc_accepted, "Out Weight"),
            "nw": safe_col_sum(_wb_proc_accepted, "Net Weight"),
        }

        # ── Excel pivot download ─────────────────────────────────────────────
        dl_bytes, dl_fname = wb_pivot_to_excel(pivot_df, filter_date=filter_date)
        pivot_tmp = os.path.join(OUTPUT_DIR, f"pivot_{uuid.uuid4().hex}.xlsx")
        with open(pivot_tmp, "wb") as fh:
            fh.write(dl_bytes)

        return jsonify({
            "ok": True,
            "rows": rows,
            "wb_stats": wb_stats,
            "errors": wb_all_errors,
            "pivot_tmp": pivot_tmp,
            "pivot_fname": dl_fname,
        })
    except Exception as e:
        import traceback as _tb
        return jsonify({"error": str(e), "trace": _tb.format_exc()}), 500


@app.route("/api/daily/wb_pivot_download")
@login_required
def daily_wb_pivot_download():
    path = request.args.get("path","")
    fname = request.args.get("fname","pivot.xlsx")
    if not path or not os.path.exists(path):
        return jsonify({"error":"File not found"}),404
    return send_file(path, as_attachment=True, download_name=fname,
                     mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTE — Rectification Report (Daily Report → Action Required)
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/daily/generate_rr", methods=["POST"])
@login_required
def daily_generate_rr():
    """
    Generate a Rectification Report .docx for a single Action Required row.

    Expected JSON payload
    ---------------------
    {
        "tmp_path"     : str,   # path to the pickled online DataFrame
        "token"        : str,   # E-Token of the affected row
        "rr_serial"    : str,   # 4-digit serial e.g. "0290"
        "accepted"     : str,   # "YES" or "NO"
        "out_weight"   : float, # operator-entered Out Weight (0 if NO)
        "net_weight"   : float, # computed Net Weight
        "weight_label" : str,   # optional page-5 heading override (YES only)
        "reason"       : str    # "Accepted Towing Vehicle" | "Rejected Towing Vehicle" | "Late Time / Breakdown"
    }
    """
    try:
        import pickle
        from datetime import timedelta

        data         = request.get_json(force=True)
        tmp_path     = data.get("tmp_path", "")
        token        = str(data.get("token", "")).strip()
        rr_serial    = str(data.get("rr_serial", "")).strip()
        accepted     = str(data.get("accepted", "")).upper().strip()
        out_weight   = data.get("out_weight", 0)
        net_weight   = data.get("net_weight", 0)
        weight_label = str(data.get("weight_label", "")).strip()
        reason       = str(data.get("reason", REASON_A)).strip() or REASON_A
        filter_date  = str(data.get("filter_date", "")).strip()  # YYYY-MM-DD from UI

        # Validate inputs
        if not tmp_path or not os.path.exists(tmp_path):
            return jsonify({"error": "Session expired — please re-upload the Online file"}), 400
        if not token:
            return jsonify({"error": "E-Token is required"}), 400
        if not rr_serial:
            return jsonify({"error": "RR Serial Number is required"}), 400
        if accepted not in ("YES", "NO"):
            return jsonify({"error": "Accepted must be YES or NO"}), 400

        # Load Online DataFrame
        with open(tmp_path, "rb") as fh:
            df = pickle.load(fh)

        # Build action_data dict
        # For NO: out_weight = in_weight (from JS corrections), net_weight = 0
        # For YES: user-entered out_weight and computed net_weight
        # Use explicit None/missing check (not truthiness) so 0 is preserved
        action_data = {
            "ACCEPTED":       accepted,
            "OUT WEIGHT":     str(out_weight) if out_weight not in (None, "", "None") else "",
            "NET WEIGHT":     str(net_weight) if net_weight not in (None, "", "None") else "0",
            "UNLADEN WEIGHT": str(out_weight) if (accepted == "YES" and out_weight) else "",
        }
        # Override for NO: ensure net_weight is always exactly 0
        if accepted == "NO":
            action_data["NET WEIGHT"] = "0"
            # Out weight for NO = In Weight (truck rejected, nothing unloaded)
            if not action_data["OUT WEIGHT"]:
                action_data["OUT WEIGHT"] = str(out_weight) if out_weight else ""

        # Fetch original row
        before_dict = fetch_row_by_token(df, token, source="online")
        after_dict  = apply_user_updates(before_dict, action_data)

        # Parse dates
        arr_dt = _parse_dt(before_dict.get("DATETIME ARRIVAL") or "")
        rpt_dt = arr_dt + timedelta(days=1)
        b_code = _MONTH_B.get(arr_dt.month, 44)
        rr_line = f"Rectification Report No. RR/B-{b_code}/{arr_dt.year}/{rr_serial}"

        # weight_label is now ALWAYS derived from reason (UI dropdown removed)
        # Option A / B → "Refer Unladen Weight"
        # Option C     → "Out Weight from Weighbridge Indicator: X T"
        from rectification_report import REASON_C
        out_w_str = str(out_weight).strip() if out_weight else ""
        if reason == REASON_C:
            weight_label = (
                f"Out Weight from Weighbridge Indicator: {out_w_str} T"
                if out_w_str else "Out Weight from Weighbridge Indicator"
            )
        else:
            weight_label = "Refer Unladen Weight"

        # Generate table images
        tbl1_jpg, _ = fetch_and_generate(
            df=df, token=token, source="online",
            is_table2=False, override_values=None, dpi=300,
            force_outnet_yellow=True,
        )
        tbl2_jpg, _ = fetch_and_generate(
            df=df, token=token, source="online",
            is_table2=True, override_values=action_data, dpi=300,
        )

        # Build .docx
        docx_bytes, filename = build_rr_docx(
            token                = token,
            rr_serial            = rr_serial,
            rr_line              = rr_line,
            before_dict          = before_dict,
            after_dict           = after_dict,
            action_data          = action_data,
            arr_dt               = arr_dt,
            rpt_dt               = rpt_dt,
            tbl1_jpg             = tbl1_jpg,
            tbl2_jpg             = tbl2_jpg,
            excel_screenshot_jpg = None,
            weight_label         = weight_label,
            reason               = reason,
            filter_date_str      = filter_date,
        )

        log_activity(session["username"], "RECTIFICATION_REPORT", filename)
        return send_file(
            io.BytesIO(docx_bytes),
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name=filename,
        )

    except ValueError as ve:
        return jsonify({"error": str(ve)}), 400
    except Exception as e:
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


# ═══════════════════════════════════════════════════════════════════════════════
#  HTML TEMPLATES
# ═══════════════════════════════════════════════════════════════════════════════

AUTH_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>APSG (Staging Ground) Report — Sign In</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box;margin:0;padding:0;}
:root{
  --bg:#080C1A;--card-bg:rgba(13,18,35,0.92);
  --indigo:#6366F1;--indigo-l:#818CF8;--cyan:#22D3EE;
  --purple:#A855F7;--green:#10B981;--red:#F87171;
  --text:#E8EEF8;--muted:#64748B;--border:rgba(99,102,241,0.18);
}
body{font-family:'Inter',system-ui,sans-serif;min-height:100vh;
  display:flex;align-items:flex-start;justify-content:center;
  padding-top:5vh;overflow-y:auto;}
.bg-mesh{position:fixed;inset:0;z-index:0;pointer-events:none;
  background:
    radial-gradient(ellipse 70% 60% at 15% 15%,rgba(99,102,241,.12) 0%,transparent 65%),
    radial-gradient(ellipse 60% 50% at 85% 85%,rgba(168,85,247,.10) 0%,transparent 65%),
    radial-gradient(ellipse 40% 35% at 50% 50%,rgba(34,211,238,.05) 0%,transparent 65%);
}
.grid-lines{position:fixed;inset:0;z-index:0;pointer-events:none;opacity:.025;
  background-image:linear-gradient(var(--indigo) 1px,transparent 1px),
    linear-gradient(90deg,var(--indigo) 1px,transparent 1px);
  background-size:60px 60px;}
.wrap{width:100%;max-width:440px;padding:1.5rem;position:relative;z-index:1;padding-top:.5rem;}
.brand{text-align:center;margin-bottom:1.8rem;}
.brand-logo{display:inline-flex;align-items:center;justify-content:center;
  width:64px;height:64px;border-radius:18px;margin-bottom:1rem;
  background:linear-gradient(135deg,rgba(99,102,241,.25),rgba(168,85,247,.15));
  border:1px solid rgba(99,102,241,.3);font-size:28px;
  box-shadow:0 0 40px rgba(99,102,241,.2);}
.brand-title{font-family:'Poppins',sans-serif;font-size:1.45rem;font-weight:800;
  color:var(--text);letter-spacing:-.03em;line-height:1.2;}
.brand-sub{font-size:.73rem;color:var(--muted);margin-top:.3rem;font-weight:400;letter-spacing:.03em;}
.card{background:rgba(8,14,38,0.72);border:1px solid rgba(255,255,255,.10);
  border-radius:20px;padding:2rem 2.2rem;
  backdrop-filter:blur(20px);-webkit-backdrop-filter:blur(20px);
  box-shadow:0 8px 48px rgba(0,0,0,.45);}
.card-header{font-size:.95rem;font-weight:700;color:var(--text);
  text-align:center;margin-bottom:1.6rem;display:flex;
  align-items:center;justify-content:center;gap:.5rem;letter-spacing:-.01em;}
.field{margin-bottom:1rem;}
.field label{display:block;font-size:.67rem;font-weight:700;color:var(--muted);
  letter-spacing:.12em;text-transform:uppercase;margin-bottom:.4rem;}
.field input{width:100%;padding:.78rem 1rem;
  background:rgba(6,10,24,.8);border:1.5px solid rgba(30,41,86,.8);
  border-radius:12px;color:var(--text);font-size:.9rem;
  font-family:'Inter',sans-serif;transition:all .2s;}
.field input:focus{outline:none;border-color:var(--indigo);
  box-shadow:0 0 0 3px rgba(99,102,241,.18);background:rgba(10,14,32,.9);}
.field input::placeholder{color:rgba(100,116,139,.5);font-weight:300;}
.btn{width:100%;padding:.85rem;border-radius:12px;border:none;cursor:pointer;
  font-size:.88rem;font-weight:700;font-family:'Inter',sans-serif;
  letter-spacing:.01em;transition:all .22s;position:relative;overflow:hidden;}
.btn-primary{
  background:linear-gradient(135deg,#4338CA 0%,#6366F1 50%,#818CF8 100%);
  color:#fff;box-shadow:0 4px 20px rgba(99,102,241,.4);margin-top:.3rem;}
.btn-primary::before{content:'';position:absolute;inset:0;
  background:linear-gradient(135deg,transparent,rgba(255,255,255,.1),transparent);
  transform:translateX(-100%);transition:transform .5s;}
.btn-primary:hover{transform:translateY(-2px);
  box-shadow:0 8px 32px rgba(99,102,241,.6);}
.btn-primary:hover::before{transform:translateX(100%);}
.btn-primary:active{transform:translateY(0) scale(.98);}
.btn-outline{background:transparent;border:1.5px solid rgba(99,102,241,.3);
  color:var(--indigo-l);margin-top:.65rem;}
.btn-outline:hover{background:rgba(99,102,241,.1);border-color:rgba(99,102,241,.55);}
.divider{display:flex;align-items:center;gap:.75rem;margin:.8rem 0;}
.divider::before,.divider::after{content:'';flex:1;height:1px;background:rgba(99,102,241,.12);}
.divider span{font-size:.68rem;color:rgba(100,116,139,.6);white-space:nowrap;}
.alert{border-radius:10px;padding:.65rem .9rem;font-size:.78rem;
  margin-bottom:.9rem;display:none;font-weight:500;}
.alert.show{display:block;}
.alert-error{background:rgba(239,68,68,.08);border:1px solid rgba(239,68,68,.2);color:#F87171;}
.alert-success{background:rgba(16,185,129,.08);border:1px solid rgba(16,185,129,.2);color:#34D399;}
.footer-links{text-align:center;margin-top:1rem;font-size:.72rem;color:var(--muted);}
.footer-links a{color:var(--indigo-l);text-decoration:none;font-weight:600;}
.footer-links a:hover{color:var(--cyan);}
.tag{display:inline-block;background:rgba(99,102,241,.1);
  border:1px solid rgba(99,102,241,.2);border-radius:6px;
  padding:.15rem .5rem;font-size:.62rem;font-weight:600;
  color:var(--indigo-l);letter-spacing:.05em;margin-bottom:1.4rem;}



/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="wrap">
  <div class="brand">
    <div class="brand-logo">📊</div>
    <div class="brand-title">APSG (Staging Ground) Report</div>
    <div class="brand-sub">Staging Ground Report System · Phase 3</div>
  </div>
  <!-- Login Card -->
  <div class="card" id="loginCard">
    <div class="card-header">🔐 Sign In to Continue</div>
    <div class="alert alert-error" id="loginError"></div>
    <div class="alert alert-success" id="loginSuccess"></div>
    <div class="field">
      <label>Username</label>
      <input type="text" id="loginUser" placeholder="Enter your username" autocomplete="username" autofocus>
    </div>
    <div class="field">
      <label>Password</label>
      <input type="password" id="loginPass" placeholder="Enter your password" autocomplete="current-password">
    </div>
    <div style="display:flex;align-items:center;gap:.6rem;margin:.5rem 0 .8rem;">
        <input type="checkbox" id="rememberMe" style="width:16px;height:16px;accent-color:#6366F1;cursor:pointer;margin:0;flex-shrink:0;">
        <label for="rememberMe" style="font-size:13px;color:#A0BAD8;cursor:pointer;font-weight:500;margin:0;user-select:none;">Remember me</label>
      </div>
      <button class="btn btn-primary" onclick="doLogin()">Sign In →</button>
    <div class="divider"><span>New to the system?</span></div>
    <button class="btn btn-outline" onclick="showRegister()">✨ Register for New User</button>
    <div class="footer-links" style="margin-top:.9rem;">Forgot password? <a href="#">Contact Admin</a></div>
  </div>
  <!-- Register Card -->
  <div class="card" id="registerCard" style="display:none">
    <div class="card-header">✨ Register New Account</div>
    <div class="alert alert-error" id="regError"></div>
    <div class="alert alert-success" id="regSuccess"></div>
    <div class="field">
      <label>Full Name</label>
      <input type="text" id="regName" placeholder="Enter your full name">
    </div>
    <div class="field">
      <label>Username</label>
      <input type="text" id="regUser" placeholder="Choose a unique username">
    </div>
    <div class="field">
      <label>Password</label>
      <input type="password" id="regPass" placeholder="Minimum 6 characters">
    </div>
    <div class="field">
      <label>Confirm Password</label>
      <input type="password" id="regPass2" placeholder="Re-enter your password">
    </div>
    <button class="btn btn-primary" onclick="doRegister()">Create Account →</button>
    <div class="divider"><span>Already registered?</span></div>
    <button class="btn btn-outline" onclick="showLogin()">← Back to Sign In</button>
  </div>
  <div class="footer-links" style="margin-top:1rem;">Need help? <a href="mailto:admin@apsg-report.com">Contact Admin</a> &nbsp;·&nbsp; Developed by <strong>Karthi</strong></div>
</div>
<script>
function show(id){document.getElementById(id).style.display='block';}
function hide(id){document.getElementById(id).style.display='none';}
function showAlert(id,msg){const el=document.getElementById(id);el.textContent=msg;el.classList.add('show');}
function clearAlerts(ids){ids.forEach(id=>{const el=document.getElementById(id);el.classList.remove('show');el.textContent='';});}
function showRegister(){hide('loginCard');show('registerCard');clearAlerts(['loginError','loginSuccess','regError','regSuccess']);}
function showLogin(){hide('registerCard');show('loginCard');clearAlerts(['loginError','loginSuccess','regError','regSuccess']);}
async function doLogin(){
  const rememberMe = document.getElementById('rememberMe')?.checked;
  if(rememberMe){
    localStorage.setItem('apsg_remember_user', document.getElementById('loginUser').value);
  } else {
    localStorage.removeItem('apsg_remember_user');
  }
  clearAlerts(['loginError','loginSuccess']);
  const user=document.getElementById('loginUser').value.trim();
  const pass=document.getElementById('loginPass').value;
  if(!user||!pass){showAlert('loginError','Please enter username and password.');return;}
  try{
    const res=await fetch('/login',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({username:user,password:pass})});
    const d=await res.json();
    if(d.ok){window.location.href='/';}
    else{showAlert('loginError',d.error||'Invalid username or password.');}
  }catch(e){showAlert('loginError','Connection error. Please try again.');}
}
async function doRegister(){
  clearAlerts(['regError','regSuccess']);
  const name=document.getElementById('regName').value.trim();
  const user=document.getElementById('regUser').value.trim();
  const pass=document.getElementById('regPass').value;
  const pass2=document.getElementById('regPass2').value;
  if(!name||!user||!pass){showAlert('regError','All fields are required.');return;}
  if(pass!==pass2){showAlert('regError','Passwords do not match.');return;}
  if(pass.length<6){showAlert('regError','Password must be at least 6 characters.');return;}
  try{
    const res=await fetch('/register',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({name,username:user,password:pass,confirm:pass2})});
    const d=await res.json();
    if(d.ok){window.location.href='/';}
    else{showAlert('regError',d.error||'Registration failed.');}
  }catch(e){showAlert('regError','Connection error. Please try again.');}
}
document.addEventListener('keydown',e=>{if(e.key==='Enter'){if(document.getElementById('loginCard').style.display!=='none')doLogin();else doRegister();}});
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""


DASHBOARD_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Dashboard - APSG (Staging Ground) Report</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&family=Poppins:wght@600;700;800&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box;margin:0;padding:0;}
:root{
  --bg:#060912;
  --surface:#0D1120;
  --surface2:#111827;
  --indigo:#6366F1;
  --indigo-d:#4F46E5;
  --indigo-l:#A5B4FC;
  --cyan:#22D3EE;
  --purple:#A855F7;
  --green:#10B981;
  --amber:#F59E0B;
  --red:#F87171;
  --text:#F0F4FF;
  --text2:#CBD5E1;
  --muted:#64748B;
  --border:rgba(148,163,184,0.08);
  --border-h:rgba(148,163,184,0.18);
}
body{font-family:"Inter",system-ui,sans-serif;min-height:100vh;
  background:var(--bg);color:var(--text);font-size:14px;overflow-x:hidden;}

/* ── Subtle aurora — reduced opacity, less distraction ── */
.aurora{position:fixed;inset:0;z-index:0;pointer-events:none;overflow:hidden;}
.aurora-blob{position:absolute;border-radius:50%;filter:blur(110px);opacity:.07;animation:blobDrift 18s ease-in-out infinite;}
.aurora-blob:nth-child(1){width:800px;height:800px;background:#6366F1;top:-300px;left:-200px;animation-delay:0s;}
.aurora-blob:nth-child(2){width:600px;height:600px;background:#A855F7;bottom:-200px;right:-150px;animation-delay:-8s;}
.aurora-blob:nth-child(3){width:400px;height:400px;background:#22D3EE;top:40%;left:40%;animation-delay:-14s;}
@keyframes blobDrift{0%,100%{transform:translate(0,0);}50%{transform:translate(30px,-20px);}}

/* ── TOP BAR ── */
.top-bar{position:sticky;top:0;z-index:300;height:56px;display:flex;
  align-items:center;padding:0 2rem;gap:.75rem;
  background: rgba(4,8,22,0.50);backdrop-filter:blur(20px);
  border-bottom: 1px solid rgba(255,255,255,0.08);}
.top-brand{display:flex;align-items:center;gap:.6rem;}
.brand-mark{width:28px;height:28px;border-radius:8px;flex-shrink:0;
  background:linear-gradient(135deg,#4F46E5,#7C3AED);
  display:flex;align-items:center;justify-content:center;font-size:14px;
  box-shadow:0 0 0 1px rgba(99,102,241,.4);}
.brand-text{font-family:"Poppins",sans-serif;font-size:.82rem;font-weight:700;
  color:var(--text);letter-spacing:-.01em;}
.brand-text span{color:var(--indigo-l);}
.top-spacer{flex:1;}
.top-right{display:flex;align-items:center;gap:.4rem;}
.top-pill{border-radius:6px;padding:.28rem .8rem;font-size:.68rem;font-weight:600;
  text-decoration:none;transition:all .18s;border:1px solid transparent;
  white-space:nowrap;letter-spacing:.01em;}
.pill-user{background:rgba(99,102,241,.09);border-color:rgba(99,102,241,.2);color:var(--indigo-l);}
.pill-admin{background:rgba(245,158,11,.08);border-color:rgba(245,158,11,.2);color:var(--amber);}
.pill-logout{background:rgba(248,113,113,.06);border-color:rgba(248,113,113,.15);color:var(--red);}
.pill-admin:hover,.pill-logout:hover{opacity:.85;transform:translateY(-1px);}

/* ── PAGE ── */
.page{padding:2.2rem 1.8rem 5rem;max-width:1160px;margin:0 auto;position:relative;z-index:1;}

/* ── HERO ── */
.hero{text-align:center;padding:2rem 1rem 2.6rem;}
.hero-eyebrow{display:inline-flex;align-items:center;gap:.5rem;padding:.28rem .9rem;
  background:rgba(99,102,241,.07);border:1px solid rgba(99,102,241,.18);
  border-radius:4px;font-size:.66rem;color:var(--indigo-l);font-weight:600;
  letter-spacing:.12em;text-transform:uppercase;margin-bottom:1.2rem;}
.live-dot{width:6px;height:6px;border-radius:50%;background:#10B981;
  animation:livePulse 2s ease-in-out infinite;flex-shrink:0;}
@keyframes livePulse{0%,100%{box-shadow:0 0 0 0 rgba(16,185,129,.5);}60%{box-shadow:0 0 0 5px rgba(16,185,129,0);}}
.hero-title{font-family:"Poppins",sans-serif;
  font-size:clamp(1.75rem,4.5vw,2.8rem);font-weight:800;
  color:var(--text);line-height:1.1;margin-bottom:.6rem;letter-spacing:-.03em;}
.hero-sub{font-size:.84rem;color:var(--muted);margin-bottom:.45rem;}
.hero-user{font-size:.8rem;color:var(--indigo-l);font-weight:600;}

/* ── SECTION LABEL ── */
.section-label{font-size:.65rem;font-weight:700;letter-spacing:.14em;
  text-transform:uppercase;color:var(--muted);margin-bottom:1rem;
  display:flex;align-items:center;gap:.6rem;}
.section-label::after{content:"";flex:1;height:1px;background:var(--border);}

/* ── CARDS GRID ── */
.cards-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(300px,1fr));gap:1rem;}

.app-card{
  position:relative;display:block;text-decoration:none;color:inherit;
  background: rgba(8,14,40,0.72);
  border: 1px solid rgba(255,255,255,0.10);
  border-radius:14px;padding:1.6rem 1.5rem;
  transition:transform .22s ease,box-shadow .22s ease,border-color .22s ease,background .22s ease;
  cursor:pointer;overflow:hidden;}

/* Left accent bar — colour identifier, always visible */
.app-card::before{
  content:"";position:absolute;left:0;top:16px;bottom:16px;width:3px;
  border-radius:0 3px 3px 0;
  background:var(--accent,#6366F1);
  opacity:.7;transition:opacity .22s,top .22s,bottom .22s;}

/* Hover state — clear, readable, not over-glowing */
.app-card:hover{
  background:var(--surface2);
  border-color:var(--border-h);
  transform:translateY(-4px);
  box-shadow:0 12px 40px rgba(0,0,0,.45),0 0 0 1px var(--accent-border,rgba(99,102,241,.25));}
.app-card:hover::before{opacity:1;top:12px;bottom:12px;}

/* Active / click feedback */
.app-card:active{transform:translateY(-2px) scale(.99);transition-duration:.08s;}

/* Hover — title becomes full white */
.app-card:hover .card-title{color:#fff;}
/* Hover — description becomes clearly readable */
.app-card:hover .card-desc{color:var(--text2);}
/* Hover — number becomes accent colour */
.app-card:hover .card-num{color:var(--accent,#6366F1);opacity:1;}
/* Hover — arrow moves and brightens */
.app-card:hover .card-arrow{transform:translateX(5px);opacity:1;}

/* Card colour themes — only sets accent, no background changes */
.card-blue{--accent:#60A5FA;--accent-border:rgba(96,165,250,.22);}
.card-green{--accent:#34D399;--accent-border:rgba(52,211,153,.22);}
.card-purple{--accent:#C084FC;--accent-border:rgba(192,132,252,.22);}
.card-amber{--accent:#FBBF24;--accent-border:rgba(251,191,36,.22);}
.card-gray{--accent:#6B7280;--accent-border:rgba(107,114,128,.15);}

.app-card.disabled{cursor:not-allowed;opacity:.35;pointer-events:none;}

/* Content */
.card-num{font-size:.58rem;font-weight:700;color:var(--muted);
  letter-spacing:.16em;text-transform:uppercase;margin-bottom:.6rem;
  transition:color .22s;}
.card-icon-wrap{display:flex;align-items:flex-start;justify-content:space-between;margin-bottom:.9rem;}
.card-icon{width:44px;height:44px;border-radius:10px;display:flex;
  align-items:center;justify-content:center;font-size:20px;flex-shrink:0;
  background:rgba(255,255,255,.04);border:1px solid var(--border);
  transition:background .22s,transform .22s;}
.app-card:hover .card-icon{background:rgba(255,255,255,.07);transform:scale(1.06);}
.card-arrow{color:var(--muted);font-size:1rem;transition:transform .22s,opacity .22s;opacity:.5;}

.card-title{font-family:"Poppins",sans-serif;font-size:.93rem;font-weight:700;
  color:var(--text);margin-bottom:.4rem;line-height:1.35;
  transition:color .22s;}
.card-desc{font-size:.76rem;color:var(--muted);line-height:1.7;
  transition:color .22s;}
.card-badge{display:inline-flex;align-items:center;gap:.28rem;
  padding:.18rem .65rem;border-radius:4px;
  font-size:.61rem;font-weight:700;margin-top:.9rem;letter-spacing:.05em;}
.badge-active{background:rgba(16,185,129,.09);color:#34D399;border:1px solid rgba(16,185,129,.2);}
.badge-soon{background:rgba(107,114,128,.07);color:var(--muted);border:1px solid var(--border);}

/* Card entry animation */
.app-card{animation:slideIn .4s cubic-bezier(.16,1,.3,1) both;}
.app-card:nth-child(1){animation-delay:.04s;}
.app-card:nth-child(2){animation-delay:.09s;}
.app-card:nth-child(3){animation-delay:.14s;}
.app-card:nth-child(4){animation-delay:.19s;}
.app-card:nth-child(5){animation-delay:.24s;}
@keyframes slideIn{from{opacity:0;transform:translateY(20px);}to{opacity:1;transform:none;}}

/* ── FOOTER ── */
.page-footer{text-align:center;padding:2rem 0 .5rem;margin-top:2.5rem;
  border-top:1px solid var(--border);}

.dev-credit{
  position:relative;display:inline-flex;align-items:center;gap:.55rem;
  background:var(--surface);border:1px solid var(--border);
  border-radius:30px;padding:.38rem 1rem .38rem .45rem;cursor:default;
  font-size:.72rem;font-weight:600;color:var(--text2);}
.dev-credit:hover{border-color:var(--border-h);}

/* Avatar */
.karthi-avatar{
  width:28px;height:28px;border-radius:50%;flex-shrink:0;
  background:linear-gradient(135deg,#4F46E5,#7C3AED,#DB2777);
  display:flex;align-items:center;justify-content:center;font-size:14px;
  animation:gentleWave 3s ease-in-out infinite;}
@keyframes gentleWave{0%,100%{transform:rotate(0deg);}35%{transform:rotate(-8deg);}65%{transform:rotate(8deg);}}

/* Tooltip on hover */
.dev-credit .tooltip{
  position:absolute;bottom:calc(100% + 10px);left:50%;transform:translateX(-50%) translateY(6px);
  background:#1E293B;border:1px solid rgba(99,102,241,.25);
  border-radius:10px;padding:.7rem 1rem;width:230px;
  font-size:.72rem;font-weight:400;color:var(--text2);line-height:1.55;
  opacity:0;pointer-events:none;transition:opacity .2s,transform .2s;
  text-align:center;white-space:normal;}
.dev-credit .tooltip::after{
  content:"";position:absolute;top:100%;left:50%;transform:translateX(-50%);
  border:6px solid transparent;border-top-color:#1E293B;}
.dev-credit:hover .tooltip{opacity:1;transform:translateX(-50%) translateY(0);}
.tooltip-quote{font-style:italic;color:var(--indigo-l);margin-top:.3rem;font-size:.7rem;}

.footer-meta{font-size:.61rem;color:rgba(100,116,139,.4);margin-top:.55rem;letter-spacing:.04em;}

/* ── WAKE SCREEN ── */
.wake{position:fixed;inset:0;z-index:9999;display:flex;flex-direction:column;
  align-items:center;justify-content:center;gap:1rem;background: rgba(4,8,22,0.85);}
.wake-spinner{width:36px;height:36px;border:2px solid rgba(99,102,241,.15);
  border-top-color:var(--indigo);border-radius:50%;animation:spin 1s linear infinite;}
@keyframes spin{to{transform:rotate(360deg);}}
.wake-label{font-family:"Poppins",sans-serif;font-size:.9rem;font-weight:700;color:var(--text2);}
.wake-sub{font-size:.72rem;color:var(--muted);}

/* ── RESPONSIVE ── */
@media(max-width:600px){
  .cards-grid{grid-template-columns:1fr;}
  .top-bar{padding:0 1rem;}
  .page{padding:1.5rem 1rem 4rem;}
}



/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>


<!-- Wake Screen -->
<div class="wake" id="wakeScreen">
  <div class="wake-spinner"></div>
  <div class="wake-label">APSG Report</div>
  <div class="wake-sub" id="wakeMsg">Loading&hellip;</div>
</div>

<!-- Top Bar -->
<div class="top-bar">
  <div class="top-brand">
    <div class="brand-mark">&#9889;</div>
    <span class="brand-text">APSG <span>Report</span></span>
  </div>
  <div class="top-spacer"></div>
  <div class="top-right">
    <a href="/admin" id="adminBtn" class="top-pill pill-admin" style="display:none;">&#9881; Admin</a>
    <span class="top-pill pill-user">&#128100; {{ name }}</span>
    <a href="/logout" class="top-pill pill-logout">Sign Out</a>
  </div>
</div>

<!-- Main -->
<div class="page">

  <!-- Hero -->
  <div class="hero">
    <div class="hero-eyebrow"><span class="live-dot"></span>Internal Reporting Platform</div>
    <div class="hero-title">APSG (Staging Ground) Report</div>
    <div class="hero-sub">Staging Ground Report System &middot; Phase 3</div>
    <div class="hero-user">Welcome back, {{ name }}</div>
  </div>

  <!-- Cards -->
  <div class="section-label">Select a report to get started</div>

  <div class="cards-grid">

    <a href="/app/daily-report" class="app-card card-blue">
      <div class="card-num">01</div>
      <div class="card-icon-wrap">
        <div class="card-icon">&#128202;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Daily Report</div>
      <div class="card-desc">Generate daily staging ground reports with Online vs WB comparison, validation, and Excel download.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/app/excel-rejection" class="app-card card-green">
      <div class="card-num">02</div>
      <div class="card-icon-wrap">
        <div class="card-icon">&#128203;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Excel Rejection / Rejection Filter</div>
      <div class="card-desc">Convert Excel data into formatted PowerPoint rejection reports grouped by source site.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/app/photo-merge" class="app-card card-purple">
      <div class="card-num">03</div>
      <div class="card-icon-wrap">
        <div class="card-icon">&#128444;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">PPT Alignment</div>
      <div class="card-desc">Merge Top Photo and Front Photo PPT files with auto token matching and layout control.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/app/ppt-rejection" class="app-card card-amber">
      <div class="card-num">04</div>
      <div class="card-icon-wrap">
        <div class="card-icon">&#128193;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Bulk Bundle Report Filter</div>
      <div class="card-desc">Upload bulk PPT files, filter by E-Token/date/reason, and generate grouped rejection reports with ZIP export.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <div class="app-card card-gray disabled">
      <div class="card-num">05</div>
      <div class="card-icon-wrap">
        <div class="card-icon">&#8987;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Cycle Time Report</div>
      <div class="card-desc">Cycle time analysis and reporting &mdash; in development. Coming soon.</div>
      <span class="card-badge badge-soon">&#8987; In Progress</span>
    </div>

  </div>

  <!-- Footer -->
  <div class="page-footer">
    <div class="dev-credit">
      <div class="karthi-avatar">&#129489;</div>
      Developed by&nbsp;<strong>Karthi</strong>
      <div class="tooltip">
        Hi &#128075; I built this tool to make your reporting faster and smarter.
        <div class="tooltip-quote">&ldquo;Data tells the truth &mdash; if you listen properly.&rdquo;</div>
      </div>
    </div>
    <div class="footer-meta">APSG (Staging Ground) Report &middot; v1.0 &middot; Internal Use Only</div>
  </div>

</div><!-- /page -->

<script>
const role="{{ role }}";
if(role==="admin") document.getElementById("adminBtn").style.display="inline-flex";

window.addEventListener("load", () => {
  const ws = document.getElementById("wakeScreen");
  const wm = document.getElementById("wakeMsg");
  fetch("/api/health").then(r => {
    wm.textContent = r.ok ? "Ready!" : "Server warming up...";
    setTimeout(() => {
      ws.style.opacity = "0";
      ws.style.transition = "opacity .5s";
      setTimeout(() => ws.style.display = "none", 520);
    }, r.ok ? 300 : 1800);
  }).catch(() => setTimeout(() => ws.style.display = "none", 1500));
});
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""


PPT_REJECTION_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>PPT Alignment — APSG (Staging Ground) Report</title>
<style>

/* ═══ MODERN UI BASE (Blue/Purple/Cyan Theme) ═══ */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
:root {
  --bg:transparent; --card-bg: rgba(13,18,35,0.88);
  --indigo: #6366F1; --indigo-l: #818CF8; --cyan: #22D3EE;
  --purple: #A855F7; --green: #10B981; --amber: #F59E0B; --red: #F87171;
  --text: #E8EEF8; --muted: #64748B; --border: rgba(99,102,241,0.15);
}
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent; color: var(--text);
  position: relative;
}
body::before {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none;
  background:
    radial-gradient(ellipse 70% 55% at 5% 10%, rgba(99,102,241,.09) 0%, transparent 65%),
    radial-gradient(ellipse 55% 45% at 95% 90%, rgba(168,85,247,.07) 0%, transparent 65%);
}
body::after {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none; opacity: .016;
  background-image: linear-gradient(var(--indigo) 1px, transparent 1px),
    linear-gradient(90deg, var(--indigo) 1px, transparent 1px);
  background-size: 60px 60px;
}

/* ── Modern Top Bar ── */
.top-bar {
  position: sticky; top: 0; z-index: 200; height: 56px;
  display: flex; align-items: center; padding: 0 1.5rem; gap: .75rem;
  background: rgba(6,9,22,.96); backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px);
  border-bottom: 1px solid var(--border);
}
.top-mini-brand {
  font-family: 'Poppins', sans-serif;
  font-size: .82rem; font-weight: 700; color: var(--text);
  letter-spacing: -.01em; white-space: nowrap; flex-shrink: 0;
  display: flex; align-items: center; gap: .45rem;
}
.top-mini-brand::before {
  content: ''; width: 8px; height: 8px; border-radius: 50%;
  background: linear-gradient(135deg, var(--indigo), var(--cyan)); flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: var(--border); flex-shrink: 0; }
.top-page-label { font-size: .75rem; font-weight: 600; color: var(--muted); white-space: nowrap; }
.top-brand-tag { font-size: .7rem; color: var(--muted); }
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.08); border: 1px solid rgba(99,102,241,.18);
  border-radius: 8px; padding: .3rem .9rem; font-size: .7rem; font-weight: 600;
  color: var(--indigo-l); text-decoration: none; transition: all .2s; white-space: nowrap;
}
.back-btn:hover { background: rgba(99,102,241,.18); transform: translateX(-2px); }

/* ── Cards & Containers ── */
.container, .page-content { position: relative; z-index: 1; }
.card, .section-card, .panel {
  background: var(--card-bg); border: 1px solid var(--border);
  border-radius: 16px; backdrop-filter: blur(16px); -webkit-backdrop-filter: blur(16px);
  box-shadow: 0 8px 32px rgba(0,0,0,.4), inset 0 1px 0 rgba(255,255,255,.03);
  transition: box-shadow .25s, border-color .25s;
}
.card:hover, .section-card:hover { border-color: rgba(99,102,241,.25); }

/* ── Upload Zone — Modern drag & drop ── */
.upload-zone, .dz {
  border: 2px dashed rgba(99,102,241,.3); border-radius: 14px;
  padding: 2rem; text-align: center; cursor: pointer;
  transition: all .22s; background: rgba(99,102,241,.03);
  position: relative;
}
.upload-zone:hover, .dz:hover, .upload-zone.drag-over, .dz.drag-over {
  border-color: var(--indigo); background: rgba(99,102,241,.08);
  box-shadow: 0 0 0 4px rgba(99,102,241,.12);
}
.upload-zone.ok, .dz.ok {
  border-color: var(--green); background: rgba(16,185,129,.06);
  border-style: solid;
}
.upload-zone.ok:hover, .dz.ok:hover {
  border-color: var(--green); background: rgba(16,185,129,.1);
  box-shadow: 0 0 0 4px rgba(16,185,129,.1);
}
.upload-icon { font-size: 2rem; margin-bottom: .5rem; display: block; }
.upload-label { font-size: .82rem; color: var(--muted); font-weight: 500; }
.upload-hint { font-size: .7rem; color: rgba(100,116,139,.6); margin-top: .25rem; }
.upload-filename { font-size: .78rem; color: var(--green); font-weight: 600; margin-top: .4rem; }

/* ── Modern Buttons ── */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  background:transparent;
  color: #fff; border: none; border-radius: 10px;
  padding: .7rem 1.4rem; font-size: .85rem; font-weight: 700;
  font-family: 'Inter', sans-serif; cursor: pointer; letter-spacing: .01em;
  box-shadow: 0 4px 18px rgba(99,102,241,.35);
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s;
}
.btn-primary::before, .btn-generate::before, .btn-teal::before, .modal-btn-primary::before {
  content: ''; position: absolute; inset: 0;
  background: linear-gradient(135deg, transparent, rgba(255,255,255,.12), transparent);
  transform: translateX(-100%); transition: transform .45s;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-2px) scale(1.01);
  box-shadow: 0 8px 28px rgba(99,102,241,.55);
}
.btn-primary:hover::before, .btn-generate:hover::before, .modal-btn-primary:hover::before {
  transform: translateX(100%);
}
.btn-primary:active, .btn-generate:active { transform: translateY(0) scale(.98); }
.btn-primary:disabled, .btn-generate:disabled { opacity: .45; cursor: not-allowed; transform: none; }

.btn-teal {
  background:transparent;
  box-shadow: 0 4px 18px rgba(16,185,129,.3);
}
.btn-teal:hover { box-shadow: 0 8px 28px rgba(16,185,129,.5); }

.btn-download, .btn-dl, .dl-btn {
  background: linear-gradient(135deg, #0D7A5F, #10B981);
  color: #fff; border: none; border-radius: 10px;
  padding: .65rem 1.3rem; font-size: .82rem; font-weight: 700;
  cursor: pointer; transition: all .22s;
  box-shadow: 0 4px 16px rgba(16,185,129,.3);
}
.btn-download:hover, .btn-dl:hover, .dl-btn:hover {
  transform: translateY(-2px); box-shadow: 0 8px 26px rgba(16,185,129,.5);
}

.btn-secondary, .btn-gray {
  background: rgba(30,41,86,.6); color: var(--muted);
  border: 1px solid rgba(99,102,241,.2); border-radius: 10px;
  padding: .65rem 1.2rem; font-size: .82rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-secondary:hover, .btn-gray:hover { background: rgba(99,102,241,.1); color: var(--indigo-l); }

.btn-del, .btn-danger {
  background: rgba(239,68,68,.08); color: var(--red);
  border: 1px solid rgba(239,68,68,.2); border-radius: 10px;
  padding: .6rem 1.2rem; font-size: .8rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-del:hover, .btn-danger:hover { background: rgba(239,68,68,.16); }
.btn-del:hover { animation: shake .3s ease; }

@keyframes shake {
  0%,100% { transform: translateX(0); }
  25% { transform: translateX(-3px); }
  75% { transform: translateX(3px); }
}

/* ── Form Inputs ── */
input[type=text], input[type=number], input[type=date], input[type=password],
select, textarea {
  background: rgba(6,10,24,.8); border: 1.5px solid rgba(30,41,86,.8);
  border-radius: 10px; color: var(--text); padding: .65rem .9rem;
  font-size: .85rem; font-family: 'Inter', sans-serif;
  transition: all .2s; width: 100%;
}
input:focus, select:focus, textarea:focus {
  outline: none; border-color: var(--indigo);
  box-shadow: 0 0 0 3px rgba(99,102,241,.16);
  background: rgba(10,14,32,.9);
}
input::placeholder, textarea::placeholder { color: rgba(100,116,139,.45); font-weight: 300; }

/* ── Alerts / Status ── */
.alert-success, .alert.success, .msg.success {
  background: rgba(16,185,129,.08); border: 1px solid rgba(16,185,129,.2);
  color: #34D399; border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-error, .alert.error, .msg.error {
  background: rgba(239,68,68,.08); border: 1px solid rgba(239,68,68,.2);
  color: var(--red); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-warn, .alert.warn { 
  background: rgba(245,158,11,.08); border: 1px solid rgba(245,158,11,.2);
  color: var(--amber); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}

/* ── Tables ── */
table { width: 100%; border-collapse: collapse; font-size: .8rem; }
thead th {
  background: rgba(99,102,241,.08); color: var(--muted);
  font-weight: 700; font-size: .68rem; letter-spacing: .06em;
  text-transform: uppercase; padding: .65rem .9rem; text-align: left;
  border-bottom: 1px solid var(--border);
}
tbody tr { border-bottom: 1px solid rgba(30,45,80,.3); transition: background .15s; }
tbody tr:hover { background: rgba(99,102,241,.04); }
tbody td { padding: .6rem .9rem; color: var(--text); }

/* ── Tabs ── */
.tab { position: relative; transition: all .2s; }
.tab.active { color: var(--indigo-l) !important; }
.tab.active::after {
  content: ''; position: absolute; bottom: -1px; left: 15%; right: 15%;
  height: 2px; border-radius: 2px; background: var(--indigo);
  animation: tabIn .2s ease;
}
@keyframes tabIn { from { left: 50%; right: 50%; } to { left: 15%; right: 15%; } }

/* ── Spinner ── */
@keyframes spin { to { transform: rotate(360deg); } }
.spinner { width: 20px; height: 20px; border: 2px solid rgba(99,102,241,.2);
  border-top-color: var(--indigo); border-radius: 50%; animation: spin 1s linear infinite; }



@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent;
  color: #E8EEF8;
  position: relative;
}


/* ══ ANIMATED BUTTONS ══════════════════════════════════════════════════════════ */
/* Primary action button — pulse + shimmer */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-primary::after, .btn-generate::after, .btn-teal::after, .modal-btn-primary::after {
  content: ''; position: absolute; top: -50%; left: -75%;
  width: 50%; height: 200%; background: rgba(255,255,255,.15);
  transform: skewX(-20deg); transition: left .5s ease;
  pointer-events: none;
}
.btn-primary:hover::after, .btn-generate:hover::after,
.btn-teal:hover::after, .modal-btn-primary:hover::after {
  left: 150%;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(99,102,241,.55) !important;
}
.btn-primary:active, .btn-generate:active, .modal-btn-primary:active {
  transform: translateY(0) scale(.98) !important;
}

/* Download button — bounce */
.btn-dl, .btn-download, .dl-btn {
  animation: gentlePulse 3s ease-in-out infinite;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-dl:hover, .btn-download:hover, .dl-btn:hover {
  animation: none;
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(16,185,129,.5) !important;
}
@keyframes gentlePulse {
  0%,100% { box-shadow: 0 4px 16px rgba(16,185,129,.3); }
  50%      { box-shadow: 0 4px 24px rgba(16,185,129,.55); }
}

/* Danger/delete button — shake on hover */
.btn-del:hover {
  animation: shake .35s ease;
}
@keyframes shake {
  0%,100% { transform: translateX(0); }
  20%     { transform: translateX(-3px); }
  40%     { transform: translateX(3px); }
  60%     { transform: translateX(-2px); }
  80%     { transform: translateX(2px); }
}

/* Secondary / back buttons — slide arrow */
.back-btn, .btn-secondary, .btn-gray, .btn-refresh {
  transition: all .2s !important;
}
.back-btn:hover {
  padding-left: .55rem !important;
  letter-spacing: .02em;
}

/* Upload zone — glow pulse when empty */
.upload-zone:not(.ok) {
  animation: uploadGlow 4s ease-in-out infinite;
}
.upload-zone:hover, .dz:hover {
  animation: none;
}
@keyframes uploadGlow {
  0%,100% { border-color: rgba(99,102,241,.3); }
  50%      { border-color: rgba(99,102,241,.65); box-shadow: 0 0 18px rgba(99,102,241,.15); }
}

/* App cards on dashboard — float in on load */
.app-card {
  animation: cardFadeIn .5s ease both;
}
.app-card:nth-child(1) { animation-delay: .05s; }
.app-card:nth-child(2) { animation-delay: .10s; }
.app-card:nth-child(3) { animation-delay: .15s; }
.app-card:nth-child(4) { animation-delay: .20s; }
.app-card:nth-child(5) { animation-delay: .25s; }
@keyframes cardFadeIn {
  from { opacity: 0; transform: translateY(18px); }
  to   { opacity: 1; transform: translateY(0); }
}

/* Stat cards pop-in */
.stat-card, .stat {
  animation: popIn .4s cubic-bezier(.34,1.56,.64,1) both;
}
.stat-card:nth-child(1), .stat:nth-child(1) { animation-delay: .05s; }
.stat-card:nth-child(2), .stat:nth-child(2) { animation-delay: .10s; }
.stat-card:nth-child(3), .stat:nth-child(3) { animation-delay: .15s; }
.stat-card:nth-child(4), .stat:nth-child(4) { animation-delay: .20s; }
@keyframes popIn {
  from { opacity: 0; transform: scale(.88); }
  to   { opacity: 1; transform: scale(1); }
}

/* Login button — ripple */
.btn.btn-primary, button.btn-primary {
  transition: transform .15s, box-shadow .15s !important;
}

/* Tab buttons — active indicator slide */
.tab {
  transition: all .2s !important;
  position: relative;
}
.tab.active::after {
  content: '';
  position: absolute; bottom: -1px; left: 10%; right: 10%;
  height: 2px; border-radius: 2px;
  background: #6366F1;
  animation: tabSlide .25s ease;
}
@keyframes tabSlide {
  from { left: 50%; right: 50%; }
  to   { left: 10%; right: 10%; }
}

/* Generate/action buttons in forms — shimmer idle */
.btn-generate:not(:disabled), .btn-teal:not(:disabled) {
  background-size: 200% 100% !important;
  transition: background-position .4s, transform .2s, box-shadow .2s !important;
}

/* ── Compact top bar (40px) for all internal pages ── */
.top-bar {
  position: sticky; top: 0; z-index: 200;
  height: 40px; display: flex; align-items: center;
  padding: 0 1.2rem; gap: .75rem;
  background: rgba(5,8,20,.96); backdrop-filter: blur(16px);
  border-bottom: 1px solid rgba(99,102,241,.12);
}
.top-mini-brand {
  font-size: .75rem; font-weight: 800; color: #6366F1;
  letter-spacing: .04em; white-space: nowrap; flex-shrink: 0;
}
.top-brand-text {
  font-size: .78rem; font-weight: 800; color: #818CF8;
  letter-spacing: .02em; white-space: nowrap; flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: rgba(99,102,241,.18); flex-shrink: 0; }
.top-page-label {
  font-size: .78rem; font-weight: 600; color: #94A3B8;
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
}
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.1); border: 1px solid rgba(99,102,241,.2);
  border-radius: 6px; padding: .22rem .7rem; font-size: .7rem; font-weight: 600;
  color: #818CF8; text-decoration: none; white-space: nowrap; transition: all .2s;
  flex-shrink: 0;
}
.back-btn:hover { background: rgba(99,102,241,.2); }

.container, .main, .page { padding: 1rem 1.2rem 2rem; max-width: 1400px; margin: 0 auto; position: relative; z-index: 1; }


/* PPT Rejection specific */
:root{--g:#00B050;--mid:#0F172A;--light:#E8EEF8;--muted:#6B7280;--danger:#ef4444;--warn:#f59e0b;--info:#3B82F6;}
.card{background:rgba(15,23,42,.8);border:1px solid rgba(99,102,241,.12);border-radius:12px;padding:1.2rem;margin-bottom:1rem;backdrop-filter:blur(10px);}
.card h2{font-size:.82rem;font-weight:700;color:#818CF8;margin-bottom:.85rem;display:flex;align-items:center;gap:.5rem;}
.upload-zone{border:2px dashed rgba(99,102,241,.3);border-radius:9px;padding:2rem;text-align:center;cursor:pointer;transition:.2s;position:relative;}
.upload-zone:hover,.upload-zone.dragover{border-color:#6366F1;background:rgba(99,102,241,.06);}
.upload-zone input{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%;}
.upload-zone .icon{font-size:2rem;margin-bottom:.4rem;}
.upload-zone p{color:var(--muted);font-size:.82rem;}
.file-list{margin-top:.6rem;display:flex;flex-wrap:wrap;gap:.35rem;}
.file-tag{background:rgba(99,102,241,.12);border:1px solid rgba(99,102,241,.25);color:#818CF8;padding:.2rem .6rem;border-radius:20px;font-size:.72rem;display:flex;align-items:center;gap:.3rem;}
.file-tag button{background:none;border:none;color:#F87171;cursor:pointer;font-size:.85rem;line-height:1;padding:0;}
.status-bar{border-radius:7px;padding:.6rem .85rem;font-size:.79rem;display:none;align-items:center;gap:.5rem;margin-top:.6rem;}
.status-bar.show{display:flex;}
.status-bar.ok{background:rgba(16,185,129,.1);border:1px solid rgba(16,185,129,.2);color:#10B981;}
.status-bar.error{background:rgba(239,68,68,.1);border:1px solid rgba(239,68,68,.3);color:#F87171;}
.status-bar.warn{background:rgba(245,158,11,.1);border:1px solid rgba(245,158,11,.3);color:#F59E0B;}
.prog-wrap{margin-top:.5rem;display:none;}.prog-wrap.show{display:block;}
.prog-bar{height:5px;border-radius:3px;background:rgba(255,255,255,.08);overflow:hidden;}
.prog-fill{height:100%;background:#6366F1;width:0;transition:width .3s;border-radius:3px;}
.prog-lbl{font-size:.72rem;color:var(--muted);margin-top:.25rem;}
.filter-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(160px,1fr));gap:.65rem;}
.field label{display:block;font-size:.72rem;color:var(--muted);margin-bottom:.3rem;font-weight:500;}
.field input,.field select{width:100%;padding:.5rem .75rem;background:rgba(8,15,40,.7);border:1px solid rgba(99,102,241,.15);border-radius:7px;color:var(--light);font-size:.82rem;transition:.2s;font-family:'Inter',sans-serif;}
.field input:focus,.field select:focus{outline:none;border-color:#6366F1;}
.field input::placeholder{color:var(--muted);}select option{background:#1a1a2e;}
.btn-row{display:flex;gap:.6rem;flex-wrap:wrap;margin-top:.6rem;align-items:center;}
.btn{padding:.5rem 1.1rem;border-radius:7px;border:none;cursor:pointer;font-size:.8rem;font-weight:600;display:flex;align-items:center;gap:.4rem;transition:.2s;font-family:'Inter',sans-serif;}
.btn-primary{background:#6366F1;color:#fff;}.btn-primary:hover{background:#4F46E5;transform:translateY(-1px);}
.btn-secondary{background:rgba(255,255,255,.06);color:var(--light);border:1px solid rgba(255,255,255,.1);}
.btn-secondary:hover{background:rgba(255,255,255,.1);}
.btn-dl{background:#3B82F6;color:#fff;}.btn-dl:hover{background:#2563EB;transform:translateY(-1px);}
.btn-xl{background:#7C3AED;color:#fff;}.btn-xl:hover{background:#6D28D9;transform:translateY(-1px);}
.btn-zip{background:#D97706;color:#fff;}.btn-zip:hover{background:#B45309;transform:translateY(-1px);}
.btn:disabled{opacity:.4;cursor:not-allowed;transform:none!important;}
.stats-row{display:flex;gap:.65rem;flex-wrap:wrap;margin-bottom:1rem;}
.stat{flex:1;min-width:80px;background:rgba(99,102,241,.08);border:1px solid rgba(99,102,241,.15);border-radius:9px;padding:.75rem;text-align:center;}
.stat .val{font-size:1.6rem;font-weight:800;color:#818CF8;}
.stat .lbl{font-size:.65rem;color:var(--muted);margin-top:.15rem;}
.result-meta{display:flex;align-items:center;gap:.65rem;margin-bottom:.5rem;flex-wrap:wrap;}
.count-badge{background:rgba(99,102,241,.12);border:1px solid rgba(99,102,241,.25);color:#818CF8;padding:.2rem .7rem;border-radius:20px;font-size:.76rem;font-weight:600;}
.tbl-wrap{overflow-x:auto;border-radius:7px;border:1px solid rgba(99,102,241,.1);}
table{width:100%;border-collapse:collapse;font-size:.76rem;}
thead{background:rgba(99,102,241,.1);}
thead th{padding:.55rem .7rem;text-align:center;font-weight:600;color:#818CF8;white-space:nowrap;border-bottom:1px solid rgba(99,102,241,.18);}
tbody tr{border-bottom:1px solid rgba(255,255,255,.03);}
tbody tr:hover{background:rgba(255,255,255,.03);}
tbody td{padding:.45rem .7rem;text-align:center;color:var(--light);}
tbody td.left{text-align:left;}
tr.group-header td{background:rgba(99,102,241,.1);color:#818CF8;font-weight:700;font-size:.72rem;text-align:left;padding:.4rem .9rem;}
.badge-no{display:inline-block;padding:.12rem .48rem;border-radius:10px;font-size:.68rem;font-weight:600;background:rgba(239,68,68,.15);color:#F87171;}
.badge-reason{display:inline-block;padding:.12rem .48rem;border-radius:10px;font-size:.68rem;font-weight:600;background:rgba(245,158,11,.12);color:#F59E0B;}
.no-data{text-align:center;padding:2rem;color:var(--muted);}

/* ══ ANIMATED BUTTONS ══════════════════════════════════════════════════════════ */
/* Primary action button — pulse + shimmer */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-primary::after, .btn-generate::after, .btn-teal::after, .modal-btn-primary::after {
  content: ''; position: absolute; top: -50%; left: -75%;
  width: 50%; height: 200%; background: rgba(255,255,255,.15);
  transform: skewX(-20deg); transition: left .5s ease;
  pointer-events: none;
}
.btn-primary:hover::after, .btn-generate:hover::after,
.btn-teal:hover::after, .modal-btn-primary:hover::after {
  left: 150%;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(99,102,241,.55) !important;
}
.btn-primary:active, .btn-generate:active, .modal-btn-primary:active {
  transform: translateY(0) scale(.98) !important;
}

/* Download button — bounce */
.btn-dl, .btn-download, .dl-btn {
  animation: gentlePulse 3s ease-in-out infinite;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-dl:hover, .btn-download:hover, .dl-btn:hover {
  animation: none;
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(16,185,129,.5) !important;
}
@keyframes gentlePulse {
  0%,100% { box-shadow: 0 4px 16px rgba(16,185,129,.3); }
  50%      { box-shadow: 0 4px 24px rgba(16,185,129,.55); }
}

/* Danger/delete button — shake on hover */
.btn-del:hover {
  animation: shake .35s ease;
}
@keyframes shake {
  0%,100% { transform: translateX(0); }
  20%     { transform: translateX(-3px); }
  40%     { transform: translateX(3px); }
  60%     { transform: translateX(-2px); }
  80%     { transform: translateX(2px); }
}

/* Secondary / back buttons — slide arrow */
.back-btn, .btn-secondary, .btn-gray, .btn-refresh {
  transition: all .2s !important;
}
.back-btn:hover {
  padding-left: .55rem !important;
  letter-spacing: .02em;
}

/* Upload zone — glow pulse when empty */
.upload-zone:not(.ok) {
  animation: uploadGlow 4s ease-in-out infinite;
}
.upload-zone:hover, .dz:hover {
  animation: none;
}
@keyframes uploadGlow {
  0%,100% { border-color: rgba(99,102,241,.3); }
  50%      { border-color: rgba(99,102,241,.65); box-shadow: 0 0 18px rgba(99,102,241,.15); }
}

/* App cards on dashboard — float in on load */
.app-card {
  animation: cardFadeIn .5s ease both;
}
.app-card:nth-child(1) { animation-delay: .05s; }
.app-card:nth-child(2) { animation-delay: .10s; }
.app-card:nth-child(3) { animation-delay: .15s; }
.app-card:nth-child(4) { animation-delay: .20s; }
.app-card:nth-child(5) { animation-delay: .25s; }
@keyframes cardFadeIn {
  from { opacity: 0; transform: translateY(18px); }
  to   { opacity: 1; transform: translateY(0); }
}

/* Stat cards pop-in */
.stat-card, .stat {
  animation: popIn .4s cubic-bezier(.34,1.56,.64,1) both;
}
.stat-card:nth-child(1), .stat:nth-child(1) { animation-delay: .05s; }
.stat-card:nth-child(2), .stat:nth-child(2) { animation-delay: .10s; }
.stat-card:nth-child(3), .stat:nth-child(3) { animation-delay: .15s; }
.stat-card:nth-child(4), .stat:nth-child(4) { animation-delay: .20s; }
@keyframes popIn {
  from { opacity: 0; transform: scale(.88); }
  to   { opacity: 1; transform: scale(1); }
}

/* Login button — ripple */
.btn.btn-primary, button.btn-primary {
  transition: transform .15s, box-shadow .15s !important;
}

/* Tab buttons — active indicator slide */
.tab {
  transition: all .2s !important;
  position: relative;
}
.tab.active::after {
  content: '';
  position: absolute; bottom: -1px; left: 10%; right: 10%;
  height: 2px; border-radius: 2px;
  background: #6366F1;
  animation: tabSlide .25s ease;
}
@keyframes tabSlide {
  from { left: 50%; right: 50%; }
  to   { left: 10%; right: 10%; }
}

/* Generate/action buttons in forms — shimmer idle */
.btn-generate:not(:disabled), .btn-teal:not(:disabled) {
  background-size: 200% 100% !important;
  transition: background-position .4s, transform .2s, box-shadow .2s !important;
}

.spinner{animation:spin 1s linear infinite;display:inline-block;}
@keyframes spin{to{transform:rotate(360deg);}}
/* ── Password management button styles ── */
.btn-pw{background:rgba(99,102,241,.12);color:#818CF8;border:1px solid rgba(99,102,241,.25);}
.btn-pw:hover{background:rgba(99,102,241,.22);}
.btn-reset{background:rgba(245,158,11,.1);color:#F59E0B;border:1px solid rgba(245,158,11,.22);}
.btn-reset:hover{background:rgba(245,158,11,.18);}
.btn-del{background:rgba(239,68,68,.08);color:#F87171;border:1px solid rgba(239,68,68,.18);}
.btn-del:hover{background:rgba(239,68,68,.16);}
/* ── Modal backdrop + card ── */
.modal-overlay{position:fixed;inset:0;z-index:9000;background:rgba(0,0,0,.65);
  display:none;align-items:center;justify-content:center;backdrop-filter:blur(4px);}
.modal-card{background:#0D1529;border:1px solid rgba(99,102,241,.25);border-radius:16px;
  padding:1.8rem 2rem;width:100%;max-width:420px;box-shadow:0 24px 60px rgba(0,0,0,.6);}
.modal-title{font-size:1rem;font-weight:800;color:#F1F5FF;margin-bottom:.25rem;
  display:flex;align-items:center;gap:.5rem;}
.modal-sub{font-size:.78rem;color:#475569;margin-bottom:1.3rem;}
.modal-field{margin-bottom:.9rem;}
.modal-field label{display:block;font-size:.68rem;font-weight:700;color:#475569;
  letter-spacing:.08em;text-transform:uppercase;margin-bottom:.35rem;}
.modal-field input{width:100%;padding:.72rem .9rem;background:rgba(8,15,40,.9);
  border:1.5px solid rgba(30,41,86,.9);border-radius:10px;color:#E8EEF8;
  font-size:.9rem;font-family:'Inter',sans-serif;transition:all .2s;}
.modal-field input:focus{outline:none;border-color:#6366F1;box-shadow:0 0 0 3px rgba(99,102,241,.18);}
.modal-field input::placeholder{color:#334155;}
.show-pw-row{display:flex;align-items:center;gap:.5rem;margin-bottom:1rem;
  font-size:.76rem;color:#475569;cursor:pointer;}
.show-pw-row input[type=checkbox]{width:14px;height:14px;accent-color:#6366F1;cursor:pointer;}
.modal-actions{display:flex;gap:.6rem;margin-top:.5rem;}
.modal-btn{flex:1;padding:.75rem;border-radius:10px;border:none;cursor:pointer;
  font-size:.88rem;font-weight:700;font-family:'Inter',sans-serif;transition:all .2s;}
.modal-btn-primary{background:linear-gradient(135deg,#4338CA,#6366F1);color:#fff;
  box-shadow:0 4px 16px rgba(99,102,241,.35);}
.modal-btn-primary:hover{transform:translateY(-1px);box-shadow:0 6px 20px rgba(99,102,241,.5);}
.modal-btn-cancel{background:rgba(255,255,255,.06);color:#94A3B8;
  border:1px solid rgba(255,255,255,.1);}
.modal-btn-cancel:hover{background:rgba(255,255,255,.1);}
.modal-result{padding:.55rem .8rem;border-radius:8px;font-size:.8rem;font-weight:500;
  margin-top:.75rem;display:none;}
.modal-result.result-ok{display:block;background:rgba(16,185,129,.1);
  border:1px solid rgba(16,185,129,.22);color:#34D399;}
.modal-result.result-err{display:block;background:rgba(239,68,68,.1);
  border:1px solid rgba(239,68,68,.22);color:#F87171;}
/* Temp PW display */
.temp-pw-box{background:rgba(8,15,40,.9);border:1.5px solid rgba(99,102,241,.3);
  border-radius:10px;padding:1rem 1.2rem;margin:.8rem 0;text-align:center;}
.temp-pw-label{font-size:.68rem;color:#475569;font-weight:600;letter-spacing:.08em;
  text-transform:uppercase;margin-bottom:.5rem;}
.temp-pw-value{font-size:1.6rem;font-weight:900;color:#818CF8;letter-spacing:.12em;
  font-family:'Courier New',monospace;}
.temp-pw-warn{font-size:.72rem;color:#F59E0B;margin-top:.6rem;line-height:1.5;}




/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="top-bar">
  <span class="top-mini-brand">APSG</span>
  <div class="top-sep"></div>
  <span class="top-page-label">📁 Bulk Bundle Report Filter</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>

<div class="container" style="padding-top:.75rem; position:relative; z-index:1; max-width:1400px; margin:0 auto;">
  <div class="card">
    <h2>📁 Upload PPT Files</h2>
    <div class="upload-zone" id="dropZone">
      <input type="file" id="fileInput" multiple accept=".ppt,.pptx">
      <div class="icon">📊</div>
      <strong>Drop files here or click to browse</strong>
      <p>Supports multiple .PPTX files — any quantity</p>
    </div>
    <div class="file-list" id="fileList"></div>
    <div class="prog-wrap" id="progWrap">
      <div class="prog-bar"><div class="prog-fill" id="progFill"></div></div>
      <div class="prog-lbl" id="progLbl"></div>
    </div>
    <div class="status-bar" id="uploadStatus"></div>
    <div class="btn-row">
      <button class="btn btn-primary" id="uploadBtn" onclick="uploadFiles()">⬆ Upload &amp; Extract</button>
      <button class="btn btn-secondary" onclick="clearAll()">🗑 Clear All</button>
    </div>
  </div>
  <div class="stats-row" id="statsRow" style="display:none">
    <div class="stat"><div class="val" id="sTotal">0</div><div class="lbl">Total Records</div></div>
    <div class="stat"><div class="val" id="sFiltered">0</div><div class="lbl">Filtered</div></div>
    <div class="stat"><div class="val" id="sVehicles">0</div><div class="lbl">Vehicles</div></div>
    <div class="stat"><div class="val" id="sSites">0</div><div class="lbl">Sites</div></div>
    <div class="stat"><div class="val" id="sFiles">0</div><div class="lbl">Files</div></div>
  </div>
  <div class="card">
    <h2>🔍 Search &amp; Filter</h2>
    <div class="filter-grid">
      <div class="field"><label>E-Token Search</label>
        <input type="text" id="searchKw" placeholder="Search E-Token…" oninput="onSearchInput()"></div>
      <div class="field"><label>E-Token</label>
        <select id="selEToken" onchange="onETokenChange()"><option value="">All</option></select></div>
      <div class="field"><label>Reject Reason</label>
        <select id="selReason" onchange="applyFilters()"><option value="">All</option></select></div>
      <div class="field"><label>Date From</label>
        <input type="date" id="dateFrom" onchange="applyFilters()" style="color-scheme:dark;"></div>
      <div class="field"><label>Date To</label>
        <input type="date" id="dateTo" onchange="applyFilters()" style="color-scheme:dark;"></div>
    </div>
    <div class="btn-row">
      <button class="btn btn-secondary" onclick="resetFilters()">↺ Reset</button>
      <button class="btn btn-primary" onclick="applyFilters()">🔎 Apply</button>
    </div>
  </div>
  <div class="card">
    <h2>📋 Preview &amp; Export</h2>
    <div class="result-meta" id="resultMeta" style="display:none">
      <span class="count-badge" id="vehicleCountBadge">0 Vehicles</span>
      <span class="count-badge" id="recordCountBadge">0 Records</span>
      <span class="count-badge" id="siteCountBadge">0 Sites</span>
    </div>
    <div class="tbl-wrap">
      <table>
        <thead><tr>
          <th>#</th><th>Ticket No</th><th>Veh No</th><th>Material</th>
          <th>Source Site</th><th>Date</th><th>Time</th>
          <th>E-Token</th><th>Accepted</th><th>Reject Reason</th>
        </tr></thead>
        <tbody id="previewBody">
          <tr><td colspan="10" class="no-data">Upload PPT files to see data</td></tr>
        </tbody>
      </table>
    </div>
    <div class="btn-row" style="margin-top:1rem">
      <button class="btn btn-dl" id="dlBtn" onclick="downloadReport()" disabled>⬇ Download PPT</button>
      <button class="btn btn-zip" id="zipBtn" onclick="downloadZip()" disabled>📦 Download ZIP</button>
      <button class="btn btn-xl" id="xlBtn" onclick="downloadExcel()" disabled>📊 Excel</button>
      <span id="dlStatus" style="font-size:.78rem;color:var(--muted);"></span>
    </div>
  </div>
</div>
<script>
let allFiles=[],allRecords=[],filtered=[],totalLoaded=0;
const fileInput=document.getElementById('fileInput');
const dropZone=document.getElementById('dropZone');
fileInput.addEventListener('change',()=>addFiles(fileInput.files));
dropZone.addEventListener('dragover',e=>{e.preventDefault();dropZone.classList.add('dragover');});
dropZone.addEventListener('dragleave',()=>dropZone.classList.remove('dragover'));
dropZone.addEventListener('drop',e=>{e.preventDefault();dropZone.classList.remove('dragover');addFiles(e.dataTransfer.files);});
function addFiles(files){for(const f of files)if(f.name.match(/\.pptx?$/i)&&!allFiles.find(x=>x.name===f.name&&x.size===f.size))allFiles.push(f);renderFileList();}
function renderFileList(){document.getElementById('fileList').innerHTML=allFiles.map((f,i)=>`<div class="file-tag">📄 ${esc(f.name)} <small>(${(f.size/1024).toFixed(0)}KB)</small><button onclick="removeFile(${i})">×</button></div>`).join('');}
function removeFile(i){allFiles.splice(i,1);renderFileList();}
function clearAll(){allFiles=[];allRecords=[];filtered=[];totalLoaded=0;renderFileList();renderTable([]);setStats(0,0,0,0,0);document.getElementById('statsRow').style.display='none';document.getElementById('resultMeta').style.display='none';setStatus('','');['dlBtn','zipBtn','xlBtn'].forEach(id=>document.getElementById(id).disabled=true);}
async function uploadFiles(){
  if(!allFiles.length){setStatus('Please select at least one PPT file.','error');return;}
  const btn=document.getElementById('uploadBtn');
  btn.disabled=true;btn.innerHTML='<span class="spinner">⟳</span> Processing…';
  setStatus('','');allRecords=[];filtered=[];totalLoaded=0;
  const BATCH=5,errs=[];
  try{
    for(let i=0;i<allFiles.length;i+=BATCH){
      const batch=allFiles.slice(i,i+BATCH);
      showProg(true,Math.round((i/allFiles.length)*90),`Processing ${i+1}–${Math.min(i+BATCH,allFiles.length)} of ${allFiles.length}…`);
      const fd=new FormData();
      batch.forEach(f=>fd.append('files',f));
      let res,data;
      try{
        res=await fetch(`/api/ppt/upload?reset=${i===0}`,{method:'POST',body:fd});
        data=await res.json();
      }catch(ne){errs.push(`Network: ${ne.message}`);continue;}
      if(data.error&&!(data.records&&data.records.length)){errs.push(data.error);continue;}
      if(data.records)allRecords=data.records;
      if(data.files_loaded)totalLoaded+=data.files_loaded;
      if(data.errors&&data.errors.length)errs.push(...data.errors);
    }
    showProg(true,100,'Done!');
    filtered=[...allRecords];
    populateDropdowns(allRecords);renderTable(filtered);updateMeta(filtered);
    const v=countUnique(filtered,'veh_no'),s=countUnique(filtered,'source_site');
    setStats(allRecords.length,filtered.length,v,s,totalLoaded);
    document.getElementById('statsRow').style.display='flex';
    ['dlBtn','zipBtn','xlBtn'].forEach(id=>document.getElementById(id).disabled=allRecords.length===0);
    let msg=`✓ Extracted ${allRecords.length} records from ${totalLoaded} file(s).`;
    let type='ok';
    if(errs.length){msg+=` ⚠ ${errs.length} issue(s).`;type='warn';}
    setStatus(msg,type);
  }catch(e){setStatus('Upload error: '+e.message,'error');}
  finally{btn.disabled=false;btn.innerHTML='⬆ Upload &amp; Extract';setTimeout(()=>showProg(false),2000);}
}
function countUnique(recs,key){return new Set(recs.map(r=>r[key]).filter(Boolean)).size;}
function onSearchInput(){const kw=document.getElementById('searchKw').value.trim();const esel=document.getElementById('selEToken');if(kw.length>=1){const m=[...new Set(allRecords.filter(r=>r.e_token.toLowerCase().includes(kw.toLowerCase())).map(r=>r.e_token))];rebuildSel(esel,m);}else{rebuildSel(esel,[...new Set(allRecords.map(r=>r.e_token))].filter(Boolean));}onETokenChange();}
function onETokenChange(){const et=document.getElementById('selEToken').value;const reasons=[...new Set(allRecords.filter(r=>!et||r.e_token===et).map(r=>r.reject_reason))].filter(Boolean);rebuildSel(document.getElementById('selReason'),reasons);applyFilters();}
function applyFilters(){const kw=document.getElementById('searchKw').value.trim().toLowerCase();const et=document.getElementById('selEToken').value;const rs=document.getElementById('selReason').value;const df=pd(document.getElementById('dateFrom').value);const dt_=pd(document.getElementById('dateTo').value);filtered=allRecords.filter(r=>{if(kw.length>=3&&!r.e_token.toLowerCase().includes(kw))return false;if(et&&r.e_token!==et)return false;if(rs&&r.reject_reason!==rs)return false;if(df||dt_){const d=pd(r.date);if(!d)return false;if(df&&d<df)return false;if(dt_&&d>dt_)return false;}return true;});filtered.sort((a,b)=>a.source_site.localeCompare(b.source_site)||a.e_token.localeCompare(b.e_token));renderTable(filtered);updateMeta(filtered);const v=countUnique(filtered,'veh_no'),s=countUnique(filtered,'source_site');setStats(allRecords.length,filtered.length,v,s,totalLoaded);['dlBtn','zipBtn','xlBtn'].forEach(id=>document.getElementById(id).disabled=filtered.length===0);}
function resetFilters(){document.getElementById('searchKw').value='';document.getElementById('dateFrom').value='';document.getElementById('dateTo').value='';populateDropdowns(allRecords);filtered=[...allRecords];filtered.sort((a,b)=>a.source_site.localeCompare(b.source_site)||a.e_token.localeCompare(b.e_token));renderTable(filtered);updateMeta(filtered);const v=countUnique(filtered,'veh_no'),s=countUnique(filtered,'source_site');setStats(allRecords.length,filtered.length,v,s,totalLoaded);}
function populateDropdowns(recs){rebuildSel(document.getElementById('selEToken'),[...new Set(recs.map(r=>r.e_token))].filter(Boolean));rebuildSel(document.getElementById('selReason'),[...new Set(recs.map(r=>r.reject_reason))].filter(Boolean));}
function rebuildSel(sel,opts){const cur=sel.value;sel.innerHTML='<option value="">All</option>'+opts.map(o=>`<option value="${esc(o)}"${o===cur?' selected':''}>${esc(o)}</option>`).join('');}
function renderTable(recs){const tb=document.getElementById('previewBody');if(!recs.length){tb.innerHTML='<tr><td colspan="10" class="no-data">No records match current filters</td></tr>';return;}let html='',lastSite='',sn=0;for(const r of recs){if(r.source_site!==lastSite){const cnt=recs.filter(x=>x.source_site===r.source_site).length;html+=`<tr class="group-header"><td colspan="10">📍 ${esc(r.source_site)} | ${cnt} record(s)</td></tr>`;lastSite=r.source_site;sn=0;}sn++;html+=`<tr><td>${sn}</td><td>${esc(r.ticket_no)}</td><td>${esc(r.veh_no)}</td><td>${esc(r.material)}</td><td>${esc(r.source_site)}</td><td>${esc(r.date)}</td><td>${esc(r.time)}</td><td style="font-size:.7rem">${esc(r.e_token)}</td><td><span class="badge-no">${esc(r.accepted||'NO')}</span></td><td class="left"><span class="badge-reason">${esc(r.reject_reason)}</span></td></tr>`;}tb.innerHTML=html;}
function updateMeta(recs){const show=recs.length>0;document.getElementById('resultMeta').style.display=show?'flex':'none';if(show){const v=countUnique(recs,'veh_no'),s=countUnique(recs,'source_site');document.getElementById('vehicleCountBadge').textContent=`🚛 ${v} Vehicle${v!==1?'s':''}`;document.getElementById('recordCountBadge').textContent=`📋 ${recs.length} Record${recs.length!==1?'s':''}`;document.getElementById('siteCountBadge').textContent=`📍 ${s} Source Site${s!==1?'s':''}`; }}
async function downloadReport(){if(!filtered.length)return;await doDownload('/api/ppt/generate','Rejection_Report.pptx','application/vnd.openxmlformats-officedocument.presentationml.presentation','⬇ Download PPT');}
async function downloadExcel(){if(!filtered.length)return;await doDownload('/api/ppt/export_excel','Rejection_Report.xlsx','application/vnd.openxmlformats-officedocument.spreadsheetml.sheet','📊 Excel');}
async function downloadZip(){if(!filtered.length)return;await doDownload('/api/ppt/generate_zip','Rejection_Report.zip','application/zip','📦 Download ZIP');}
async function doDownload(url,defaultName,mime,btnLabel){const btn=url.includes('zip')?document.getElementById('zipBtn'):url.includes('generate')?document.getElementById('dlBtn'):document.getElementById('xlBtn');const st=document.getElementById('dlStatus');btn.disabled=true;btn.innerHTML=`<span class="spinner">⟳</span> Generating…`;st.textContent='Building…';try{const res=await fetch(url,{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({records:filtered})});if(!res.ok){const e=await res.json().catch(()=>({error:`HTTP ${res.status}`}));st.textContent='Error: '+(e.error||'failed');return;}const blob=await res.blob();const dl=URL.createObjectURL(blob);const a=document.createElement('a');a.href=dl;a.download=defaultName;a.click();URL.revokeObjectURL(dl);st.textContent=`✓ Downloaded (${filtered.length} records)`;} catch(e){st.textContent='Error: '+e.message;}finally{btn.disabled=filtered.length===0;btn.innerHTML=btnLabel;}}
function setStats(total,filt,v,s,files){document.getElementById('sTotal').textContent=total;document.getElementById('sFiltered').textContent=filt;document.getElementById('sVehicles').textContent=v;document.getElementById('sSites').textContent=s;document.getElementById('sFiles').textContent=files;}
function setStatus(msg,type){const el=document.getElementById('uploadStatus');el.textContent=msg;el.className='status-bar'+(msg?' show':'')+(type?' '+type:'');}
function showProg(show,pct,lbl){document.getElementById('progWrap').className='prog-wrap'+(show?' show':'');if(show){document.getElementById('progFill').style.width=(pct||0)+'%';document.getElementById('progLbl').textContent=lbl||'';}}
function pd(str){if(!str)return null;const iso=str.match(/^(\d{4})-(\d{2})-(\d{2})$/);if(iso)return new Date(+iso[1],+iso[2]-1,+iso[3]);const dmy=str.match(/^(\d{1,2})-(\d{1,2})-(\d{4})$/);return dmy?new Date(+dmy[3],+dmy[2]-1,+dmy[1]):null;}
function esc(s){return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""

EXCEL_REJECTION_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Excel Rejection Report — APSG (Staging Ground) Report</title>
<style>

/* ═══ MODERN UI BASE (Blue/Purple/Cyan Theme) ═══ */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
:root {
  --bg:transparent; --card-bg: rgba(13,18,35,0.88);
  --indigo: #6366F1; --indigo-l: #818CF8; --cyan: #22D3EE;
  --purple: #A855F7; --green: #10B981; --amber: #F59E0B; --red: #F87171;
  --text: #E8EEF8; --muted: #64748B; --border: rgba(99,102,241,0.15);
}
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent; color: var(--text);
  position: relative;
}
body::before {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none;
  background:
    radial-gradient(ellipse 70% 55% at 5% 10%, rgba(99,102,241,.09) 0%, transparent 65%),
    radial-gradient(ellipse 55% 45% at 95% 90%, rgba(168,85,247,.07) 0%, transparent 65%);
}
body::after {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none; opacity: .016;
  background-image: linear-gradient(var(--indigo) 1px, transparent 1px),
    linear-gradient(90deg, var(--indigo) 1px, transparent 1px);
  background-size: 60px 60px;
}

/* ── Modern Top Bar ── */
.top-bar {
  position: sticky; top: 0; z-index: 200; height: 56px;
  display: flex; align-items: center; padding: 0 1.5rem; gap: .75rem;
  background: rgba(6,9,22,.96); backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px);
  border-bottom: 1px solid var(--border);
}
.top-mini-brand {
  font-family: 'Poppins', sans-serif;
  font-size: .82rem; font-weight: 700; color: var(--text);
  letter-spacing: -.01em; white-space: nowrap; flex-shrink: 0;
  display: flex; align-items: center; gap: .45rem;
}
.top-mini-brand::before {
  content: ''; width: 8px; height: 8px; border-radius: 50%;
  background: linear-gradient(135deg, var(--indigo), var(--cyan)); flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: var(--border); flex-shrink: 0; }
.top-page-label { font-size: .75rem; font-weight: 600; color: var(--muted); white-space: nowrap; }
.top-brand-tag { font-size: .7rem; color: var(--muted); }
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.08); border: 1px solid rgba(99,102,241,.18);
  border-radius: 8px; padding: .3rem .9rem; font-size: .7rem; font-weight: 600;
  color: var(--indigo-l); text-decoration: none; transition: all .2s; white-space: nowrap;
}
.back-btn:hover { background: rgba(99,102,241,.18); transform: translateX(-2px); }

/* ── Cards & Containers ── */
.container, .page-content { position: relative; z-index: 1; }
.card, .section-card, .panel {
  background: var(--card-bg); border: 1px solid var(--border);
  border-radius: 16px; backdrop-filter: blur(16px); -webkit-backdrop-filter: blur(16px);
  box-shadow: 0 8px 32px rgba(0,0,0,.4), inset 0 1px 0 rgba(255,255,255,.03);
  transition: box-shadow .25s, border-color .25s;
}
.card:hover, .section-card:hover { border-color: rgba(99,102,241,.25); }

/* ── Upload Zone — Modern drag & drop ── */
.upload-zone, .dz {
  border: 2px dashed rgba(99,102,241,.3); border-radius: 14px;
  padding: 2rem; text-align: center; cursor: pointer;
  transition: all .22s; background: rgba(99,102,241,.03);
  position: relative;
}
.upload-zone:hover, .dz:hover, .upload-zone.drag-over, .dz.drag-over {
  border-color: var(--indigo); background: rgba(99,102,241,.08);
  box-shadow: 0 0 0 4px rgba(99,102,241,.12);
}
.upload-zone.ok, .dz.ok {
  border-color: var(--green); background: rgba(16,185,129,.06);
  border-style: solid;
}
.upload-zone.ok:hover, .dz.ok:hover {
  border-color: var(--green); background: rgba(16,185,129,.1);
  box-shadow: 0 0 0 4px rgba(16,185,129,.1);
}
.upload-icon { font-size: 2rem; margin-bottom: .5rem; display: block; }
.upload-label { font-size: .82rem; color: var(--muted); font-weight: 500; }
.upload-hint { font-size: .7rem; color: rgba(100,116,139,.6); margin-top: .25rem; }
.upload-filename { font-size: .78rem; color: var(--green); font-weight: 600; margin-top: .4rem; }

/* ── Modern Buttons ── */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  background:transparent;
  color: #fff; border: none; border-radius: 10px;
  padding: .7rem 1.4rem; font-size: .85rem; font-weight: 700;
  font-family: 'Inter', sans-serif; cursor: pointer; letter-spacing: .01em;
  box-shadow: 0 4px 18px rgba(99,102,241,.35);
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s;
}
.btn-primary::before, .btn-generate::before, .btn-teal::before, .modal-btn-primary::before {
  content: ''; position: absolute; inset: 0;
  background: linear-gradient(135deg, transparent, rgba(255,255,255,.12), transparent);
  transform: translateX(-100%); transition: transform .45s;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-2px) scale(1.01);
  box-shadow: 0 8px 28px rgba(99,102,241,.55);
}
.btn-primary:hover::before, .btn-generate:hover::before, .modal-btn-primary:hover::before {
  transform: translateX(100%);
}
.btn-primary:active, .btn-generate:active { transform: translateY(0) scale(.98); }
.btn-primary:disabled, .btn-generate:disabled { opacity: .45; cursor: not-allowed; transform: none; }

.btn-teal {
  background:transparent;
  box-shadow: 0 4px 18px rgba(16,185,129,.3);
}
.btn-teal:hover { box-shadow: 0 8px 28px rgba(16,185,129,.5); }

.btn-download, .btn-dl, .dl-btn {
  background: linear-gradient(135deg, #0D7A5F, #10B981);
  color: #fff; border: none; border-radius: 10px;
  padding: .65rem 1.3rem; font-size: .82rem; font-weight: 700;
  cursor: pointer; transition: all .22s;
  box-shadow: 0 4px 16px rgba(16,185,129,.3);
}
.btn-download:hover, .btn-dl:hover, .dl-btn:hover {
  transform: translateY(-2px); box-shadow: 0 8px 26px rgba(16,185,129,.5);
}

.btn-secondary, .btn-gray {
  background: rgba(30,41,86,.6); color: var(--muted);
  border: 1px solid rgba(99,102,241,.2); border-radius: 10px;
  padding: .65rem 1.2rem; font-size: .82rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-secondary:hover, .btn-gray:hover { background: rgba(99,102,241,.1); color: var(--indigo-l); }

.btn-del, .btn-danger {
  background: rgba(239,68,68,.08); color: var(--red);
  border: 1px solid rgba(239,68,68,.2); border-radius: 10px;
  padding: .6rem 1.2rem; font-size: .8rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-del:hover, .btn-danger:hover { background: rgba(239,68,68,.16); }
.btn-del:hover { animation: shake .3s ease; }

@keyframes shake {
  0%,100% { transform: translateX(0); }
  25% { transform: translateX(-3px); }
  75% { transform: translateX(3px); }
}

/* ── Form Inputs ── */
input[type=text], input[type=number], input[type=date], input[type=password],
select, textarea {
  background: rgba(6,10,24,.8); border: 1.5px solid rgba(30,41,86,.8);
  border-radius: 10px; color: var(--text); padding: .65rem .9rem;
  font-size: .85rem; font-family: 'Inter', sans-serif;
  transition: all .2s; width: 100%;
}
input:focus, select:focus, textarea:focus {
  outline: none; border-color: var(--indigo);
  box-shadow: 0 0 0 3px rgba(99,102,241,.16);
  background: rgba(10,14,32,.9);
}
input::placeholder, textarea::placeholder { color: rgba(100,116,139,.45); font-weight: 300; }

/* ── Alerts / Status ── */
.alert-success, .alert.success, .msg.success {
  background: rgba(16,185,129,.08); border: 1px solid rgba(16,185,129,.2);
  color: #34D399; border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-error, .alert.error, .msg.error {
  background: rgba(239,68,68,.08); border: 1px solid rgba(239,68,68,.2);
  color: var(--red); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-warn, .alert.warn { 
  background: rgba(245,158,11,.08); border: 1px solid rgba(245,158,11,.2);
  color: var(--amber); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}

/* ── Tables ── */
table { width: 100%; border-collapse: collapse; font-size: .8rem; }
thead th {
  background: rgba(99,102,241,.08); color: var(--muted);
  font-weight: 700; font-size: .68rem; letter-spacing: .06em;
  text-transform: uppercase; padding: .65rem .9rem; text-align: left;
  border-bottom: 1px solid var(--border);
}
tbody tr { border-bottom: 1px solid rgba(30,45,80,.3); transition: background .15s; }
tbody tr:hover { background: rgba(99,102,241,.04); }
tbody td { padding: .6rem .9rem; color: var(--text); }

/* ── Tabs ── */
.tab { position: relative; transition: all .2s; }
.tab.active { color: var(--indigo-l) !important; }
.tab.active::after {
  content: ''; position: absolute; bottom: -1px; left: 15%; right: 15%;
  height: 2px; border-radius: 2px; background: var(--indigo);
  animation: tabIn .2s ease;
}
@keyframes tabIn { from { left: 50%; right: 50%; } to { left: 15%; right: 15%; } }

/* ── Spinner ── */
@keyframes spin { to { transform: rotate(360deg); } }
.spinner { width: 20px; height: 20px; border: 2px solid rgba(99,102,241,.2);
  border-top-color: var(--indigo); border-radius: 50%; animation: spin 1s linear infinite; }



@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent;
  color: #E8EEF8;
  position: relative;
}


/* ══ ANIMATED BUTTONS ══════════════════════════════════════════════════════════ */
/* Primary action button — pulse + shimmer */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-primary::after, .btn-generate::after, .btn-teal::after, .modal-btn-primary::after {
  content: ''; position: absolute; top: -50%; left: -75%;
  width: 50%; height: 200%; background: rgba(255,255,255,.15);
  transform: skewX(-20deg); transition: left .5s ease;
  pointer-events: none;
}
.btn-primary:hover::after, .btn-generate:hover::after,
.btn-teal:hover::after, .modal-btn-primary:hover::after {
  left: 150%;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(99,102,241,.55) !important;
}
.btn-primary:active, .btn-generate:active, .modal-btn-primary:active {
  transform: translateY(0) scale(.98) !important;
}

/* Download button — bounce */
.btn-dl, .btn-download, .dl-btn {
  animation: gentlePulse 3s ease-in-out infinite;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-dl:hover, .btn-download:hover, .dl-btn:hover {
  animation: none;
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(16,185,129,.5) !important;
}
@keyframes gentlePulse {
  0%,100% { box-shadow: 0 4px 16px rgba(16,185,129,.3); }
  50%      { box-shadow: 0 4px 24px rgba(16,185,129,.55); }
}

/* Danger/delete button — shake on hover */
.btn-del:hover {
  animation: shake .35s ease;
}
@keyframes shake {
  0%,100% { transform: translateX(0); }
  20%     { transform: translateX(-3px); }
  40%     { transform: translateX(3px); }
  60%     { transform: translateX(-2px); }
  80%     { transform: translateX(2px); }
}

/* Secondary / back buttons — slide arrow */
.back-btn, .btn-secondary, .btn-gray, .btn-refresh {
  transition: all .2s !important;
}
.back-btn:hover {
  padding-left: .55rem !important;
  letter-spacing: .02em;
}

/* Upload zone — glow pulse when empty */
.upload-zone:not(.ok) {
  animation: uploadGlow 4s ease-in-out infinite;
}
.upload-zone:hover, .dz:hover {
  animation: none;
}
@keyframes uploadGlow {
  0%,100% { border-color: rgba(99,102,241,.3); }
  50%      { border-color: rgba(99,102,241,.65); box-shadow: 0 0 18px rgba(99,102,241,.15); }
}

/* App cards on dashboard — float in on load */
.app-card {
  animation: cardFadeIn .5s ease both;
}
.app-card:nth-child(1) { animation-delay: .05s; }
.app-card:nth-child(2) { animation-delay: .10s; }
.app-card:nth-child(3) { animation-delay: .15s; }
.app-card:nth-child(4) { animation-delay: .20s; }
.app-card:nth-child(5) { animation-delay: .25s; }
@keyframes cardFadeIn {
  from { opacity: 0; transform: translateY(18px); }
  to   { opacity: 1; transform: translateY(0); }
}

/* Stat cards pop-in */
.stat-card, .stat {
  animation: popIn .4s cubic-bezier(.34,1.56,.64,1) both;
}
.stat-card:nth-child(1), .stat:nth-child(1) { animation-delay: .05s; }
.stat-card:nth-child(2), .stat:nth-child(2) { animation-delay: .10s; }
.stat-card:nth-child(3), .stat:nth-child(3) { animation-delay: .15s; }
.stat-card:nth-child(4), .stat:nth-child(4) { animation-delay: .20s; }
@keyframes popIn {
  from { opacity: 0; transform: scale(.88); }
  to   { opacity: 1; transform: scale(1); }
}

/* Login button — ripple */
.btn.btn-primary, button.btn-primary {
  transition: transform .15s, box-shadow .15s !important;
}

/* Tab buttons — active indicator slide */
.tab {
  transition: all .2s !important;
  position: relative;
}
.tab.active::after {
  content: '';
  position: absolute; bottom: -1px; left: 10%; right: 10%;
  height: 2px; border-radius: 2px;
  background: #6366F1;
  animation: tabSlide .25s ease;
}
@keyframes tabSlide {
  from { left: 50%; right: 50%; }
  to   { left: 10%; right: 10%; }
}

/* Generate/action buttons in forms — shimmer idle */
.btn-generate:not(:disabled), .btn-teal:not(:disabled) {
  background-size: 200% 100% !important;
  transition: background-position .4s, transform .2s, box-shadow .2s !important;
}

/* ── Compact top bar (40px) for all internal pages ── */
.top-bar {
  position: sticky; top: 0; z-index: 200;
  height: 40px; display: flex; align-items: center;
  padding: 0 1.2rem; gap: .75rem;
  background: rgba(5,8,20,.96); backdrop-filter: blur(16px);
  border-bottom: 1px solid rgba(99,102,241,.12);
}
.top-mini-brand {
  font-size: .75rem; font-weight: 800; color: #6366F1;
  letter-spacing: .04em; white-space: nowrap; flex-shrink: 0;
}
.top-brand-text {
  font-size: .78rem; font-weight: 800; color: #818CF8;
  letter-spacing: .02em; white-space: nowrap; flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: rgba(99,102,241,.18); flex-shrink: 0; }
.top-page-label {
  font-size: .78rem; font-weight: 600; color: #94A3B8;
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
}
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.1); border: 1px solid rgba(99,102,241,.2);
  border-radius: 6px; padding: .22rem .7rem; font-size: .7rem; font-weight: 600;
  color: #818CF8; text-decoration: none; white-space: nowrap; transition: all .2s;
  flex-shrink: 0;
}
.back-btn:hover { background: rgba(99,102,241,.2); }

.container, .main, .page { padding: 1rem 1.2rem 2rem; max-width: 1400px; margin: 0 auto; }


.main{padding:.75rem 1.1rem 2rem;max-width:900px;margin:0 auto;position:relative;z-index:1;}
.hero{text-align:center;padding:1.2rem 1rem 1.4rem;}
.hero-icon{font-size:2.2rem;display:block;margin-bottom:.4rem;}
.hero-title{font-size:1.6rem;font-weight:900;color:#F1F5FF;letter-spacing:-.03em;text-shadow:0 0 60px rgba(99,102,241,.5);}
.hero-sub{font-size:.78rem;color:#374167;margin:.3rem 0;}
.hero-pill{display:inline-block;padding:.18rem .8rem;background:rgba(99,102,241,.1);border:1px solid rgba(99,102,241,.3);border-radius:20px;font-size:.68rem;color:#818CF8;font-style:italic;}
.hbar{height:1px;margin:.4rem 0 1.4rem;background:linear-gradient(90deg,transparent,#1E2456 40%,#1E2456 60%,transparent);}
.step-card{background:rgba(15,23,42,.8);border:1px solid rgba(99,102,241,.12);border-radius:12px;padding:1.1rem 1.2rem;margin-bottom:1rem;}
.step-label{font-size:.6rem;font-weight:700;letter-spacing:.1em;text-transform:uppercase;color:#3B82F6;margin-bottom:.35rem;}
.step-title{font-size:.9rem;font-weight:800;color:#E4EDF8;margin-bottom:.2rem;}
.step-hint{font-size:.74rem;color:#374167;margin-bottom:.65rem;}
.upload-zone{border:2px dashed rgba(99,102,241,.3);border-radius:10px;padding:1.5rem;text-align:center;cursor:pointer;position:relative;transition:.2s;}
.upload-zone:hover{border-color:#3B82F6;background:rgba(59,130,246,.04);}
.upload-zone input{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%;}
.chip-ok{display:inline-block;padding:.22rem .75rem;background:rgba(34,197,94,.1);color:#4ADE80;border:1px solid rgba(34,197,94,.25);border-radius:7px;font-size:.72rem;font-weight:700;}
.chip-wait{display:inline-block;padding:.22rem .75rem;background:rgba(251,191,36,.08);color:#FBBF24;border:1px solid rgba(251,191,36,.22);border-radius:7px;font-size:.72rem;font-weight:700;}
.date-select{width:100%;padding:.6rem .85rem;background:rgba(8,15,40,.8);border:1.5px solid rgba(30,58,110,.8);border-radius:9px;color:#C8D8F8;font-size:.87rem;font-family:'Inter',sans-serif;}
.date-select:focus{outline:none;border-color:#3B82F6;}
.stat-row{display:flex;gap:.65rem;margin:.65rem 0 .8rem;}
.stat-box{flex:1;background:rgba(8,14,30,.8);border:1px solid rgba(21,52,112,.6);border-radius:10px;padding:.75rem .5rem;text-align:center;}
.stat-num{font-size:1.6rem;font-weight:900;color:#3B82F6;line-height:1.1;}
.stat-lbl{font-size:.6rem;color:#374167;text-transform:uppercase;letter-spacing:.09em;margin-top:.15rem;}
.btn-generate{width:100%;padding:.82rem;border-radius:11px;border:none;cursor:pointer;font-size:.95rem;font-weight:800;font-family:'Inter',sans-serif;transition:all .18s;background:linear-gradient(135deg,#1740C0,#2563EB 55%,#3B82F6);color:#fff;box-shadow:0 4px 22px rgba(29,78,216,.5);}
.btn-generate:hover{box-shadow:0 6px 30px rgba(37,99,235,.65);transform:translateY(-2px);}
.btn-generate:disabled{opacity:.4;cursor:not-allowed;transform:none;}
.btn-download{width:100%;padding:.82rem;border-radius:11px;border:none;cursor:pointer;font-size:.95rem;font-weight:800;font-family:'Inter',sans-serif;transition:all .18s;background:linear-gradient(135deg,#065F46,#059669 55%,#10B981);color:#fff;box-shadow:0 4px 22px rgba(5,150,105,.4);margin-top:.65rem;}
.btn-download:hover{box-shadow:0 6px 30px rgba(5,150,105,.6);transform:translateY(-2px);}
.status-msg{padding:.6rem .85rem;border-radius:9px;font-size:.82rem;margin:.65rem 0;display:none;}
.status-msg.show{display:block;}
.status-msg.ok{background:rgba(16,185,129,.08);border:1px solid rgba(16,185,129,.22);color:#10B981;}
.status-msg.error{background:rgba(239,68,68,.08);border:1px solid rgba(239,68,68,.22);color:#F87171;}
.spinner{animation:spin 1s linear infinite;display:inline-block;}
@keyframes spin{to{transform:rotate(360deg);}}




/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="top-bar">
  <span class="top-mini-brand">APSG</span>
  <div class="top-sep"></div>
  <span class="top-page-label">📋 Excel Rejection / Rejection Filter</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>

<div class="main">
  <div class="hero">
    <span class="hero-icon">📊</span>
    <div class="hero-title">Rejection Report Generator</div>
    <div class="hero-sub">Management of Staging Ground and Infilling Works (Phase 3)</div>
    <span class="hero-pill">✦ Convert Excel to PowerPoint Rejection Report</span>
  </div>
  <div class="hbar"></div>
  <div class="step-card">
    <div class="step-label">Step 1 of 3</div>
    <div class="step-title">📂 Upload Excel File</div>
    <div class="step-hint">Drag & drop or click — .xlsx, .xls, .csv accepted</div>
    <div class="upload-zone" id="uploadZone" onclick="document.getElementById('excelFile').click()"
         ondragover="event.preventDefault();this.classList.add('drag-over')"
         ondragleave="this.classList.remove('drag-over')"
         ondrop="event.preventDefault();this.classList.remove('drag-over');const f=event.dataTransfer.files[0];if(f){document.getElementById('excelFile').files=event.dataTransfer.files;onFileChange();}">
      <input type="file" id="excelFile" accept=".xlsx,.xls,.csv" onchange="onFileChange()" style="display:none">
      <div style="font-size:2rem;margin-bottom:.4rem">📂</div>
      <strong>Drop file here or click to browse</strong>
      <p style="font-size:.8rem;color:#374167;margin-top:.3rem">Required columns: Ticket, Vehicle, Material, Source Site, Date, Time, E-Token, Accepted, Reject Reason</p>
    </div>
    <div id="fileStatus" style="margin-top:.6rem"></div>
  </div>
  <div class="step-card" id="dateCard" style="display:none">
    <div class="step-label">Step 2 of 3</div>
    <div class="step-title">📅 Select Report Date</div>
    <select class="date-select" id="dateSelect" onchange="onDateChange()">
      <option value="">All dates</option>
    </select>
    <div id="dateStatus" style="margin-top:.6rem"></div>
  </div>
  <div class="step-card" id="generateCard" style="display:none">
    <div class="step-label">Step 3 of 3</div>
    <div class="step-title">🚀 Generate Report</div>
    <div class="stat-row" id="statRow"></div>
    <div class="badge-row" id="badgeRow" style="display:none;flex-wrap:wrap;gap:.35rem;margin:.4rem 0 .8rem"></div>
    <div class="status-msg" id="genStatus"></div>
    <button class="btn-generate" id="genBtn" onclick="generateReport()">⚡ Generate PowerPoint</button>
    <button class="btn-download" id="dlBtn" style="display:none" onclick="downloadResult()">⬇️ Download PowerPoint</button>
  </div>
</div>
<script>
let selectedFile=null, selectedDate='', resultBlob=null, resultName='';
const esc=s=>String(s).replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');
function onFileChange(){
  const f=document.getElementById('excelFile').files[0];
  if(!f)return;
  selectedFile=f;
  document.getElementById('fileStatus').innerHTML=`<span class="chip-ok">✓ ${f.name}</span>`;
  loadDates();
}
async function loadDates(){
  const fd=new FormData(); fd.append('file',selectedFile);
  // Reset badge row when loading new file
  const badgeEl=document.getElementById('badgeRow');
  if(badgeEl){badgeEl.innerHTML='';badgeEl.style.display='none';}
  try{
    const res=await fetch('/api/excel/dates',{method:'POST',body:fd});
    const data=await res.json();
    if(data.error){document.getElementById('fileStatus').innerHTML=`<span style="color:#F87171">❌ ${data.error}</span>`;return;}
    const sel=document.getElementById('dateSelect');
    sel.innerHTML='<option value="">All dates</option>'+data.dates.map(d=>`<option value="${d.value}">${d.label}</option>`).join('');
    if(data.dates.length>0){sel.value=data.dates[0].value;selectedDate=data.dates[0].value;}
    document.getElementById('dateCard').style.display='block';
    document.getElementById('generateCard').style.display='block';
    await updateStats();
  }catch(e){document.getElementById('fileStatus').innerHTML=`<span style="color:#F87171">❌ ${e.message}</span>`;}
}
function onDateChange(){selectedDate=document.getElementById('dateSelect').value;updateStats();}
async function updateStats(){
  if(!selectedFile){
    document.getElementById('statRow').innerHTML=
      `<div class="stat-box"><div class="stat-num">—</div><div class="stat-lbl">Rejections</div></div>
       <div class="stat-box"><div class="stat-num">—</div><div class="stat-lbl">Sites</div></div>
       <div class="stat-box"><div class="stat-num">—</div><div class="stat-lbl">Est. Slides</div></div>`;
    return;
  }
  try{
    const fd=new FormData();fd.append('file',selectedFile);fd.append('date',selectedDate);
    const res=await fetch('/api/excel/preview',{method:'POST',body:fd});
    const d=await res.json();
    if(d.error){return;}
    document.getElementById('statRow').innerHTML=
      `<div class="stat-box"><div class="stat-num">${d.rejections}</div><div class="stat-lbl">Rejections</div></div>
       <div class="stat-box"><div class="stat-num">${d.sites}</div><div class="stat-lbl">Sites</div></div>
       <div class="stat-box"><div class="stat-num">${d.est_slides}</div><div class="stat-lbl">Est. Slides</div></div>`;
    if(d.date_label){
      document.getElementById('dateStatus').innerHTML=`<span class="chip-ok">✓ ${esc(d.date_label)}</span>`;
    }
    if(d.badges&&d.badges.length>0){
      const bdg=d.badges.map(b=>`<span class="badge">${esc(b.site)} <span style="opacity:.45">×${b.count}</span></span>`).join('');
      const badgeEl=document.getElementById('badgeRow');
      if(badgeEl){badgeEl.innerHTML=bdg;badgeEl.style.display='flex';}
    }
  }catch(e){}
}
async function generateReport(){
  if(!selectedFile)return;
  const btn=document.getElementById('genBtn');
  btn.disabled=true;btn.innerHTML='<span class="spinner">⟳</span> Generating…';
  const st=document.getElementById('genStatus');st.className='status-msg show ok';st.textContent='Building report from template…';
  try{
    const fd=new FormData();fd.append('file',selectedFile);fd.append('date',selectedDate);
    const res=await fetch('/api/excel/generate',{method:'POST',body:fd});
    if(!res.ok){const e=await res.json().catch(()=>({error:`HTTP ${res.status}`}));st.className='status-msg show error';st.textContent='Error: '+(e.error||'failed');return;}
    resultBlob=await res.blob();
    const cd=res.headers.get('Content-Disposition')||'';
    resultName=cd.match(/filename="([^"]+)"/)?.[1]||`APSG-Loads_Rejected-${selectedDate.replace(/-/g,'').replace(/(\d{4})(\d{2})(\d{2})/,'$3$2$1')}.pptx`;
    st.textContent=`✅ Report ready — ${(resultBlob.size/1024).toFixed(0)} KB`;
    document.getElementById('dlBtn').style.display='block';
  }catch(e){st.className='status-msg show error';st.textContent='Error: '+e.message;}
  finally{btn.disabled=false;btn.innerHTML='⚡ Generate PowerPoint';}
}
function downloadResult(){if(!resultBlob)return;const a=document.createElement('a');a.href=URL.createObjectURL(resultBlob);a.download=resultName;a.click();}
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""

PHOTO_MERGE_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>APSG PowerPoint Merger</title>
<link href="https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=DM+Sans:ital,opsz,wght@0,9..40,300;0,9..40,400;0,9..40,500;0,9..40,600&display=swap" rel="stylesheet"/>
<style>
*,*::before,*::after{box-sizing:border-box;margin:0;padding:0;}
:root{
  --bg:#0A0E1A; --s1:#111827; --s2:#1A2234; --s3:#222E42;
  --b1:rgba(255,255,255,.07); --b2:rgba(255,255,255,.13);
  --teal:#00D4AA; --teal-d:rgba(0,212,170,.12); --teal-g:rgba(0,212,170,.28);
  --amber:#FFB340; --amber-d:rgba(255,179,64,.1);
  --indigo:#818CF8; --indigo-d:rgba(129,140,248,.12);
  --emerald:#34D399; --emerald-d:rgba(52,211,153,.12);
  --red:#FF5B5B; --text:#E8EDF5; --text2:#8A94A6; --text3:#3D4A5C;
  --r:14px; --rs:8px;
}
html{scroll-behavior:smooth;}
body{font-family:'DM Sans',sans-serif;background:transparent;color:var(--text);
  min-height:100vh;font-size:14px;line-height:1.6;overflow-x:hidden;}
body::before{content:'';position:fixed;inset:0;z-index:0;pointer-events:none;
  background:
    radial-gradient(ellipse 70% 50% at 5% 5%,rgba(0,212,170,.05) 0%,transparent 55%),
    radial-gradient(ellipse 50% 40% at 95% 85%,rgba(129,140,248,.05) 0%,transparent 50%),
    radial-gradient(ellipse 40% 60% at 50% 50%,rgba(255,179,64,.025) 0%,transparent 60%);}
.page{max-width:960px;margin:0 auto;padding:0 20px 60px;position:relative;z-index:1;}

/* ── HEADER ── */
header{display:flex;align-items:center;justify-content:space-between;
  padding:26px 0 28px;border-bottom:1px solid var(--b1);margin-bottom:32px;}
.logo{display:flex;align-items:center;gap:14px;}
.logo-mark{width:44px;height:44px;border-radius:12px;flex-shrink:0;
  background:linear-gradient(135deg,var(--teal),#00A882);font-size:20px;
  display:flex;align-items:center;justify-content:center;
  box-shadow:0 0 28px var(--teal-g);}
.logo-text h1{font-family:'Syne',sans-serif;font-size:19px;font-weight:800;
  letter-spacing:-.3px;}
.logo-text p{font-size:11.5px;color:var(--text2);margin-top:1px;}
.by-tag{font-size:11px;color:var(--text3);display:flex;align-items:center;gap:6px;}
.by-tag b{color:var(--teal);background:var(--teal-d);padding:3px 9px;border-radius:20px;
  border:1px solid rgba(0,212,170,.2);font-weight:600;}

/* ── CARDS ── */
.card{background:var(--s1);border:1px solid var(--b1);border-radius:var(--r);
  margin-bottom:18px;overflow:hidden;transition:border-color .2s;}
.card:hover{border-color:var(--b2);}
.card-head{display:flex;align-items:center;gap:12px;padding:16px 20px;
  border-bottom:1px solid var(--b1);}
.badge{width:34px;height:34px;border-radius:9px;flex-shrink:0;
  display:flex;align-items:center;justify-content:center;font-size:15px;}
.badge.teal{background:var(--teal-d);} .badge.amber{background:var(--amber-d);}
.badge.indigo{background:var(--indigo-d);}
.ch h3{font-family:'Syne',sans-serif;font-size:13.5px;font-weight:700;}
.ch p{font-size:11px;color:var(--text2);margin-top:2px;}
.card-body{padding:20px;}

/* ── UPLOAD ── */
.drop-grid{display:grid;grid-template-columns:1fr 1fr;gap:12px;}
@media(max-width:560px){.drop-grid{grid-template-columns:1fr;}}
.dz{border:1.5px dashed var(--b2);border-radius:var(--r);padding:26px 16px;
  text-align:center;cursor:pointer;transition:all .2s;background:var(--s2);
  display:flex;flex-direction:column;align-items:center;gap:5px;min-height:130px;justify-content:center;}
.dz:hover,.dz.drag{border-color:var(--teal);background:var(--teal-d);}
.dz.ok{border-color:var(--teal);border-style:solid;background:rgba(0,212,170,.07);}
.dz-icon{font-size:24px;} .dz-label{font-family:'Syne',sans-serif;font-size:13px;font-weight:700;}
.dz-hint{font-size:11px;color:var(--text3);}
.dz-name{font-size:11px;color:var(--teal);font-weight:600;word-break:break-all;max-width:180px;}
input[type=file]{display:none;}

/* ── PREVIEW CARD ── */
.prev-card{background:var(--s1);border:1.5px solid rgba(0,212,170,.22);
  border-radius:var(--r);margin-bottom:18px;overflow:hidden;
  box-shadow:0 0 48px rgba(0,212,170,.05);}
.prev-head{display:flex;align-items:center;gap:12px;padding:16px 20px;
  background:linear-gradient(90deg,rgba(0,212,170,.07),transparent);
  border-bottom:1px solid rgba(0,212,170,.13);}

/* Slide canvas */
.slide-wrap{background:var(--s2);border-radius:10px;overflow:hidden;
  margin-bottom:20px;border:1px solid var(--b1);}
.slide-bar{background:var(--s1);padding:7px 14px;font-size:10px;font-weight:600;
  color:var(--text2);letter-spacing:.8px;text-transform:uppercase;
  border-bottom:1px solid var(--b1);display:flex;align-items:center;gap:7px;}
.live-dot{width:6px;height:6px;border-radius:50%;background:var(--teal);
  box-shadow:0 0 6px var(--teal);animation:pulse 2s ease-in-out infinite;}
@keyframes pulse{0%,100%{opacity:1;}50%{opacity:.3;}}
.slide-canvas{width:100%;padding:10px;background:#F8FAFC;}
.slide-inner{position:relative;width:100%;}

/* ── CONTROL GROUPS ── */
.ctrl-grid{display:grid;grid-template-columns:1fr 1fr;gap:16px;}
@media(max-width:640px){.ctrl-grid{grid-template-columns:1fr;}}

.ctrl-group{background:var(--s2);border-radius:var(--r);padding:18px;
  border:1px solid var(--b1);transition:border-color .2s;}
.ctrl-group:hover{border-color:var(--b2);}
.ctrl-group.top-g{border-top:3px solid var(--indigo);}
.ctrl-group.frt-g{border-top:3px solid var(--emerald);}
.ctrl-group.gap-g{border-top:3px solid var(--amber);grid-column:1/-1;}

.g-title{font-family:'Syne',sans-serif;font-size:11.5px;font-weight:700;
  color:var(--text2);letter-spacing:.5px;text-transform:uppercase;
  margin-bottom:14px;display:flex;align-items:center;gap:7px;}
.g-dot{width:8px;height:8px;border-radius:50%;flex-shrink:0;}
.top-g .g-dot{background:var(--indigo);} .frt-g .g-dot{background:var(--emerald);}
.gap-g .g-dot{background:var(--amber);}

/* independence badge */
.ind-badge{margin-left:auto;font-size:9.5px;padding:2px 7px;border-radius:20px;
  font-weight:600;letter-spacing:.3px;}
.top-g .ind-badge{background:rgba(129,140,248,.15);color:var(--indigo);border:1px solid rgba(129,140,248,.2);}
.frt-g .ind-badge{background:rgba(52,211,153,.15);color:var(--emerald);border:1px solid rgba(52,211,153,.2);}

.field-2col{display:grid;grid-template-columns:1fr 1fr;gap:10px;margin-bottom:10px;}
.field{display:flex;flex-direction:column;gap:4px;}
.field label{font-size:11px;font-weight:500;color:var(--text2);}
.sr{display:flex;align-items:center;gap:7px;}
.sr input[type=range]{flex:1;height:4px;appearance:none;background:var(--b2);
  border-radius:99px;outline:none;cursor:pointer;}
.top-g .sr input[type=range]{accent-color:var(--indigo);}
.frt-g .sr input[type=range]{accent-color:var(--emerald);}
.gap-g .sr input[type=range]{accent-color:var(--amber);}
input[type=range]::-webkit-slider-thumb{appearance:none;width:15px;height:15px;
  border-radius:50%;cursor:pointer;box-shadow:0 0 6px rgba(0,0,0,.4);}
.top-g input[type=range]::-webkit-slider-thumb{background:var(--indigo);}
.frt-g input[type=range]::-webkit-slider-thumb{background:var(--emerald);}
.gap-g input[type=range]::-webkit-slider-thumb{background:var(--amber);}
.vbadge{min-width:42px;padding:2px 6px;border-radius:5px;font-size:12px;font-weight:700;
  text-align:center;background:var(--s1);border:1px solid var(--b2);}
.top-g .vbadge{color:var(--indigo);} .frt-g .vbadge{color:var(--emerald);}
.gap-g .vbadge{color:var(--amber);}
.num-in{width:100%;padding:7px 10px;background:var(--s1);border:1px solid var(--b2);
  border-radius:var(--rs);font-size:13px;font-weight:500;color:var(--text);
  outline:none;transition:border-color .2s;}
.num-in:focus{border-color:var(--teal);}
.fhint{font-size:10px;color:var(--text3);}

/* gap row */
.gap-row{display:grid;grid-template-columns:1fr auto;gap:12px;align-items:end;}

/* save default row */
.save-row{display:flex;gap:8px;align-items:center;margin-top:14px;
  padding-top:12px;border-top:1px solid var(--b1);}
.save-btn{padding:7px 14px;background:var(--teal-d);border:1px solid rgba(0,212,170,.25);
  border-radius:var(--rs);color:var(--teal);font-size:12px;font-weight:600;cursor:pointer;
  transition:all .2s;}
.save-btn:hover{background:rgba(0,212,170,.2);}
.reset-btn{padding:7px 14px;background:transparent;border:1px solid var(--b2);
  border-radius:var(--rs);color:var(--text2);font-size:12px;font-weight:600;cursor:pointer;
  transition:all .2s;}
.reset-btn:hover{border-color:var(--text2);color:var(--text);}
.save-status{font-size:11px;color:var(--teal);opacity:0;transition:opacity .4s;}

/* ── ADVANCED ── */
.adv-grid{display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px;}
@media(max-width:600px){.adv-grid{grid-template-columns:1fr;}}
.text-in,.sel-in{width:100%;padding:8px 11px;background:var(--s2);
  border:1px solid var(--b1);border-radius:var(--rs);font-size:13px;
  color:var(--text);outline:none;transition:border-color .2s;}
.text-in:focus,.sel-in:focus{border-color:var(--teal);}

/* ── RUN BUTTON ── */
.run-btn{width:100%;padding:13px;background:linear-gradient(135deg,var(--teal),#00A882);
  color:#041A14;border:none;border-radius:var(--r);
  font-family:'Syne',sans-serif;font-size:13px;font-weight:700;cursor:pointer;
  display:flex;align-items:center;justify-content:center;gap:8px;
  transition:all .2s;box-shadow:0 4px 28px var(--teal-g);margin-bottom:18px;letter-spacing:.2px;}
.run-btn:hover:not(:disabled){transform:translateY(-1px);box-shadow:0 6px 36px var(--teal-g);}
.run-btn:disabled{opacity:.35;cursor:not-allowed;transform:none;}

/* ── PROGRESS ── */
.prog-card{display:none;background:var(--s1);border:1px solid var(--b1);
  border-radius:var(--r);padding:20px;margin-bottom:18px;}
.prog-head{display:flex;align-items:center;gap:10px;margin-bottom:12px;}
.spinner{width:20px;height:20px;border-radius:50%;flex-shrink:0;
  border:2px solid rgba(0,212,170,.2);border-top-color:var(--teal);
  animation:spin .85s linear infinite;}
@keyframes spin{to{transform:rotate(360deg);}}
.prog-title{font-family:'Syne',sans-serif;font-size:14px;font-weight:700;}
.prog-track{background:var(--s2);border-radius:99px;height:5px;overflow:hidden;margin-bottom:8px;}
.prog-bar{height:100%;width:0%;border-radius:99px;transition:width .4s;
  background:linear-gradient(90deg,var(--teal),#00F5D4);}
.prog-st{font-size:11.5px;color:var(--text2);margin-bottom:10px;}
.log-term{background:#050A12;border-radius:var(--rs);padding:12px;
  font-family:'Courier New',monospace;font-size:11px;line-height:1.8;
  max-height:180px;overflow-y:auto;border:1px solid var(--b1);}
.log-ok{color:var(--teal);} .log-err{color:var(--red);} .log-li{color:#3D5A80;}

/* ── RESULT ── */
.res-card{display:none;background:var(--s1);border:1px solid rgba(0,212,170,.18);
  border-radius:var(--r);padding:20px;margin-bottom:18px;
  box-shadow:0 0 36px rgba(0,212,170,.05);}
.res-title{font-family:'Syne',sans-serif;font-size:16px;font-weight:800;
  display:flex;align-items:center;gap:10px;margin-bottom:6px;}
.res-sub{font-size:12px;color:var(--text2);margin-bottom:14px;}
.stats-row{display:flex;flex-wrap:wrap;gap:8px;margin-bottom:16px;}
.stat{background:var(--s2);border:1px solid var(--b1);border-radius:10px;
  padding:9px 14px;text-align:center;min-width:72px;}
.stat strong{display:block;font-size:21px;font-weight:800;line-height:1.1;}
.stat span{font-size:10px;color:var(--text2);text-transform:uppercase;letter-spacing:.5px;}
.dl-btn{display:inline-flex;align-items:center;gap:8px;
  background:linear-gradient(135deg,var(--teal),#00A882);color:#041A14;
  padding:10px 22px;border-radius:var(--rs);border:none;
  font-family:'Syne',sans-serif;font-size:13px;font-weight:800;cursor:pointer;
  margin-right:10px;box-shadow:0 4px 20px var(--teal-g);transition:all .2s;}
.dl-btn:hover{transform:translateY(-1px);}
.new-btn{display:inline-flex;align-items:center;gap:6px;
  background:transparent;border:1px solid var(--b2);color:var(--text2);
  padding:10px 18px;border-radius:var(--rs);font-size:13px;font-weight:600;cursor:pointer;
  transition:all .2s;}
.new-btn:hover{border-color:var(--text2);color:var(--text);}



/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>

<div class="top-bar">
  <span class="top-mini-brand">APSG</span>
  <div class="top-sep"></div>
  <span class="top-page-label">🖼 PPT Alignment — Top/Front Photo Merge</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>

<div class="page">

<!-- HEADER -->
<header>
  <div class="logo">
    <div class="logo-mark">📊</div>
    <div class="logo-text">
      <h1>APSG PowerPoint Merger</h1>
      <p>Independent Photo Layout · v6.0</p>
    </div>
  </div>
  <div class="by-tag">Created by <b>Karthi</b></div>
</header>

<!-- UPLOAD -->
<div class="card">
  <div class="card-head">
    <div class="badge teal">📁</div>
    <div class="ch"><h3>Upload PPTX Files</h3><p>Drop both files — reference auto-detected</p></div>
  </div>
  <div class="card-body">
    <div class="drop-grid">
      <div>
        <div class="dz" id="dzA" onclick="document.getElementById('iA').click()"
             ondragover="ev(event,'A',true)" ondragleave="ev(event,'A',false)" ondrop="drp(event,'A')">
          <div class="dz-icon">🪨</div><div class="dz-label">Top Photo File</div>
          <div class="dz-hint">Click or drag .pptx</div><div class="dz-name" id="nA"></div>
        </div>
        <input type="file" accept=".pptx" id="iA" onchange="onFile('A',this)"/>
      </div>
      <div>
        <div class="dz" id="dzB" onclick="document.getElementById('iB').click()"
             ondragover="ev(event,'B',true)" ondragleave="ev(event,'B',false)" ondrop="drp(event,'B')">
          <div class="dz-icon">🚛</div><div class="dz-label">Front Photo File</div>
          <div class="dz-hint">Click or drag .pptx</div><div class="dz-name" id="nB"></div>
        </div>
        <input type="file" accept=".pptx" id="iB" onchange="onFile('B',this)"/>
      </div>
    </div>
  </div>
</div>

<!-- PREVIEW + CONTROLS -->
<div class="prev-card">
  <div class="prev-head">
    <div class="badge amber">🖼️</div>
    <div class="ch">
      <h3>Photo Layout — Live Preview &amp; Independent Controls</h3>
      <p>Top Photo and Front Photo are completely independent — adjusting one never moves the other</p>
    </div>
  </div>
  <div class="card-body">

    <!-- Live preview -->
    <div class="slide-wrap">
      <div class="slide-bar"><div class="live-dot"></div>Live Slide Preview</div>
      <div class="slide-canvas" id="slideCanvas">
        <div class="slide-inner" id="slideInner"></div>
      </div>
    </div>

    <!-- Controls -->
    <div class="ctrl-grid">

      <!-- TOP PHOTO — independent -->
      <div class="ctrl-group top-g">
        <div class="g-title">
          <div class="g-dot"></div>🪨 Top Photo (Soil/Material)
          <span class="ind-badge">INDEPENDENT</span>
        </div>
        <div class="field-2col">
          <div class="field">
            <label>Height (cm)</label>
            <div class="sr">
              <input type="range" id="sTopH" min="4" max="14" step="0.01" value="9.11"
                     oninput="set('TopH',this.value)"/>
              <span class="vbadge" id="vTopH">9.11</span>
            </div>
            <input type="number" class="num-in" id="nTopH" min="4" max="14" step="0.01" value="9.11"
                   oninput="set('TopH',this.value)"/>
          </div>
          <div class="field">
            <label>Width (cm)</label>
            <div class="sr">
              <input type="range" id="sTopW" min="4" max="22" step="0.01" value="15.28"
                     oninput="set('TopW',this.value)"/>
              <span class="vbadge" id="vTopW">15.28</span>
            </div>
            <input type="number" class="num-in" id="nTopW" min="4" max="22" step="0.01" value="15.28"
                   oninput="set('TopW',this.value)"/>
          </div>
        </div>
        <div class="field">
          <label>Left Position (cm) — independent absolute position</label>
          <div class="sr">
            <input type="range" id="sTopL" min="0" max="20" step="0.01" value="0.95"
                   oninput="set('TopL',this.value)"/>
            <span class="vbadge" id="vTopL">0.95</span>
          </div>
          <input type="number" class="num-in" id="nTopL" min="0" max="20" step="0.01" value="0.95"
                 oninput="set('TopL',this.value)"/>
          <span class="fhint">Move left edge — Front Photo stays where it is</span>
        </div>
        <div class="save-row">
          <button class="save-btn" onclick="saveDefault('top')">💾 Save as Default</button>
          <button class="reset-btn" onclick="resetTop()">↺ Reset Top</button>
          <span class="save-status" id="topSaved">✓ Saved!</span>
        </div>
      </div>

      <!-- FRONT PHOTO — independent -->
      <div class="ctrl-group frt-g">
        <div class="g-title">
          <div class="g-dot"></div>🚛 Front Photo (Truck Plate)
          <span class="ind-badge">INDEPENDENT</span>
        </div>
        <div class="field-2col">
          <div class="field">
            <label>Height (cm)</label>
            <div class="sr">
              <input type="range" id="sFrtH" min="4" max="14" step="0.01" value="9.11"
                     oninput="set('FrtH',this.value)"/>
              <span class="vbadge" id="vFrtH">9.11</span>
            </div>
            <input type="number" class="num-in" id="nFrtH" min="4" max="14" step="0.01" value="9.11"
                   oninput="set('FrtH',this.value)"/>
          </div>
          <div class="field">
            <label>Width (cm)</label>
            <div class="sr">
              <input type="range" id="sFrtW" min="4" max="22" step="0.01" value="15.51"
                     oninput="set('FrtW',this.value)"/>
              <span class="vbadge" id="vFrtW">15.51</span>
            </div>
            <input type="number" class="num-in" id="nFrtW" min="4" max="22" step="0.01" value="15.51"
                   oninput="set('FrtW',this.value)"/>
          </div>
        </div>
        <div class="field">
          <label>Left Position (cm) — independent absolute position</label>
          <div class="sr">
            <input type="range" id="sFrtL" min="0" max="32" step="0.01" value="18.73"
                   oninput="set('FrtL',this.value)"/>
            <span class="vbadge" id="vFrtL">18.73</span>
          </div>
          <input type="number" class="num-in" id="nFrtL" min="0" max="32" step="0.01" value="18.73"
                 oninput="set('FrtL',this.value)"/>
          <span class="fhint">Move left edge — Top Photo stays where it is</span>
        </div>
        <div class="save-row">
          <button class="save-btn" onclick="saveDefault('frt')">💾 Save as Default</button>
          <button class="reset-btn" onclick="resetFrt()">↺ Reset Front</button>
          <span class="save-status" id="frtSaved">✓ Saved!</span>
        </div>
      </div>

      <!-- GAP (visual-only helper) -->
      <div class="ctrl-group gap-g">
        <div class="g-title"><div class="g-dot"></div>⬌ Visual Gap Helper</div>
        <div class="gap-row">
          <div class="field">
            <label>Desired gap between photos (cm) — sets Front Photo left automatically</label>
            <div class="sr">
              <input type="range" id="sGap" min="0" max="6" step="0.1" value="2.5"
                     oninput="applyGap(this.value)"/>
              <span class="vbadge" id="vGap">2.5</span>
            </div>
            <input type="number" class="num-in" id="nGap" min="0" max="6" step="0.1" value="2.5"
                   oninput="applyGap(this.value)"/>
            <span class="fhint">Shortcut: sets Front Photo left = Top left + Top width + gap. After applying, both photos remain independent.</span>
          </div>
          <div style="display:flex;flex-direction:column;gap:8px;padding-bottom:24px;">
            <button class="save-btn" onclick="applyGapToFrt()">Apply Gap →</button>
            <button class="reset-btn" onclick="resetAll()">↺ Reset All</button>
          </div>
        </div>
      </div>

    </div>
  </div>
</div>

<!-- ADVANCED -->
<div class="card">
  <div class="card-head">
    <div class="badge indigo">⚙️</div>
    <div class="ch"><h3>Advanced Settings</h3><p>Only change if your template uses different shape names</p></div>
  </div>
  <div class="card-body">
    <div class="adv-grid">
      <div class="field"><label>Top Photo shape name</label>
        <input type="text" class="text-in" id="topPh" value="Rectangle 14"/></div>
      <div class="field"><label>Front Photo shape name</label>
        <input type="text" class="text-in" id="frontPh" value="Rectangle 15"/></div>
      <div class="field"><label>Log level</label>
        <select class="sel-in" id="verbSel">
          <option value="0">Standard</option>
          <option value="1">Verbose (debug)</option>
        </select></div>
    </div>
  </div>
</div>

<!-- RUN -->
<button class="run-btn" id="runBtn" onclick="runMerge()" disabled>
  ⚡ Generate Report
</button>

<!-- PROGRESS -->
<div class="prog-card" id="progCard">
  <div class="prog-head">
    <div class="spinner"></div>
    <div class="prog-title" id="progTitle">Processing…</div>
  </div>
  <div class="prog-track"><div class="prog-bar" id="progBar"></div></div>
  <div class="prog-st" id="progSt"></div>
  <div class="log-term" id="logBox"></div>
</div>

<!-- RESULT -->
<div class="res-card" id="resCard">
  <div class="res-title">✅ Report Generated Successfully</div>
  <div class="res-sub" id="resMsg"></div>
  <div class="stats-row" id="statsRow"></div>
  <button class="dl-btn" id="dlBtn" style="display:none">⬇ Download PPTX</button>
  <button class="new-btn" onclick="doReset()">↺ New Merge</button>
</div>

</div><!-- /page -->
<script>
// ════════════════════════════════════════════════════════════════════════════
//  STATE — each photo fully independent, no shared calculations
// ════════════════════════════════════════════════════════════════════════════
const FACTORY = { TopH:9.11, TopW:15.28, TopL:0.95, FrtH:9.11, FrtW:15.51, FrtL:18.73, Gap:2.5 };

// Load saved defaults from localStorage (or fall back to factory)
function loadDefaults() {
  const saved = JSON.parse(localStorage.getItem('apsg_defaults') || '{}');
  return Object.assign({}, FACTORY, saved);
}

let v = loadDefaults();

const files = { A: null, B: null };
const SLIDE_W = 33.867, SLIDE_H = 19.05;
const LABEL_TOP = 6.73, LABEL_H = 0.9, IMG_TOP = LABEL_TOP + LABEL_H + 0.1;

// ── File upload ──────────────────────────────────────────────────────────────
function onFile(tag, inp) {
  const f = inp.files[0];
  if (!f) return;
  if (!f.name.toLowerCase().endsWith('.pptx')) { alert('Please upload a .pptx file.'); return; }
  files[tag] = f;
  document.getElementById('n'+tag).textContent = f.name;
  document.getElementById('dz'+tag).classList.add('ok');
  checkReady();
}
function ev(e,tag,on){ e.preventDefault(); document.getElementById('dz'+tag).classList.toggle('drag',on); }
function drp(e,tag) {
  e.preventDefault(); document.getElementById('dz'+tag).classList.remove('drag');
  const f = e.dataTransfer.files[0];
  if (f) { files[tag]=f; document.getElementById('n'+tag).textContent=f.name; document.getElementById('dz'+tag).classList.add('ok'); checkReady(); }
}
function checkReady() { document.getElementById('runBtn').disabled = !(files.A && files.B); }

// ── CRITICAL: set() updates ONLY the named field — zero effect on others ────
function set(key, rawVal) {
  const val = parseFloat(rawVal);
  if (isNaN(val)) return;

  // Update ONLY this field's own state
  v[key] = val;

  const dec = key === 'Gap' ? 1 : 2;
  const disp = val.toFixed(dec);

  // Update ONLY this field's own UI elements
  const s = document.getElementById('s'+key);
  const n = document.getElementById('n'+key);
  const b = document.getElementById('v'+key);
  if (s && String(s.value) !== String(val)) s.value = val;
  if (n && n.value !== disp) n.value = disp;
  if (b) b.textContent = disp;

  // Redraw preview from current v[] state
  renderPreview();
}

// ── Gap helper: applies gap by moving FRONT LEFT only ─────────────────────
function applyGap(raw) {
  const g = parseFloat(raw);
  if (isNaN(g)) return;
  v.Gap = g;
  document.getElementById('sGap').value = g;
  document.getElementById('nGap').value = g.toFixed(1);
  document.getElementById('vGap').textContent = g.toFixed(1);
  // Only updates the gap display — does NOT auto-apply to front left
  renderPreview();
}
function applyGapToFrt() {
  // One-shot: set Front Left = Top Left + Top Width + Gap
  // After this, both remain independent — Gap slider no longer linked
  const newFrtL = v.TopL + v.TopW + v.Gap;
  set('FrtL', parseFloat(newFrtL.toFixed(2)));
}

// ── Save defaults ─────────────────────────────────────────────────────────
function saveDefault(who) {
  const saved = JSON.parse(localStorage.getItem('apsg_defaults') || '{}');
  if (who === 'top') {
    saved.TopH = v.TopH; saved.TopW = v.TopW; saved.TopL = v.TopL;
    flashSaved('topSaved');
  } else {
    saved.FrtH = v.FrtH; saved.FrtW = v.FrtW; saved.FrtL = v.FrtL;
    flashSaved('frtSaved');
  }
  localStorage.setItem('apsg_defaults', JSON.stringify(saved));
}
function flashSaved(id) {
  const el = document.getElementById(id);
  el.style.opacity = '1';
  setTimeout(() => { el.style.opacity = '0'; }, 2000);
}

// ── Reset helpers ─────────────────────────────────────────────────────────
function resetTop()  { ['TopH','TopW','TopL'].forEach(k => set(k, FACTORY[k])); }
function resetFrt()  { ['FrtH','FrtW','FrtL'].forEach(k => set(k, FACTORY[k])); }
function resetAll()  { Object.keys(FACTORY).forEach(k => set(k, FACTORY[k])); }

// ── Preview renderer — pure function of v[], no side effects ──────────────
function renderPreview() {
  const canvas = document.getElementById('slideCanvas');
  const inner  = document.getElementById('slideInner');
  const CW = canvas.offsetWidth - 20;
  const CH = CW * (SLIDE_H / SLIDE_W);
  inner.style.width  = CW + 'px';
  inner.style.height = CH + 'px';
  canvas.style.height = (CH + 20) + 'px';

  const sc = CW / SLIDE_W;
  const cm = n => n * sc;

  // Top photo uses ONLY v.TopL, v.TopW, v.TopH — no front values
  const tL = v.TopL, tW = v.TopW, tH = v.TopH;
  // Front photo uses ONLY v.FrtL, v.FrtW, v.FrtH — no top values
  const fL = v.FrtL, fW = v.FrtW, fH = v.FrtH;
  // Gap shown for visual reference only
  const gap = fL - (tL + tW);

  let html = `<div style="position:absolute;inset:0;background:#F8FAFC;border-radius:6px;"></div>`;

  // Title
  html += `<div style="position:absolute;left:${cm(.86)}px;top:${cm(.63)}px;
    width:${cm(30)}px;height:${cm(1.35)}px;background:#E5E7EB;border-radius:3px;
    display:flex;align-items:center;padding:0 ${cm(.28)}px;
    font-size:${cm(.31)}px;font-weight:700;color:#374151;white-space:nowrap;overflow:hidden;">
    1. Rejection Due To Stone Found at Material Platform</div>`;

  // Table
  html += `<div style="position:absolute;left:${cm(.95)}px;top:${cm(2.67)}px;
    width:${cm(31.9)}px;height:${cm(2.38)}px;background:#F3F4F6;
    border:1px solid #D1D5DB;border-radius:3px;
    display:flex;align-items:center;justify-content:center;
    font-size:${cm(.24)}px;color:#9CA3AF;">Ticket Data Table</div>`;

  // TOP PHOTO label — uses only top values
  html += `<div style="position:absolute;left:${cm(tL)}px;top:${cm(LABEL_TOP)}px;
    width:${cm(tW)}px;height:${cm(LABEL_H)}px;
    background:#EEF2FF;border-radius:3px;
    display:flex;align-items:center;justify-content:center;
    font-size:${cm(.27)}px;font-weight:700;color:#4338CA;letter-spacing:.8px;">TOP PHOTO</div>`;

  // FRONT PHOTO label — uses only front values (independent position)
  html += `<div style="position:absolute;left:${cm(fL)}px;top:${cm(LABEL_TOP)}px;
    width:${cm(fW)}px;height:${cm(LABEL_H)}px;
    background:#ECFDF5;border-radius:3px;
    display:flex;align-items:center;justify-content:center;
    font-size:${cm(.27)}px;font-weight:700;color:#065F46;letter-spacing:.8px;">FRONT PHOTO</div>`;

  // TOP PHOTO image — ONLY tL, tW, tH used
  html += `<div style="position:absolute;
    left:${cm(tL)}px;top:${cm(IMG_TOP)}px;
    width:${cm(tW)}px;height:${cm(tH)}px;
    background:linear-gradient(135deg,#EEF2FF,#C7D2FE);
    border:2px solid #818CF8;border-radius:4px;overflow:hidden;
    display:flex;flex-direction:column;align-items:center;justify-content:center;gap:2px;">
    <div style="font-size:${cm(.5)}px;">🪨</div>
    <div style="font-size:${cm(.21)}px;font-weight:700;color:#3730A3;">Top Photo</div>
    <div style="font-size:${cm(.175)}px;color:#818CF8;background:rgba(129,140,248,.12);
      padding:1px ${cm(.08)}px;border-radius:3px;margin-top:2px;">
      ${tW.toFixed(2)} × ${tH.toFixed(2)} cm &nbsp;|&nbsp; left: ${tL.toFixed(2)} cm</div>
  </div>`;

  // Gap visual (amber) — purely decorative, non-linked
  if (gap > 0.05) {
    const gL = tL + tW, gW = gap;
    const gH = Math.max(tH, fH);
    html += `<div style="position:absolute;
      left:${cm(gL)}px;top:${cm(IMG_TOP)}px;
      width:${cm(gW)}px;height:${cm(gH)}px;
      background:repeating-linear-gradient(45deg,
        rgba(255,179,64,.07) 0,rgba(255,179,64,.07) 3px,transparent 3px,transparent 7px);
      border-left:1.5px dashed #F59E0B;border-right:1.5px dashed #F59E0B;
      display:flex;align-items:center;justify-content:center;overflow:hidden;">
      <div style="font-size:${cm(.17)}px;font-weight:700;color:#D97706;
        writing-mode:vertical-lr;transform:rotate(180deg);white-space:nowrap;">
        ${gW.toFixed(2)} cm</div></div>`;
  }

  // FRONT PHOTO image — ONLY fL, fW, fH used — no dependency on top
  html += `<div style="position:absolute;
    left:${cm(fL)}px;top:${cm(IMG_TOP)}px;
    width:${cm(fW)}px;height:${cm(fH)}px;
    background:linear-gradient(135deg,#ECFDF5,#A7F3D0);
    border:2px solid #34D399;border-radius:4px;overflow:hidden;
    display:flex;flex-direction:column;align-items:center;justify-content:center;gap:2px;">
    <div style="font-size:${cm(.5)}px;">🚛</div>
    <div style="font-size:${cm(.21)}px;font-weight:700;color:#065F46;">Front Photo</div>
    <div style="font-size:${cm(.175)}px;color:#34D399;background:rgba(52,211,153,.12);
      padding:1px ${cm(.08)}px;border-radius:3px;margin-top:2px;">
      ${fW.toFixed(2)} × ${fH.toFixed(2)} cm &nbsp;|&nbsp; left: ${fL.toFixed(2)} cm</div>
  </div>`;

  // Footer
  html += `<div style="position:absolute;left:0;bottom:${cm(.25)}px;width:100%;
    height:${cm(1.15)}px;background:#F3F4F6;border-top:1px solid #E5E7EB;
    display:flex;align-items:center;padding:0 ${cm(.35)}px;
    font-size:${cm(.21)}px;color:#9CA3AF;font-weight:500;">
    TOA-SAMSUNG C&amp;T JOINT VENTURE</div>`;

  inner.innerHTML = html;
}

// ── Merge ─────────────────────────────────────────────────────────────────
let curJob=null, pollT=null;

async function runMerge() {
  document.getElementById('runBtn').disabled = true;
  document.getElementById('progCard').style.display = 'block';
  document.getElementById('resCard').style.display  = 'none';
  document.getElementById('progBar').style.width    = '0%';
  document.getElementById('logBox').innerHTML        = '';
  document.getElementById('progTitle').textContent  = 'Processing…';
  document.getElementById('progSt').textContent     = 'Uploading files…';

  logLine(`Top: left=${v.TopL.toFixed(2)} w=${v.TopW.toFixed(2)} h=${v.TopH.toFixed(2)} cm`,'ok');
  logLine(`Front: left=${v.FrtL.toFixed(2)} w=${v.FrtW.toFixed(2)} h=${v.FrtH.toFixed(2)} cm`,'ok');
  logLine(`(Both fully independent — no shared values)`,'li');

  const fd = new FormData();
  fd.append('file_a', files.A); fd.append('file_b', files.B);
  fd.append('top_placeholder',  document.getElementById('topPh').value.trim()||'Rectangle 14');
  fd.append('front_placeholder',document.getElementById('frontPh').value.trim()||'Rectangle 15');
  fd.append('verbose', document.getElementById('verbSel').value==='1'?'1':'0');
  // Each photo sends its own independent values
  fd.append('top_h_cm',     v.TopH);
  fd.append('top_w_cm',     v.TopW);
  fd.append('top_left_cm',  v.TopL);
  fd.append('front_h_cm',   v.FrtH);
  fd.append('front_w_cm',   v.FrtW);
  fd.append('front_left_cm',v.FrtL);
  fd.append('center_gap_cm',v.Gap);

  try {
    const r = await fetch('/api/merge',{method:'POST',body:fd});
    const j = await r.json();
    if (j.error){showErr(j.error);return;}
    curJob=j.job_id; pollT=setInterval(poll,700);
  } catch(e){showErr(String(e));}
}

async function poll() {
  try {
    const r=await fetch('/api/job/'+curJob); const j=await r.json();
    document.getElementById('progBar').style.width=(j.progress||0)+'%';
    document.getElementById('progSt').textContent=j.status||'';
    if(j.log_lines) j.log_lines.slice(-5).forEach(([lv,msg])=>logLine(msg,lv==='ERROR'?'err':'li'));
    if(j.status==='complete'){clearInterval(pollT);showResult(j);}
    else if(j.status==='error'){clearInterval(pollT);showErr(j.error||'Error');}
  }catch(e){clearInterval(pollT);showErr(String(e));}
}

function showResult(j) {
  document.getElementById('progCard').style.display='none';
  document.getElementById('resCard').style.display='block';
  const st=j.stats||{};
  document.getElementById('resMsg').textContent=`Output: ${j.output_name||''}  ·  Reference: ${j.reference_file||'?'}`;
  const sr=document.getElementById('statsRow'); sr.innerHTML='';
  [['Total',st.total||0],['Merged',st.merged||0],['Both ✓',st.both||0],
   ['Top only',st.top_only||0],['Front only',st.front_only||0],['None',st.none_found||0]]
  .forEach(([l,n])=>{const d=document.createElement('div');d.className='stat';
    d.innerHTML=`<strong>${n}</strong><span>${l}</span>`;sr.appendChild(d);});
  const db=document.getElementById('dlBtn');
  db.style.display=j.has_result?'inline-flex':'none';
  db.onclick=()=>{window.location='/api/download/'+curJob;};
}
function showErr(msg){
  document.getElementById('progTitle').textContent='❌ Error';
  document.getElementById('progSt').textContent=msg;
  document.getElementById('runBtn').disabled=false;
}
function doReset(){
  files.A=null;files.B=null;curJob=null;
  ['A','B'].forEach(t=>{document.getElementById('n'+t).textContent='';
    document.getElementById('dz'+t).classList.remove('ok','drag');});
  document.getElementById('progCard').style.display='none';
  document.getElementById('resCard').style.display='none';
  document.getElementById('logBox').innerHTML='';
  document.getElementById('runBtn').disabled=true;
}
function logLine(msg,cls){
  const b=document.getElementById('logBox');
  const d=document.createElement('div');d.className='log-'+cls;
  d.textContent=msg;b.appendChild(d);b.scrollTop=b.scrollHeight;
}

// ── Init: load saved defaults into UI ────────────────────────────────────
window.addEventListener('load',()=>{
  Object.keys(v).forEach(k=>set(k,v[k]));
  renderPreview();
});
window.addEventListener('resize', renderPreview);
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>
"""

DAILY_REPORT_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Daily Report — APSG Report</title>
<style>

/* ═══ MODERN UI BASE (Blue/Purple/Cyan Theme) ═══ */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
:root {
  --bg:transparent; --card-bg: rgba(13,18,35,0.88);
  --indigo: #6366F1; --indigo-l: #818CF8; --cyan: #22D3EE;
  --purple: #A855F7; --green: #10B981; --amber: #F59E0B; --red: #F87171;
  --text: #E8EEF8; --muted: #64748B; --border: rgba(99,102,241,0.15);
}
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent; color: var(--text);
  position: relative;
}
body::before {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none;
  background:
    radial-gradient(ellipse 70% 55% at 5% 10%, rgba(99,102,241,.09) 0%, transparent 65%),
    radial-gradient(ellipse 55% 45% at 95% 90%, rgba(168,85,247,.07) 0%, transparent 65%);
}
body::after {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none; opacity: .016;
  background-image: linear-gradient(var(--indigo) 1px, transparent 1px),
    linear-gradient(90deg, var(--indigo) 1px, transparent 1px);
  background-size: 60px 60px;
}

/* ── Modern Top Bar ── */
.top-bar {
  position: sticky; top: 0; z-index: 200; height: 56px;
  display: flex; align-items: center; padding: 0 1.5rem; gap: .75rem;
  background: rgba(6,9,22,.96); backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px);
  border-bottom: 1px solid var(--border);
}
.top-mini-brand {
  font-family: 'Poppins', sans-serif;
  font-size: .82rem; font-weight: 700; color: var(--text);
  letter-spacing: -.01em; white-space: nowrap; flex-shrink: 0;
  display: flex; align-items: center; gap: .45rem;
}
.top-mini-brand::before {
  content: ''; width: 8px; height: 8px; border-radius: 50%;
  background: linear-gradient(135deg, var(--indigo), var(--cyan)); flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: var(--border); flex-shrink: 0; }
.top-page-label { font-size: .75rem; font-weight: 600; color: var(--muted); white-space: nowrap; }
.top-brand-tag { font-size: .7rem; color: var(--muted); }
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.08); border: 1px solid rgba(99,102,241,.18);
  border-radius: 8px; padding: .3rem .9rem; font-size: .7rem; font-weight: 600;
  color: var(--indigo-l); text-decoration: none; transition: all .2s; white-space: nowrap;
}
.back-btn:hover { background: rgba(99,102,241,.18); transform: translateX(-2px); }
.main { position: relative; z-index: 1; max-width: 1400px; margin: 0 auto; padding: .7rem 1rem 2rem; }

/* ── Cards & Containers ── */
.container, .page-content { position: relative; z-index: 1; }
.card, .section-card, .panel {
  background: var(--card-bg); border: 1px solid var(--border);
  border-radius: 16px; backdrop-filter: blur(16px); -webkit-backdrop-filter: blur(16px);
  box-shadow: 0 8px 32px rgba(0,0,0,.4), inset 0 1px 0 rgba(255,255,255,.03);
  transition: box-shadow .25s, border-color .25s;
}
.card:hover, .section-card:hover { border-color: rgba(99,102,241,.25); }

/* ── Upload Zone — Modern drag & drop ── */
.upload-zone, .dz {
  border: 2px dashed rgba(99,102,241,.3); border-radius: 14px;
  padding: 2rem; text-align: center; cursor: pointer;
  transition: all .22s; background: rgba(99,102,241,.03);
  position: relative;
}
.upload-zone:hover, .dz:hover, .upload-zone.drag-over, .dz.drag-over {
  border-color: var(--indigo); background: rgba(99,102,241,.08);
  box-shadow: 0 0 0 4px rgba(99,102,241,.12);
}
.upload-zone.ok, .dz.ok {
  border-color: var(--green); background: rgba(16,185,129,.06);
  border-style: solid;
}
.upload-zone.ok:hover, .dz.ok:hover {
  border-color: var(--green); background: rgba(16,185,129,.1);
  box-shadow: 0 0 0 4px rgba(16,185,129,.1);
}
.upload-icon { font-size: 2rem; margin-bottom: .5rem; display: block; }
.upload-label { font-size: .82rem; color: var(--muted); font-weight: 500; }
.upload-hint { font-size: .7rem; color: rgba(100,116,139,.6); margin-top: .25rem; }
.upload-filename { font-size: .78rem; color: var(--green); font-weight: 600; margin-top: .4rem; }

/* ── Modern Buttons ── */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  background:transparent;
  color: #fff; border: none; border-radius: 10px;
  padding: .7rem 1.4rem; font-size: .85rem; font-weight: 700;
  font-family: 'Inter', sans-serif; cursor: pointer; letter-spacing: .01em;
  box-shadow: 0 4px 18px rgba(99,102,241,.35);
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s;
}
.btn-primary::before, .btn-generate::before, .btn-teal::before, .modal-btn-primary::before {
  content: ''; position: absolute; inset: 0;
  background: linear-gradient(135deg, transparent, rgba(255,255,255,.12), transparent);
  transform: translateX(-100%); transition: transform .45s;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-2px) scale(1.01);
  box-shadow: 0 8px 28px rgba(99,102,241,.55);
}
.btn-primary:hover::before, .btn-generate:hover::before, .modal-btn-primary:hover::before {
  transform: translateX(100%);
}
.btn-primary:active, .btn-generate:active { transform: translateY(0) scale(.98); }
.btn-primary:disabled, .btn-generate:disabled { opacity: .45; cursor: not-allowed; transform: none; }

.btn-teal {
  background:transparent;
  box-shadow: 0 4px 18px rgba(16,185,129,.3);
}
.btn-teal:hover { box-shadow: 0 8px 28px rgba(16,185,129,.5); }

.btn-download, .btn-dl, .dl-btn {
  background: linear-gradient(135deg, #0D7A5F, #10B981);
  color: #fff; border: none; border-radius: 10px;
  padding: .65rem 1.3rem; font-size: .82rem; font-weight: 700;
  cursor: pointer; transition: all .22s;
  box-shadow: 0 4px 16px rgba(16,185,129,.3);
}
.btn-download:hover, .btn-dl:hover, .dl-btn:hover {
  transform: translateY(-2px); box-shadow: 0 8px 26px rgba(16,185,129,.5);
}

.btn-secondary, .btn-gray {
  background: rgba(30,41,86,.6); color: var(--muted);
  border: 1px solid rgba(99,102,241,.2); border-radius: 10px;
  padding: .65rem 1.2rem; font-size: .82rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-secondary:hover, .btn-gray:hover { background: rgba(99,102,241,.1); color: var(--indigo-l); }

.btn-del, .btn-danger {
  background: rgba(239,68,68,.08); color: var(--red);
  border: 1px solid rgba(239,68,68,.2); border-radius: 10px;
  padding: .6rem 1.2rem; font-size: .8rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-del:hover, .btn-danger:hover { background: rgba(239,68,68,.16); }
.btn-del:hover { animation: shake .3s ease; }

@keyframes shake {
  0%,100% { transform: translateX(0); }
  25% { transform: translateX(-3px); }
  75% { transform: translateX(3px); }
}

/* ── Form Inputs ── */
input[type=text], input[type=number], input[type=date], input[type=password],
select, textarea {
  background: rgba(6,10,24,.8); border: 1.5px solid rgba(30,41,86,.8);
  border-radius: 10px; color: var(--text); padding: .65rem .9rem;
  font-size: .85rem; font-family: 'Inter', sans-serif;
  transition: all .2s; width: 100%;
}
input:focus, select:focus, textarea:focus {
  outline: none; border-color: var(--indigo);
  box-shadow: 0 0 0 3px rgba(99,102,241,.16);
  background: rgba(10,14,32,.9);
}
input::placeholder, textarea::placeholder { color: rgba(100,116,139,.45); font-weight: 300; }

/* ── Alerts / Status ── */
.alert-success, .alert.success, .msg.success {
  background: rgba(16,185,129,.08); border: 1px solid rgba(16,185,129,.2);
  color: #34D399; border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-error, .alert.error, .msg.error {
  background: rgba(239,68,68,.08); border: 1px solid rgba(239,68,68,.2);
  color: var(--red); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-warn, .alert.warn { 
  background: rgba(245,158,11,.08); border: 1px solid rgba(245,158,11,.2);
  color: var(--amber); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}

/* ── Tables ── */
table { width: 100%; border-collapse: collapse; font-size: .8rem; }
thead th {
  background: rgba(99,102,241,.08); color: var(--muted);
  font-weight: 700; font-size: .68rem; letter-spacing: .06em;
  text-transform: uppercase; padding: .65rem .9rem; text-align: left;
  border-bottom: 1px solid var(--border);
}
tbody tr { border-bottom: 1px solid rgba(30,45,80,.3); transition: background .15s; }
tbody tr:hover { background: rgba(99,102,241,.04); }
tbody td { padding: .6rem .9rem; color: var(--text); }

/* ── Tabs ── */
.tab { position: relative; transition: all .2s; }
.tab.active { color: var(--indigo-l) !important; }
.tab.active::after {
  content: ''; position: absolute; bottom: -1px; left: 15%; right: 15%;
  height: 2px; border-radius: 2px; background: var(--indigo);
  animation: tabIn .2s ease;
}
@keyframes tabIn { from { left: 50%; right: 50%; } to { left: 15%; right: 15%; } }

/* ── Spinner ── */
@keyframes spin { to { transform: rotate(360deg); } }
.spinner { width: 20px; height: 20px; border: 2px solid rgba(99,102,241,.2);
  border-top-color: var(--indigo); border-radius: 50%; animation: spin 1s linear infinite; }

/* ══ ANIMATED BUTTONS — Daily Report ══════════════════════════════════════════ */
.btn-primary::after{content:'';position:absolute;top:-50%;left:-75%;width:50%;height:200%;
  background:rgba(255,255,255,.15);transform:skewX(-20deg);transition:left .5s;pointer-events:none;}
.btn-primary:hover::after{left:150%;}
.btn-dl{animation:gentlePulse 3s ease-in-out infinite;transition:transform .2s,box-shadow .2s !important;}
.btn-dl:hover{animation:none;transform:translateY(-3px) !important;box-shadow:0 8px 28px rgba(16,185,129,.5) !important;}
@keyframes gentlePulse{0%,100%{box-shadow:0 4px 16px rgba(16,185,129,.3);}50%{box-shadow:0 4px 24px rgba(16,185,129,.55);}}
.upload-zone:not(.ok){animation:uploadGlow 4s ease-in-out infinite;}
.upload-zone:hover{animation:none;}
@keyframes uploadGlow{0%,100%{border-color:rgba(99,102,241,.3);}50%{border-color:rgba(99,102,241,.65);box-shadow:0 0 18px rgba(99,102,241,.12);}}

.hero{text-align:center;padding:.8rem 1rem .9rem;}
.hero-icon{font-size:2rem;display:block;margin-bottom:.25rem;}
.hero-title{font-size:1.5rem;font-weight:900;color:#F1F5FF;letter-spacing:-.03em;}
.hero-sub{font-size:.82rem;color:#374167;margin:.4rem 0;}
.hero-pill{display:inline-block;padding:.22rem .9rem;background:rgba(99,102,241,.1);border:1px solid rgba(99,102,241,.3);border-radius:20px;font-size:.7rem;color:#818CF8;font-style:italic;}
.hbar{height:1px;margin:.3rem 0 .9rem;background:linear-gradient(90deg,transparent,#1E2456 40%,#1E2456 60%,transparent);}
.two-col{display:grid;grid-template-columns:1fr 1fr;gap:1.5rem;}
@media(max-width:900px){.two-col{grid-template-columns:1fr;}}
.panel{border-radius:14px;overflow:hidden;margin-bottom:1rem;}
.panel-head{padding:.75rem 1.1rem;font-size:1rem;font-weight:800;letter-spacing:.015em;}
.panel-online .panel-head{background:linear-gradient(90deg,#0D1B40,#0A1530);color:#818CF8;border-left:4px solid #6366F1;}
.panel-wb .panel-head{background:linear-gradient(90deg,#0B1F18,#081610);color:#34D399;border-left:4px solid #10B981;}
.panel-body{background:#060C1C;border:1px solid #111827;border-top:none;padding:1.2rem;border-radius:0 0 14px 14px;}
.sec-card{background:#080F24;border:1px solid #162040;border-radius:12px;padding:.85rem 1.1rem;margin:.4rem 0 .85rem;}
.sec-label{font-size:.62rem;font-weight:700;letter-spacing:.12em;text-transform:uppercase;color:#6366F1;margin-bottom:.3rem;}
.sec-label-wb{color:#10B981;}
.sec-title{font-size:.95rem;font-weight:800;color:#E8EEF8;margin-bottom:.15rem;}
.sec-hint{font-size:.75rem;color:#2D3F60;line-height:1.5;}
.upload-zone{border:2px dashed rgba(99,102,241,.35);border-radius:12px;padding:1.8rem;text-align:center;cursor:pointer;position:relative;transition:.2s;background:rgba(8,15,40,.6);}
.upload-zone:hover{border-color:#6366F1;background:rgba(99,102,241,.06);}
.upload-zone input{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%;}
.upload-zone-wb{border-color:rgba(16,185,129,.35);}
.upload-zone-wb:hover{border-color:#10B981;background:rgba(16,185,129,.05);}
.chip-ok{display:inline-block;padding:.25rem .8rem;background:rgba(16,185,129,.12);color:#10B981;border:1px solid rgba(16,185,129,.3);border-radius:7px;font-size:.75rem;font-weight:700;}
.chip-wait{display:inline-block;padding:.25rem .8rem;background:rgba(251,191,36,.08);color:#F59E0B;border:1px solid rgba(251,191,36,.22);border-radius:7px;font-size:.75rem;font-weight:700;}
.chip-info{display:inline-block;padding:.25rem .8rem;background:rgba(99,102,241,.1);color:#818CF8;border:1px solid rgba(99,102,241,.25);border-radius:7px;font-size:.75rem;font-weight:700;}
.date-input{width:100%;padding:.6rem .9rem;background:#070D1E;border:1.5px solid #1E2456;border-radius:9px;color:#C8D8F8;font-size:.88rem;font-family:'Inter',sans-serif;}
.date-input:focus{outline:none;border-color:#6366F1;}
.stats-box{background:#060C1C;border:1px solid #111827;border-radius:12px;padding:.85rem 1rem;margin:.5rem 0 .85rem;}
.stats-box-wb{border-color:#0D2A1A;}
.stats-title{font-size:.65rem;font-weight:700;letter-spacing:.12em;text-transform:uppercase;margin-bottom:.55rem;}
.stats-title-online{color:#6366F1;}
.stats-title-wb{color:#10B981;}
.stats-inner table{width:100%;border-collapse:collapse;font-size:.82rem;}
.stats-inner td{padding:.28rem 0;color:#6B7280;}
.stats-inner td:last-child{text-align:right;font-weight:700;}
.val-ok{background:#031A10;border:1px solid #065F38;border-radius:11px;padding:.7rem 1rem;margin:.5rem 0;color:#10B981;font-size:.82rem;}
.val-err{background:#1A0508;border:1px solid #7F1D1D;border-radius:11px;padding:.7rem 1rem;margin:.5rem 0;color:#F87171;font-size:.82rem;}
.val-warn{background:#150E02;border:1px solid #78450A;border-radius:11px;padding:.7rem 1rem;margin:.5rem 0;color:#F59E0B;font-size:.82rem;}
.cmp-table{width:100%;border-collapse:collapse;font-size:.82rem;margin:.5rem 0 .8rem;}
.cmp-table th{background:#080F24;color:#818CF8;font-weight:700;padding:.5rem .8rem;border:1px solid #1E2456;text-align:center;}
.cmp-table td{padding:.42rem .8rem;border:1px solid #111827;color:#C8D8F8;text-align:center;}
.cmp-match{color:#10B981;font-weight:900;}
.cmp-miss{color:#F87171;font-weight:900;}
.pivot-outer{background:#060C1C;border:1px solid #111827;border-radius:13px;padding:1rem;margin:.8rem 0;overflow-x:auto;}
.pivot-title{font-size:.88rem;font-weight:800;color:#10B981;letter-spacing:.06em;text-align:center;padding:.4rem 1rem;background:rgba(16,185,129,.07);border-radius:8px;border:1px solid rgba(16,185,129,.18);text-transform:uppercase;margin-bottom:.8rem;}
.pivot-tbl{width:100%;border-collapse:collapse;font-size:.8rem;min-width:500px;}
.pivot-tbl th{background:#080F24;color:#818CF8;font-weight:700;padding:.5rem .7rem;border:1px solid #1E2456;text-align:center;}
.pivot-tbl td{padding:.42rem .7rem;border:1px solid #111827;color:#C8D8F8;text-align:center;}
.pivot-mat td{background:#07101E;color:#818CF8;font-weight:800;text-align:left;font-size:.84rem;border-top:2px solid #1E2456;}
.pivot-grand td{background:#08122A;color:#10B981;font-weight:900;border-top:2px solid #1E2456;}
.btn{width:100%;padding:.82rem 1rem;border-radius:12px;border:none;cursor:pointer;font-size:.95rem;font-weight:700;font-family:'Inter',sans-serif;transition:all .18s;margin-top:.5rem;}
.btn-primary{background:linear-gradient(135deg,#4338CA,#6366F1 55%,#818CF8);color:#fff;box-shadow:0 4px 20px rgba(99,102,241,.4);}
.btn-primary:hover{transform:translateY(-2px);box-shadow:0 6px 28px rgba(99,102,241,.55);}
.btn-primary:disabled{opacity:.4;cursor:not-allowed;transform:none;}
.btn-dl{background:linear-gradient(135deg,#065F38,#059669 55%,#10B981);color:#fff;box-shadow:0 4px 20px rgba(5,150,105,.4);}
.btn-dl:hover{transform:translateY(-2px);box-shadow:0 6px 28px rgba(16,185,129,.5);}
.btn-dl:disabled{opacity:.4;cursor:not-allowed;transform:none;}
.action-card{background:#0A1020;border:1.5px solid #F59E0B;border-radius:12px;padding:1rem;margin:.6rem 0;}
.action-title{font-size:.82rem;font-weight:800;color:#F59E0B;margin-bottom:.5rem;}
.action-select,.input-num{width:100%;padding:.55rem .8rem;background:#070D1E;border:1.5px solid #1E2456;border-radius:8px;color:#C8D8F8;font-size:.85rem;font-family:'Inter',sans-serif;margin-bottom:.3rem;}
.rr-panel{margin-top:.8rem;padding:.75rem;background:rgba(99,102,241,.07);border:1.5px solid rgba(99,102,241,.3);border-radius:10px;}
.rr-panel-title{font-size:.76rem;font-weight:800;color:#818CF8;margin-bottom:.5rem;letter-spacing:.03em;}
.rr-input{width:100%;padding:.48rem .75rem;background:#070D1E;border:1.5px solid #1E2456;border-radius:8px;color:#C8D8F8;font-size:.83rem;font-family:'Inter',sans-serif;margin-bottom:.4rem;}
.rr-select{width:100%;padding:.48rem .75rem;background:#070D1E;border:1.5px solid #1E2456;border-radius:8px;color:#C8D8F8;font-size:.83rem;font-family:'Inter',sans-serif;margin-bottom:.4rem;}
.rr-btn{display:inline-flex;align-items:center;gap:.4rem;padding:.45rem 1.1rem;background:linear-gradient(135deg,#4F46E5,#7C3AED);color:#fff;border:none;border-radius:8px;font-size:.8rem;font-weight:700;cursor:pointer;transition:opacity .2s;}
.rr-btn:hover{opacity:.85;}
.rr-btn:disabled{opacity:.4;cursor:not-allowed;}
.rr-status{font-size:.74rem;margin-top:.3rem;min-height:1.2rem;}
.step-row{display:flex;gap:.4rem;align-items:center;margin:.2rem 0 1rem;flex-wrap:wrap;}
.step-pip{display:inline-flex;align-items:center;gap:.3rem;padding:.2rem .7rem;border-radius:20px;font-size:.7rem;font-weight:700;}
.step-done{background:rgba(16,185,129,.15);color:#10B981;border:1px solid rgba(16,185,129,.35);}
.step-wait{background:rgba(245,158,11,.08);color:#F59E0B;border:1px solid rgba(245,158,11,.22);}
.step-arr{color:#1E2456;font-size:.7rem;}
hr{border:none;border-top:1px solid #0C1325;margin:1rem 0;}
.spinner{animation:spin 1s linear infinite;display:inline-block;}
@keyframes spin{to{transform:rotate(360deg);}}



/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="top-bar">
  <span class="top-mini-brand">APSG</span>
  <div class="top-sep"></div>
  <span class="top-page-label">📊 Daily Report Generator</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>

<div class="main">
  <div class="hero">
    <span class="hero-icon">📊</span>
    <div class="hero-title">Daily Report Generator</div>
    <div class="hero-sub">Staging Ground Report System · Phase 3</div>
    <span class="hero-pill">✦ WB Summary Contract Report · WB Server Comparison</span>
  </div>
  <div class="hbar"></div>
  <div class="two-col">
    <div>
      <div class="panel panel-online">
        <div class="panel-head">🌐 ONLINE SUMMARY / ONLINE DATA</div>
        <div class="panel-body">
          <div class="sec-card">
            <div class="sec-label">Step 1 — Upload</div>
            <div class="sec-title">📂 Upload Online Export</div>
            <div class="sec-hint">Drag & drop or click — .xlsx, .xls, .csv accepted</div>
          </div>
          <div class="upload-zone" onclick="document.getElementById('onlineFile').click()">
            <input type="file" id="onlineFile" accept=".xlsx,.xls,.csv" onchange="onOnlineFileChange()">
            <div style="font-size:2rem;margin-bottom:.4rem">📂</div>
            <strong style="font-size:.9rem">Drop or click to upload Online export</strong>
            <p style="font-size:.76rem;color:#2D3F60;margin-top:.3rem">.xlsx · .xls · .csv</p>
          </div>
          <div id="onlineFileStatus" style="margin:.5rem 0"></div>
          <hr>
          <div id="onlineDateSection" style="display:none">
            <div class="sec-card">
              <div class="sec-label">Step 2 — Filter Date</div>
              <div class="sec-title">📅 Filter Date (Online)</div>
              <div class="sec-hint">Filtered on <strong style="color:#60A5FA">WB In Time</strong></div>
            </div>
            <input type="date" class="date-input" id="onlineDate" onchange="onOnlineDateChange()" style="color-scheme:dark;">
            <div id="onlineDateStatus" style="margin:.5rem 0"></div>
            <hr>
          </div>
          <div id="onlineValidationSection" style="display:none">
            <div id="onlineActionCards"></div>
            <div id="onlineValResult"></div>
            <div id="onlineStats" style="display:none"></div>
            <hr>
            <div class="step-row" id="stepTracker"></div>
            <button class="btn btn-primary" id="generateBtn" onclick="generateReport()" disabled>⚡ Generate Report</button>
            <div id="generateStatus" style="margin:.5rem 0;font-size:.82rem;color:#818CF8"></div>
            <button class="btn btn-dl" id="downloadBtn" onclick="downloadReport()" disabled style="display:none">⬇️ Download Report</button>
          </div>
        </div>
      </div>
    </div>
    <div>
      <div class="panel panel-wb">
        <div class="panel-head">⚖️ WB DATA <span style="font-size:.72rem;font-weight:600;opacity:.65">WB Server · Comparison View</span></div>
        <div class="panel-body">
          <div class="sec-card">
            <div class="sec-label sec-label-wb">Step 1 — Upload</div>
            <div class="sec-title">📂 Upload WB Transaction List</div>
            <div class="sec-hint">.xlsx / .xls / .csv — Weighbridge server export</div>
          </div>
          <div class="upload-zone upload-zone-wb" onclick="document.getElementById('wbFile').click()">
            <input type="file" id="wbFile" accept=".xlsx,.xls,.csv" onchange="onWbFileChange()">
            <div style="font-size:2rem;margin-bottom:.4rem">⚖️</div>
            <strong style="font-size:.9rem">Drop or click to upload WB file</strong>
            <p style="font-size:.76rem;color:#2D3F60;margin-top:.3rem">.xlsx · .xls · .csv</p>
          </div>
          <div id="wbFileStatus" style="margin:.5rem 0"></div>
          <hr>
          <div id="wbDateSection" style="display:none">
            <div class="sec-card">
              <div class="sec-label sec-label-wb">Step 2 — Filter Date</div>
              <div class="sec-title">📅 Filter Date (WB)</div>
              <div class="sec-hint">Aligns with Online for comparison</div>
            </div>
            <input type="date" class="date-input" id="wbDate" onchange="onWbDateChange()" style="color-scheme:dark;border-color:#1E3A2A;">
            <hr>
          </div>
          <div id="wbResultSection" style="display:none">
            <div id="wbStats"></div>
            <div id="comparisonTable"></div>
            <div id="pivotSection"></div>
            <button class="btn btn-dl" id="pivotDlBtn" onclick="downloadPivot()" style="display:none;margin-top:.5rem">⬇️ Download Pivot (Excel)</button>
          </div>
        </div>
      </div>
    </div>
  </div>
</div>
<script>
// ── State ─────────────────────────────────────────────────────────────────────
let onlineTmpPath='',onlineDate='',wbTmpPath='',wbDate='';
let corrections={},wbDecisions={},reportBlob=null,reportName='',pivotTmpPath='',pivotFname='';
let onlineStats=null,wbStats=null;
const _incompleteItems={};   // idx → item metadata (token, in_weight)
const _rrSerials={};         // idx → saved RR serial string
const _rrReasons={};         // idx → saved RR reason string
let _validationTimer=null;   // debounce handle
let _wbPivotTimer=null;

// ═══════════════════════════════════════════════════════════════════════════════
// CORE FIX: scheduleValidation NEVER clears cards while user is typing.
// It only fires AFTER a 900 ms idle gap, and _syncActionCards NEVER destroys
// a card that still has an active correction or a focused input.
// ═══════════════════════════════════════════════════════════════════════════════

function _userIsTyping(){
  const active=document.activeElement;
  if(!active)return false;
  const tag=active.tagName.toLowerCase();
  if(tag==='input'||tag==='textarea')return true;
  return false;
}

function scheduleValidation(){
  clearTimeout(_validationTimer);
  _validationTimer=setTimeout(function(){
    if(_userIsTyping()){scheduleValidation();return;} // re-defer while typing
    runOnlineValidation();
  },900);
}

// ── File upload ───────────────────────────────────────────────────────────────
async function onOnlineFileChange(){
  const f=document.getElementById('onlineFile').files[0];if(!f)return;
  setEl('onlineFileStatus','<span class="chip-info">⏳ Uploading…</span>');
  const fd=new FormData();fd.append('file',f);
  try{
    const res=await fetch('/api/daily/upload',{method:'POST',body:fd});
    const d=await res.json();
    if(d.error){setEl('onlineFileStatus',`<span class="chip-wait">❌ ${esc(d.error)}</span>`);return;}
    onlineTmpPath=d.tmp_path;
    setEl('onlineFileStatus',`<span class="chip-ok">✓ ${esc(f.name)} — ${d.rows.toLocaleString()} rows loaded</span>`);
    const di=document.getElementById('onlineDate');
    di.min=d.min_date;di.max=d.max_date;di.value=d.max_date;onlineDate=d.max_date;
    show('onlineDateSection');show('onlineValidationSection');
    corrections={};await runOnlineValidation();
  }catch(e){setEl('onlineFileStatus',`<span class="chip-wait">❌ ${e.message}</span>`);}
}

async function onOnlineDateChange(){
  onlineDate=document.getElementById('onlineDate').value;corrections={};
  clearEl('onlineActionCards');clearEl('onlineValResult');clearEl('onlineStats');
  hide('downloadBtn');document.getElementById('generateBtn').disabled=true;
  await runOnlineValidation();
}

// ── Validation ────────────────────────────────────────────────────────────────
async function runOnlineValidation(){
  if(!onlineTmpPath||!onlineDate)return;
  setEl('onlineValResult','<div class="chip-info">🔍 Running validation…</div>');
  try{
    const res=await fetch('/api/daily/validate',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({tmp_path:onlineTmpPath,filter_date:onlineDate,corrections})});
    const d=await res.json();
    if(d.error){setEl('onlineValResult',`<div class="val-err">❌ ${esc(d.error)}</div>`);return;}
    if(d.rows===0){setEl('onlineValResult','<div class="val-warn">⚠️ No records for selected date.</div>');return;}
    _syncActionCards(d.incomplete);
    if(d.all_resolved||d.incomplete.length===0){
      if(d.errors&&d.errors.length>0){
        setEl('onlineValResult',d.errors.map(e=>`<div class="val-err">❌ ${esc(e)}</div>`).join(''));
      }else{
        setEl('onlineValResult','<div class="val-ok">✅ All validations completed successfully</div>');
      }
      onlineStats=d.stats;renderOnlineStats(d.stats);show('onlineStats');
      renderComparisonTable();
    }else{clearEl('onlineValResult');}
    updateStepTracker();updateGenerateBtn();
  }catch(e){setEl('onlineValResult',`<div class="val-err">❌ ${e.message}</div>`);}
}

// ── Smart card sync ───────────────────────────────────────────────────────────
// KEY RULES:
//  1. Never destroy a card that has a focused input inside it.
//  2. Never destroy a card that has a pending correction (user has made a choice).
//  3. Only add cards for new incomplete items; skip existing ones entirely.
function _syncActionCards(incomplete){
  const container=document.getElementById('onlineActionCards');
  if(!container)return;

  // Store metadata for any new items
  if(incomplete&&incomplete.length){
    for(const item of incomplete)_incompleteItems[item.idx]=item;
  }

  // Build set of idx values the server says are still incomplete
  const serverIncomplete=new Set((incomplete||[]).map(i=>i.idx));

  // Remove cards ONLY if: server no longer lists them AND no correction AND not focused
  container.querySelectorAll('[id^="acard-"]').forEach(el=>{
    const idx=parseInt(el.id.replace('acard-',''));
    const hasFocus=el.contains(document.activeElement);
    const hasCorrection=!!corrections[idx];
    if(!serverIncomplete.has(idx)&&!hasFocus&&!hasCorrection){
      el.remove();
    }
  });

  // Add cards for items not yet in the DOM
  for(const item of (incomplete||[])){
    if(document.getElementById(`acard-${item.idx}`))continue;
    const corr=corrections[item.idx]||{};
    const sel=corr.Accepted||'';
    const card=document.createElement('div');
    card.className='action-card';card.id=`acard-${item.idx}`;
    card.innerHTML=`
      <div class="action-title">⚠️ Action Required — E-Token: <span style="color:#F59E0B;font-weight:900">${esc(item.token)}</span></div>
      <div style="font-size:12px;color:#9BB4D4;margin-bottom:.4rem;">Missing Date Out / Time Out — Select action to resolve:</div>
      <select class="action-select" id="dec-${item.idx}">
        <option value="" ${sel===''?'selected':''}>— select —</option>
        <option value="YES" ${sel==='YES'?'selected':''}>YES — Accept</option>
        <option value="NO" ${sel==='NO'?'selected':''}>NO — Reject</option>
      </select>
      <div id="dec-detail-${item.idx}"></div>
      <div id="rr-panel-${item.idx}"></div>`;
    card.querySelector('select').addEventListener('change',function(){
      onDecisionChange(item.idx,item.in_weight,this.value);
    });
    container.appendChild(card);
    if(corrections[item.idx]){
      renderDecisionDetail(item.idx,item.in_weight,corrections[item.idx].Accepted,true);
    }
  }
}

// ── Decision handlers ─────────────────────────────────────────────────────────
function onDecisionChange(idx,inW,decision){
  renderDecisionDetail(idx,inW,decision,false);
}

function renderDecisionDetail(idx,inW,decision,restoring){
  const el=document.getElementById(`dec-detail-${idx}`);if(!el)return;

  if(decision==='NO'){
    corrections[idx]={Accepted:'NO','Out Weight':inW,'Net Weight':0};
    el.innerHTML='<span class="chip-ok">✓ Marked Rejected</span>';
    _ensureRRPanel(idx,inW,'NO',0,0);
    if(!restoring)scheduleValidation();

  }else if(decision==='YES'){
    const prev=(corrections[idx]&&corrections[idx]['Out Weight'])||0;
    // Only build the input DOM once — never replace it
    if(!document.getElementById(`outw-${idx}`)){
      el.innerHTML=`
        <div style="font-size:12px;color:#9BB4D4;font-weight:700;letter-spacing:.06em;text-transform:uppercase;margin-bottom:.4rem;">📥 Enter Net Weight → Out Weight</div>
        <div style="display:flex;align-items:center;gap:.5rem;margin-bottom:.3rem;">
          <label style="font-size:13px;color:#C8D8F8;font-weight:600;white-space:nowrap;min-width:90px;">Out Weight (T):</label>
          <input type="number" class="input-num" id="outw-${idx}"
            placeholder="e.g. 18.250" min="0" max="${inW}" step="0.001"
            value="${prev||''}" autocomplete="off"
            style="font-size:14px;padding:.6rem .9rem;border:1.5px solid rgba(99,102,241,.5);background:rgba(6,10,28,.9);color:#E8F0FF;border-radius:9px;flex:1;">
        </div>
        <div style="font-size:12px;color:#6B7EC8;margin-bottom:.2rem;">In Weight: <strong style="color:#10B981">${inW} T</strong></div>
        <div id="outw-s-${idx}" style="margin-top:.35rem;font-size:13px;font-weight:600;min-height:1.4rem;"></div>`;
      const inp=document.getElementById(`outw-${idx}`);
      // Use 'input' event only — fires on each keystroke, updates status label only
      inp.addEventListener('input',function(){
        _applyOutWValue(idx,inW,parseFloat(this.value)||0,false);
      });
      // Only schedule validation on blur (when user leaves the field)
      inp.addEventListener('blur',function(){
        const v=parseFloat(this.value)||0;
        if(v>0)scheduleValidation();
      });
      // Stop any key event from bubbling out of this input
      inp.addEventListener('keydown',function(e){e.stopPropagation();});
    }
    if(prev>0)_applyOutWValue(idx,inW,prev,true);
    else _ensureRRPanel(idx,inW,'YES',0,0);

  }else{
    delete corrections[idx];
    el.innerHTML='';
    const rp=document.getElementById(`rr-panel-${idx}`);if(rp)rp.innerHTML='';
  }
}

function onOutW(idx,inW,outW){_applyOutWValue(idx,inW,outW,false);}

function _applyOutWValue(idx,inW,outW,restoring){
  const statusEl=document.getElementById(`outw-s-${idx}`);if(!statusEl)return;
  if(outW<=0){
    statusEl.innerHTML='<span style="color:#F59E0B">Enter Out Weight to proceed</span>';
    delete corrections[idx];
    _updateRRPanelWeightLabel(idx,0);
    return;
  }
  const net=Math.round((inW-outW)*1000)/1000;
  corrections[idx]={Accepted:'YES','Out Weight':outW,'Net Weight':net};
  statusEl.innerHTML=`<span class="chip-ok">✓ Out: ${outW}T, Net: ${net}T</span>`;
  // Show RR panel & update weight label — never destroys the serial input
  _ensureRRPanel(idx,inW,'YES',outW,net);
  _updateRRPanelWeightLabel(idx,outW);
  // NOTE: do NOT call scheduleValidation() here — the blur handler does that
  // so typing never triggers a re-render that kills focus
}

// ── RR Panel ──────────────────────────────────────────────────────────────────
function _ensureRRPanel(idx,inW,decision,outW,net){
  const panel=document.getElementById(`rr-panel-${idx}`);
  if(!panel)return;

  // Already rendered — just update weight label, preserve all inputs
  if(panel.querySelector('.rr-panel')){
    _updateRRPanelWeightLabel(idx,outW);
    return;
  }

  const savedSerial=_rrSerials[idx]||'';

  // All three options always visible; default pre-selection based on decision
  const reasonOptions=`
    <option value="Accepted Towing Vehicle" ${decision==='YES'?'selected':''}>Option A — Accepted Towing Vehicle</option>
    <option value="Rejected Towing Vehicle" ${decision==='NO'?'selected':''}>Option B — Rejected Towing Vehicle</option>
    <option value="Late Time / Breakdown">Option C — Late Time / Breakdown</option>`;

  const wrapper=document.createElement('div');
  wrapper.className='rr-panel';
  wrapper.innerHTML=`
    <div class="rr-panel-title">📄 Rectification Report — ${decision==='YES'?'✅ Acceptance Format':'❌ Rejection Format'}</div>

    <div style="margin:.6rem 0 .25rem;font-size:11px;color:#9BB4D4;font-weight:700;letter-spacing:.07em;text-transform:uppercase;">📋 Reason / Narrative Type</div>
    <select class="rr-select" id="rr-reason-${idx}"
      style="font-size:13px;padding:.55rem .9rem;margin-bottom:.6rem;border:1.5px solid rgba(99,102,241,.5);background:rgba(6,10,28,.9);color:#E8F0FF;border-radius:9px;width:100%;">
      ${reasonOptions}
    </select>

    <div style="margin:.3rem 0 .25rem;font-size:11px;color:#9BB4D4;font-weight:700;letter-spacing:.07em;text-transform:uppercase;">🔢 Serial Number</div>
    <input class="rr-input" id="rr-serial-${idx}" type="text"
      placeholder="Enter RR Serial No. (e.g. 0290)" maxlength="10"
      value="${esc(savedSerial)}" autocomplete="off" spellcheck="false"
      style="font-size:14px;padding:.6rem .9rem;border:1.5px solid rgba(99,102,241,.5);background:rgba(6,10,28,.9);color:#E8F0FF;border-radius:9px;width:100%;margin-bottom:.6rem;">



    <button class="rr-btn" id="rr-btn-${idx}" type="button"
      style="width:100%;padding:.7rem;margin-top:.4rem;font-size:13px;font-weight:700;letter-spacing:.02em;border-radius:9px;">
      ⬇ Download Rectification Report
    </button>
    <div class="rr-status" id="rr-status-${idx}" style="font-size:13px;margin-top:.45rem;min-height:1.4rem;"></div>`;

  panel.appendChild(wrapper);

  // Attach events AFTER DOM insertion
  const serialEl=panel.querySelector(`#rr-serial-${idx}`);
  if(serialEl){
    serialEl.addEventListener('input',function(){_rrSerials[idx]=this.value;});
    serialEl.addEventListener('keydown',function(e){e.stopPropagation();});
    serialEl.addEventListener('click',function(e){e.stopPropagation();});
    serialEl.addEventListener('mousedown',function(e){e.stopPropagation();});
  }
  const reasonDropEl=panel.querySelector(`#rr-reason-${idx}`);
  if(reasonDropEl){
    reasonDropEl.addEventListener('keydown',function(e){e.stopPropagation();});
    reasonDropEl.addEventListener('click',function(e){e.stopPropagation();});
  }
  const btn=panel.querySelector(`#rr-btn-${idx}`);
  if(btn)btn.addEventListener('click',function(e){e.stopPropagation();generateRR(idx);});
}
function _updateRRPanelWeightLabel(idx,outW){
  // Weight reference UI removed — computed server-side from reason
}

// ── Generate Rectification Report ─────────────────────────────────────────────
async function generateRR(idx){
  const item=_incompleteItems[idx]||{};
  const token=item.token||'';
  const corr=corrections[idx]||{};
  const decision=corr.Accepted||'';
  const outW=corr['Out Weight']||0;
  const net=corr['Net Weight']||0;
  const serialEl=document.getElementById(`rr-serial-${idx}`);
  const serial=(serialEl?serialEl.value:'').trim();
  // Read reason dropdown (weight_label computed server-side)  // Read reason dropdown
  const reasonEl=document.getElementById(`rr-reason-${idx}`);
  const reason=reasonEl?reasonEl.value:(decision==='YES'?'Accepted Towing Vehicle':'Rejected Towing Vehicle');

  const statusEl=document.getElementById(`rr-status-${idx}`);
  const btn=document.getElementById(`rr-btn-${idx}`);

  if(!serial){if(statusEl)statusEl.innerHTML='<span style="color:#F59E0B">⚠ Enter RR Serial Number first</span>';return;}
  if(!decision){if(statusEl)statusEl.innerHTML='<span style="color:#F59E0B">⚠ Select YES or NO first</span>';return;}
  if(!onlineTmpPath){if(statusEl)statusEl.innerHTML='<span style="color:#F87171">❌ Online file not loaded</span>';return;}

  if(btn){btn.disabled=true;btn.textContent='⏳ Generating…';}
  if(statusEl)statusEl.innerHTML='<span style="color:#818CF8">⏳ Building report…</span>';
  try{
    const res=await fetch('/api/daily/generate_rr',{
      method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({
        tmp_path:onlineTmpPath,token,rr_serial:serial,accepted:decision,
        out_weight:outW,net_weight:net,
        reason:reason,filter_date:onlineDate
      })
    });
    if(!res.ok){
      const e=await res.json().catch(()=>({error:`HTTP ${res.status}`}));
      if(statusEl)statusEl.innerHTML=`<span style="color:#F87171">❌ ${esc(e.error||'Failed')}</span>`;
      return;
    }
    const blob=await res.blob();
    const cd=res.headers.get('Content-Disposition')||'';
    // Try quoted filename first, then unquoted
    const fnMatch=cd.match(/filename="([^"]+)"/) || cd.match(/filename=([^;\s]+)/);
    const fname=fnMatch?decodeURIComponent(fnMatch[1]):`RR-B-${serial}-${token}.docx`;
    const url=URL.createObjectURL(blob);
    const a=document.createElement('a');a.href=url;a.download=fname;a.click();
    URL.revokeObjectURL(url);
    if(statusEl)statusEl.innerHTML=`<span style="color:#4ADE80">✅ Downloaded: ${esc(fname)}</span>`;
  }catch(ex){
    if(statusEl)statusEl.innerHTML=`<span style="color:#F87171">❌ ${esc(ex.message)}</span>`;
  }finally{
    if(btn){btn.disabled=false;btn.textContent='⬇ Download Rectification Report';}
  }
}
// ── Stats / tracker ───────────────────────────────────────────────────────────
function renderOnlineStats(s){
  setEl('onlineStats',`<div class="stats-box"><div class="stats-title stats-title-online">Online Data — Summary</div>
    <div class="stats-inner"><table>
      <tr><td>Accepted Loads</td><td style="color:#818CF8">${s.accepted}</td></tr>
      <tr><td>Rejected Loads</td><td style="color:#F87171">${s.rejected}</td></tr>
      <tr><td>Total Weight In (T)</td><td style="color:#10B981">${s.wi.toFixed(2)}</td></tr>
      <tr><td>Total Weight Out (T)</td><td style="color:#10B981">${s.wo.toFixed(2)}</td></tr>
      <tr><td>Total Net Weight (T)</td><td style="color:#10B981;font-weight:800">${s.nw.toFixed(2)}</td></tr>
    </table></div></div>`);show('onlineStats');
}

function updateStepTracker(){
  const s=(ok,l)=>`<span class="step-pip ${ok?'step-done':'step-wait'}">${ok?'✔':'○'} ${l}</span>`;
  setEl('stepTracker',`${s(!!onlineStats,'Online Validated')}<span class="step-arr">→</span>${s(!!wbStats,'WB Compared')}<span class="step-arr">→</span>${s(!!onlineStats&&!!wbStats,'Ready to Download')}`);
}
function updateGenerateBtn(){
  const btn=document.getElementById('generateBtn');
  if(btn)btn.disabled=!(onlineTmpPath&&onlineDate&&onlineStats);
}

// ── WB Side ───────────────────────────────────────────────────────────────────
async function onWbFileChange(){
  const f=document.getElementById('wbFile').files[0];if(!f)return;
  setEl('wbFileStatus','<span class="chip-info">⏳ Uploading…</span>');
  const fd=new FormData();fd.append('file',f);
  try{
    const res=await fetch('/api/daily/wb_upload',{method:'POST',body:fd});
    const d=await res.json();
    if(d.error){setEl('wbFileStatus',`<span class="chip-wait">❌ ${esc(d.error)}</span>`);return;}
    wbTmpPath=d.tmp_path;wbDecisions={};
    setEl('wbFileStatus',`<span class="chip-ok">✓ ${esc(f.name)} — ${d.rows.toLocaleString()} rows loaded</span>`);
    const di=document.getElementById('wbDate');
    di.min=d.min_date;di.max=d.max_date;di.value=d.max_date;wbDate=d.max_date;
    show('wbDateSection');show('wbResultSection');await runWbPivot();
  }catch(e){setEl('wbFileStatus',`<span class="chip-wait">❌ ${e.message}</span>`);}
}

async function onWbDateChange(){wbDate=document.getElementById('wbDate').value;wbDecisions={};await runWbPivot();}

function scheduleWbPivot(){
  clearTimeout(_wbPivotTimer);
  _wbPivotTimer=setTimeout(function(){
    if(_userIsTyping()){scheduleWbPivot();return;}
    runWbPivot();
  },900);
}

async function runWbPivot(){
  if(!wbTmpPath||!wbDate)return;
  setEl('wbStats','<div class="chip-info">⏳ Processing WB data…</div>');
  try{
    const res=await fetch('/api/daily/wb_pivot',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({
        wb_tmp_path:wbTmpPath,filter_date:wbDate,wb_decisions:wbDecisions,
        online_tmp_path:onlineTmpPath,online_date:onlineDate,corrections
      })});
    const d=await res.json();
    if(d.error){setEl('wbStats',`<div class="val-err">❌ ${esc(d.error)}</div>`);return;}
    if(d.incomplete&&d.incomplete.length>0){renderWbActionCards(d.incomplete);return;}
    wbStats=d.wb_stats;pivotTmpPath=d.pivot_tmp;pivotFname=d.pivot_fname;
    if(d.errors&&d.errors.length>0){
      setEl('wbStats',d.errors.map(e=>`<div class="val-err">❌ ${esc(e)}</div>`).join('')+renderWbStatsHtml(d.wb_stats));
    }else{
      setEl('wbStats','<div class="val-ok">✅ All validations completed successfully</div>'+renderWbStatsHtml(d.wb_stats));
    }
    renderComparisonTable();renderPivot(d.rows);
    show('pivotDlBtn');updateStepTracker();
  }catch(e){setEl('wbStats',`<div class="val-err">❌ ${e.message}</div>`);}
}

function renderWbStatsHtml(s){
  if(!s||typeof s.accepted==='undefined')return '';
  return `<div class="stats-box stats-box-wb"><div class="stats-title stats-title-wb">Weighbridge Data — Summary</div>
    <div class="stats-inner"><table>
      <tr><td>Accepted Loads</td><td style="color:#34D399">${s.accepted}</td></tr>
      <tr><td>Rejected Loads</td><td style="color:#F87171">${s.rejected}</td></tr>
      <tr><td>Total Weight In (T)</td><td style="color:#10B981">${(s.wi||0).toFixed(2)}</td></tr>
      <tr><td>Total Weight Out (T)</td><td style="color:#10B981">${(s.wo||0).toFixed(2)}</td></tr>
      <tr><td>Total Net Weight (T)</td><td style="color:#10B981;font-weight:800">${(s.nw||0).toFixed(2)}</td></tr>
    </table></div></div>`;
}

function renderWbActionCards(incomplete){
  if(!incomplete||!incomplete.length)return;
  let html='<div class="val-warn">⚠️ Action Required — WB rows with missing Date Out / Time Out</div>';
  for(const item of incomplete){
    const dec=wbDecisions[item.idx]||{};
    html+=`<div class="action-card" id="wb-acard-${item.idx}">
      <div class="action-title">⚠️ E-Token [${esc(item.etoken)}] – Missing Date Out / Time Out</div>
      <select class="action-select" id="wb-dec-${item.idx}" onchange="onWbDecisionChange(${item.idx},${item.in_weight})">
        <option value="" ${!dec.decision&&dec.decision!==0?'selected':''}>— Select —</option>
        <option value="0" ${dec.decision===0?'selected':''}>0 — Rejected (remove row)</option>
        <option value="1" ${dec.decision===1?'selected':''}>1 — Accepted (enter Out Weight)</option>
      </select>
      <div id="wb-dec-detail-${item.idx}"></div>
    </div>`;
  }
  setEl('wbStats',html);
  for(const item of incomplete){if(wbDecisions[item.idx])renderWbDecisionDetail(item.idx,item.in_weight,wbDecisions[item.idx].decision);}
}

function onWbDecisionChange(idx,inW){
  const v=document.getElementById(`wb-dec-${idx}`).value;
  renderWbDecisionDetail(idx,inW,v===''?null:parseInt(v));
}

function renderWbDecisionDetail(idx,inW,decision){
  const el=document.getElementById(`wb-dec-detail-${idx}`);if(!el)return;
  if(decision===0||decision==='0'){
    wbDecisions[idx]={decision:0};
    el.innerHTML='<span class="chip-ok">✓ Marked Rejected — row will be removed.</span>';
    scheduleWbPivot();
  }else if(decision===1||decision==='1'){
    if(!document.getElementById(`wb-outw-${idx}`)){
      const prev=(wbDecisions[idx]&&wbDecisions[idx].out_weight)||0;
      el.innerHTML=`<input type="number" class="input-num" id="wb-outw-${idx}"
        placeholder="Out Weight (T) [In Weight=${inW}T]" min="0" max="${inW||9999}"
        step="0.001" value="${prev||''}" autocomplete="off">
        <div id="wb-outw-s-${idx}" style="margin-top:.3rem;font-size:.78rem"></div>`;
      const inp=document.getElementById(`wb-outw-${idx}`);
      inp.addEventListener('input',function(){onWbOutW(idx,inW,parseFloat(this.value)||0);});
      inp.addEventListener('blur',function(){
        const v=parseFloat(this.value)||0;
        if(v>0)scheduleWbPivot();
      });
      inp.addEventListener('keydown',function(e){e.stopPropagation();});
      if(prev>0)_applyWbOutWValue(idx,inW,prev,true);
    }
  }else{delete wbDecisions[idx];el.innerHTML='';}
}

function onWbOutW(idx,inW,outW){_applyWbOutWValue(idx,inW,outW,false);}
function _applyWbOutWValue(idx,inW,outW,restoring){
  const el=document.getElementById(`wb-outw-s-${idx}`);if(!el)return;
  if(outW<=0){el.innerHTML='<span style="color:#F59E0B">Enter Out Weight to proceed</span>';delete wbDecisions[idx];return;}
  const net=Math.round((inW-outW)*1000)/1000;
  wbDecisions[idx]={decision:1,out_weight:outW};
  el.innerHTML=`<span class="chip-ok">✓ Out: ${outW}T, Net: ${net}T</span>`;
  // validation only on blur, not here
}

function renderWbStats(s){setEl('wbStats',renderWbStatsHtml(s));}

function renderComparisonTable(){
  if(!onlineStats||!wbStats)return;
  const on=onlineStats,wb=wbStats;
  const m=(a,b,wt)=>Math.abs(a-b)<(wt?.01:.5)?`<td class="cmp-match">✔ Validated</td>`:`<td class="cmp-miss">✘ Mismatch</td>`;
  setEl('comparisonTable',`<table class="cmp-table">
    <thead><tr><th style="text-align:left">Metric</th><th>Online</th><th>Weighbridge</th><th>Remark</th></tr></thead>
    <tbody>
      <tr><td style="text-align:left;color:#C8D8F8">Accepted Loads</td><td style="color:#818CF8;font-weight:700">${on.accepted}</td><td style="color:#34D399;font-weight:700">${wb.accepted}</td>${m(on.accepted,wb.accepted,false)}</tr>
      <tr><td style="text-align:left;color:#C8D8F8">Rejected Loads</td><td style="color:#818CF8;font-weight:700">${on.rejected}</td><td style="color:#34D399;font-weight:700">${wb.rejected}</td>${m(on.rejected,wb.rejected,false)}</tr>
      <tr><td style="text-align:left;color:#C8D8F8">Weight In (T)</td><td style="color:#818CF8;font-weight:700">${on.wi.toFixed(2)}</td><td style="color:#34D399;font-weight:700">${wb.wi.toFixed(2)}</td>${m(on.wi,wb.wi,true)}</tr>
      <tr><td style="text-align:left;color:#C8D8F8">Weight Out (T)</td><td style="color:#818CF8;font-weight:700">${on.wo.toFixed(2)}</td><td style="color:#34D399;font-weight:700">${wb.wo.toFixed(2)}</td>${m(on.wo,wb.wo,true)}</tr>
      <tr><td style="text-align:left;color:#C8D8F8">Net Weight (T)</td><td style="color:#818CF8;font-weight:700;font-size:.95rem">${on.nw.toFixed(2)}</td><td style="color:#34D399;font-weight:700;font-size:.95rem">${wb.nw.toFixed(2)}</td>${m(on.nw,wb.nw,true)}</tr>
    </tbody></table>`);show('comparisonTable');
}

function renderPivot(rows){
  if(!rows||!rows.length){clearEl('pivotSection');return;}
  let html=`<div class="pivot-outer"><div class="pivot-title">📊 WB Pivot — Net Weight by Material & Site</div><table class="pivot-tbl"><thead><tr>`;
  const cols=Object.keys(rows[0]);
  cols.forEach(c=>{html+=`<th>${esc(c)}</th>`;});
  html+=`</tr></thead><tbody>`;
  rows.forEach(r=>{
    const cls=r['SITE CODE']==='GRAND TOTAL'?'pivot-grand':r['SITE CODE']===''?'pivot-mat':'';
    html+=`<tr class="${cls}">`;
    cols.forEach(c=>{html+=`<td>${esc(String(r[c]??''))}</td>`;});
    html+=`</tr>`;
  });
  html+=`</tbody></table></div>`;
  setEl('pivotSection',html);
}

async function generateReport(){
  const btn=document.getElementById('generateBtn');
  if(btn){btn.disabled=true;btn.textContent='⏳ Generating…';}
  setEl('generateStatus','<span style="color:#818CF8">⏳ Generating report…</span>');
  try{
    const res=await fetch('/api/daily/generate',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({tmp_path:onlineTmpPath,filter_date:onlineDate,corrections})});
    if(!res.ok){const e=await res.json().catch(()=>({error:`HTTP ${res.status}`}));setEl('generateStatus',`<span style="color:#F87171">❌ ${esc(e.error||'Failed')}</span>`);return;}
    const blob=await res.blob();
    const cd=res.headers.get('Content-Disposition')||'';
    const fnMatch=cd.match(/filename="([^"]+)"/);
    reportName=fnMatch?fnMatch[1]:'daily_report.xlsx';
    reportBlob=blob;
    setEl('generateStatus',`<span style="color:#4ADE80">✅ Ready: ${esc(reportName)}</span>`);
    show('downloadBtn');document.getElementById('downloadBtn').disabled=false;
  }catch(e){setEl('generateStatus',`<span style="color:#F87171">❌ ${e.message}</span>`);}
  finally{if(btn){btn.disabled=false;btn.textContent='⚡ Generate Report';}}
}

function downloadReport(){
  if(!reportBlob)return;
  const url=URL.createObjectURL(reportBlob);
  const a=document.createElement('a');a.href=url;a.download=reportName;a.click();
  URL.revokeObjectURL(url);
}

function downloadPivot(){if(!pivotTmpPath)return;window.location.href=`/api/daily/wb_pivot_download?path=${encodeURIComponent(pivotTmpPath)}&fname=${encodeURIComponent(pivotFname)}`;}

function setEl(id,html){const el=document.getElementById(id);if(el)el.innerHTML=html;}
function clearEl(id){const el=document.getElementById(id);if(el)el.innerHTML='';}
function show(id){const el=document.getElementById(id);if(el)el.style.display='block';}
function hide(id){const el=document.getElementById(id);if(el)el.style.display='none';}
function esc(s){return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""


ADMIN_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Admin Panel — APSG Report</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap" rel="stylesheet">
<style>

/* ═══ MODERN UI BASE (Blue/Purple/Cyan Theme) ═══ */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
:root {
  --bg:transparent; --card-bg: rgba(13,18,35,0.88);
  --indigo: #6366F1; --indigo-l: #818CF8; --cyan: #22D3EE;
  --purple: #A855F7; --green: #10B981; --amber: #F59E0B; --red: #F87171;
  --text: #E8EEF8; --muted: #64748B; --border: rgba(99,102,241,0.15);
}
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent; color: var(--text);
  position: relative;
}
body::before {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none;
  background:
    radial-gradient(ellipse 70% 55% at 5% 10%, rgba(99,102,241,.09) 0%, transparent 65%),
    radial-gradient(ellipse 55% 45% at 95% 90%, rgba(168,85,247,.07) 0%, transparent 65%);
}
body::after {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none; opacity: .016;
  background-image: linear-gradient(var(--indigo) 1px, transparent 1px),
    linear-gradient(90deg, var(--indigo) 1px, transparent 1px);
  background-size: 60px 60px;
}

/* ── Modern Top Bar ── */
.top-bar {
  position: sticky; top: 0; z-index: 200; height: 56px;
  display: flex; align-items: center; padding: 0 1.5rem; gap: .75rem;
  background: rgba(6,9,22,.96); backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px);
  border-bottom: 1px solid var(--border);
}
.top-mini-brand {
  font-family: 'Poppins', sans-serif;
  font-size: .82rem; font-weight: 700; color: var(--text);
  letter-spacing: -.01em; white-space: nowrap; flex-shrink: 0;
  display: flex; align-items: center; gap: .45rem;
}
.top-mini-brand::before {
  content: ''; width: 8px; height: 8px; border-radius: 50%;
  background: linear-gradient(135deg, var(--indigo), var(--cyan)); flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: var(--border); flex-shrink: 0; }
.top-page-label { font-size: .75rem; font-weight: 600; color: var(--muted); white-space: nowrap; }
.top-brand-tag { font-size: .7rem; color: var(--muted); }
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.08); border: 1px solid rgba(99,102,241,.18);
  border-radius: 8px; padding: .3rem .9rem; font-size: .7rem; font-weight: 600;
  color: var(--indigo-l); text-decoration: none; transition: all .2s; white-space: nowrap;
}
.back-btn:hover { background: rgba(99,102,241,.18); transform: translateX(-2px); }
/* ── Cards & Containers ── */
.container, .page-content { position: relative; z-index: 1; }
.card, .section-card, .panel {
  background: var(--card-bg); border: 1px solid var(--border);
  border-radius: 16px; backdrop-filter: blur(16px); -webkit-backdrop-filter: blur(16px);
  box-shadow: 0 8px 32px rgba(0,0,0,.4), inset 0 1px 0 rgba(255,255,255,.03);
  transition: box-shadow .25s, border-color .25s;
}
.card:hover, .section-card:hover { border-color: rgba(99,102,241,.25); }

/* ── Upload Zone — Modern drag & drop ── */
.upload-zone, .dz {
  border: 2px dashed rgba(99,102,241,.3); border-radius: 14px;
  padding: 2rem; text-align: center; cursor: pointer;
  transition: all .22s; background: rgba(99,102,241,.03);
  position: relative;
}
.upload-zone:hover, .dz:hover, .upload-zone.drag-over, .dz.drag-over {
  border-color: var(--indigo); background: rgba(99,102,241,.08);
  box-shadow: 0 0 0 4px rgba(99,102,241,.12);
}
.upload-zone.ok, .dz.ok {
  border-color: var(--green); background: rgba(16,185,129,.06);
  border-style: solid;
}
.upload-zone.ok:hover, .dz.ok:hover {
  border-color: var(--green); background: rgba(16,185,129,.1);
  box-shadow: 0 0 0 4px rgba(16,185,129,.1);
}
.upload-icon { font-size: 2rem; margin-bottom: .5rem; display: block; }
.upload-label { font-size: .82rem; color: var(--muted); font-weight: 500; }
.upload-hint { font-size: .7rem; color: rgba(100,116,139,.6); margin-top: .25rem; }
.upload-filename { font-size: .78rem; color: var(--green); font-weight: 600; margin-top: .4rem; }

/* ── Modern Buttons ── */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  background:transparent;
  color: #fff; border: none; border-radius: 10px;
  padding: .7rem 1.4rem; font-size: .85rem; font-weight: 700;
  font-family: 'Inter', sans-serif; cursor: pointer; letter-spacing: .01em;
  box-shadow: 0 4px 18px rgba(99,102,241,.35);
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s;
}
.btn-primary::before, .btn-generate::before, .btn-teal::before, .modal-btn-primary::before {
  content: ''; position: absolute; inset: 0;
  background: linear-gradient(135deg, transparent, rgba(255,255,255,.12), transparent);
  transform: translateX(-100%); transition: transform .45s;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-2px) scale(1.01);
  box-shadow: 0 8px 28px rgba(99,102,241,.55);
}
.btn-primary:hover::before, .btn-generate:hover::before, .modal-btn-primary:hover::before {
  transform: translateX(100%);
}
.btn-primary:active, .btn-generate:active { transform: translateY(0) scale(.98); }
.btn-primary:disabled, .btn-generate:disabled { opacity: .45; cursor: not-allowed; transform: none; }

.btn-teal {
  background:transparent;
  box-shadow: 0 4px 18px rgba(16,185,129,.3);
}
.btn-teal:hover { box-shadow: 0 8px 28px rgba(16,185,129,.5); }

.btn-download, .btn-dl, .dl-btn {
  background: linear-gradient(135deg, #0D7A5F, #10B981);
  color: #fff; border: none; border-radius: 10px;
  padding: .65rem 1.3rem; font-size: .82rem; font-weight: 700;
  cursor: pointer; transition: all .22s;
  box-shadow: 0 4px 16px rgba(16,185,129,.3);
}
.btn-download:hover, .btn-dl:hover, .dl-btn:hover {
  transform: translateY(-2px); box-shadow: 0 8px 26px rgba(16,185,129,.5);
}

.btn-secondary, .btn-gray {
  background: rgba(30,41,86,.6); color: var(--muted);
  border: 1px solid rgba(99,102,241,.2); border-radius: 10px;
  padding: .65rem 1.2rem; font-size: .82rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-secondary:hover, .btn-gray:hover { background: rgba(99,102,241,.1); color: var(--indigo-l); }

.btn-del, .btn-danger {
  background: rgba(239,68,68,.08); color: var(--red);
  border: 1px solid rgba(239,68,68,.2); border-radius: 10px;
  padding: .6rem 1.2rem; font-size: .8rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-del:hover, .btn-danger:hover { background: rgba(239,68,68,.16); }
.btn-del:hover { animation: shake .3s ease; }

@keyframes shake {
  0%,100% { transform: translateX(0); }
  25% { transform: translateX(-3px); }
  75% { transform: translateX(3px); }
}

/* ── Form Inputs ── */
input[type=text], input[type=number], input[type=date], input[type=password],
select, textarea {
  background: rgba(6,10,24,.8); border: 1.5px solid rgba(30,41,86,.8);
  border-radius: 10px; color: var(--text); padding: .65rem .9rem;
  font-size: .85rem; font-family: 'Inter', sans-serif;
  transition: all .2s; width: 100%;
}
input:focus, select:focus, textarea:focus {
  outline: none; border-color: var(--indigo);
  box-shadow: 0 0 0 3px rgba(99,102,241,.16);
  background: rgba(10,14,32,.9);
}
input::placeholder, textarea::placeholder { color: rgba(100,116,139,.45); font-weight: 300; }

/* ── Alerts / Status ── */
.alert-success, .alert.success, .msg.success {
  background: rgba(16,185,129,.08); border: 1px solid rgba(16,185,129,.2);
  color: #34D399; border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-error, .alert.error, .msg.error {
  background: rgba(239,68,68,.08); border: 1px solid rgba(239,68,68,.2);
  color: var(--red); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-warn, .alert.warn { 
  background: rgba(245,158,11,.08); border: 1px solid rgba(245,158,11,.2);
  color: var(--amber); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}

/* ── Tables ── */
table { width: 100%; border-collapse: collapse; font-size: .8rem; }
thead th {
  background: rgba(99,102,241,.08); color: var(--muted);
  font-weight: 700; font-size: .68rem; letter-spacing: .06em;
  text-transform: uppercase; padding: .65rem .9rem; text-align: left;
  border-bottom: 1px solid var(--border);
}
tbody tr { border-bottom: 1px solid rgba(30,45,80,.3); transition: background .15s; }
tbody tr:hover { background: rgba(99,102,241,.04); }
tbody td { padding: .6rem .9rem; color: var(--text); }

/* ── Tabs ── */
.tab { position: relative; transition: all .2s; }
.tab.active { color: var(--indigo-l) !important; }
.tab.active::after {
  content: ''; position: absolute; bottom: -1px; left: 15%; right: 15%;
  height: 2px; border-radius: 2px; background: var(--indigo);
  animation: tabIn .2s ease;
}
@keyframes tabIn { from { left: 50%; right: 50%; } to { left: 15%; right: 15%; } }

/* ── Spinner ── */
@keyframes spin { to { transform: rotate(360deg); } }
.spinner { width: 20px; height: 20px; border: 2px solid rgba(99,102,241,.2);
  border-top-color: var(--indigo); border-radius: 50%; animation: spin 1s linear infinite; }



@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent;
  color: #E8EEF8;
}


/* ══ ANIMATED BUTTONS ══════════════════════════════════════════════════════════ */
/* Primary action button — pulse + shimmer */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-primary::after, .btn-generate::after, .btn-teal::after, .modal-btn-primary::after {
  content: ''; position: absolute; top: -50%; left: -75%;
  width: 50%; height: 200%; background: rgba(255,255,255,.15);
  transform: skewX(-20deg); transition: left .5s ease;
  pointer-events: none;
}
.btn-primary:hover::after, .btn-generate:hover::after,
.btn-teal:hover::after, .modal-btn-primary:hover::after {
  left: 150%;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(99,102,241,.55) !important;
}
.btn-primary:active, .btn-generate:active, .modal-btn-primary:active {
  transform: translateY(0) scale(.98) !important;
}

/* Download button — bounce */
.btn-dl, .btn-download, .dl-btn {
  animation: gentlePulse 3s ease-in-out infinite;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-dl:hover, .btn-download:hover, .dl-btn:hover {
  animation: none;
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(16,185,129,.5) !important;
}
@keyframes gentlePulse {
  0%,100% { box-shadow: 0 4px 16px rgba(16,185,129,.3); }
  50%      { box-shadow: 0 4px 24px rgba(16,185,129,.55); }
}

/* Danger/delete button — shake on hover */
.btn-del:hover {
  animation: shake .35s ease;
}
@keyframes shake {
  0%,100% { transform: translateX(0); }
  20%     { transform: translateX(-3px); }
  40%     { transform: translateX(3px); }
  60%     { transform: translateX(-2px); }
  80%     { transform: translateX(2px); }
}

/* Secondary / back buttons — slide arrow */
.back-btn, .btn-secondary, .btn-gray, .btn-refresh {
  transition: all .2s !important;
}
.back-btn:hover {
  padding-left: .55rem !important;
  letter-spacing: .02em;
}

/* Upload zone — glow pulse when empty */
.upload-zone:not(.ok) {
  animation: uploadGlow 4s ease-in-out infinite;
}
.upload-zone:hover, .dz:hover {
  animation: none;
}
@keyframes uploadGlow {
  0%,100% { border-color: rgba(99,102,241,.3); }
  50%      { border-color: rgba(99,102,241,.65); box-shadow: 0 0 18px rgba(99,102,241,.15); }
}

/* App cards on dashboard — float in on load */
.app-card {
  animation: cardFadeIn .5s ease both;
}
.app-card:nth-child(1) { animation-delay: .05s; }
.app-card:nth-child(2) { animation-delay: .10s; }
.app-card:nth-child(3) { animation-delay: .15s; }
.app-card:nth-child(4) { animation-delay: .20s; }
.app-card:nth-child(5) { animation-delay: .25s; }
@keyframes cardFadeIn {
  from { opacity: 0; transform: translateY(18px); }
  to   { opacity: 1; transform: translateY(0); }
}

/* Stat cards pop-in */
.stat-card, .stat {
  animation: popIn .4s cubic-bezier(.34,1.56,.64,1) both;
}
.stat-card:nth-child(1), .stat:nth-child(1) { animation-delay: .05s; }
.stat-card:nth-child(2), .stat:nth-child(2) { animation-delay: .10s; }
.stat-card:nth-child(3), .stat:nth-child(3) { animation-delay: .15s; }
.stat-card:nth-child(4), .stat:nth-child(4) { animation-delay: .20s; }
@keyframes popIn {
  from { opacity: 0; transform: scale(.88); }
  to   { opacity: 1; transform: scale(1); }
}

/* Login button — ripple */
.btn.btn-primary, button.btn-primary {
  transition: transform .15s, box-shadow .15s !important;
}

/* Tab buttons — active indicator slide */
.tab {
  transition: all .2s !important;
  position: relative;
}
.tab.active::after {
  content: '';
  position: absolute; bottom: -1px; left: 10%; right: 10%;
  height: 2px; border-radius: 2px;
  background: #6366F1;
  animation: tabSlide .25s ease;
}
@keyframes tabSlide {
  from { left: 50%; right: 50%; }
  to   { left: 10%; right: 10%; }
}

/* Generate/action buttons in forms — shimmer idle */
.btn-generate:not(:disabled), .btn-teal:not(:disabled) {
  background-size: 200% 100% !important;
  transition: background-position .4s, transform .2s, box-shadow .2s !important;
}

/* ── Compact top bar (40px) for all internal pages ── */
.top-bar {
  position: sticky; top: 0; z-index: 200;
  height: 40px; display: flex; align-items: center;
  padding: 0 1.2rem; gap: .75rem;
  background: rgba(5,8,20,.96); backdrop-filter: blur(16px);
  border-bottom: 1px solid rgba(99,102,241,.12);
}
.top-mini-brand {
  font-size: .75rem; font-weight: 800; color: #6366F1;
  letter-spacing: .04em; white-space: nowrap; flex-shrink: 0;
}
.top-brand-text {
  font-size: .78rem; font-weight: 800; color: #818CF8;
  letter-spacing: .02em; white-space: nowrap; flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: rgba(99,102,241,.18); flex-shrink: 0; }
.top-page-label {
  font-size: .78rem; font-weight: 600; color: #94A3B8;
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
}
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.1); border: 1px solid rgba(99,102,241,.2);
  border-radius: 6px; padding: .22rem .7rem; font-size: .7rem; font-weight: 600;
  color: #818CF8; text-decoration: none; white-space: nowrap; transition: all .2s;
  flex-shrink: 0;
}
.back-btn:hover { background: rgba(99,102,241,.2); }

.container, .main, .page { padding: 1rem 1.2rem 2rem; max-width: 1400px; margin: 0 auto; }


.page-content{padding:.75rem 1.2rem 2.5rem;max-width:1200px;margin:0 auto;}
.admin-hero{display:flex;align-items:center;gap:.8rem;margin-bottom:1.3rem;padding:.9rem 1.2rem;
  background:rgba(245,158,11,.06);border:1px solid rgba(245,158,11,.13);border-radius:12px;}
.admin-hero-icon{font-size:1.6rem;}
.admin-hero-title{font-size:1rem;font-weight:800;color:#F1F5FF;}
.admin-hero-sub{font-size:.74rem;color:#475569;margin-top:.1rem;}
.tabs{display:flex;gap:.4rem;margin-bottom:1.2rem;flex-wrap:wrap;}
.tab{padding:.42rem 1rem;border-radius:8px;border:1px solid rgba(99,102,241,.18);background:transparent;
  color:#475569;cursor:pointer;font-size:.78rem;font-weight:600;font-family:'Inter',sans-serif;transition:all .2s;}
.tab:hover{background:rgba(99,102,241,.06);color:#818CF8;}
.tab.active{background:rgba(99,102,241,.14);border-color:rgba(99,102,241,.35);color:#818CF8;}
.panel{display:none;}.panel.active{display:block;}
.stats-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(140px,1fr));gap:.85rem;margin-bottom:1.2rem;}
.stat-card{background:rgba(13,19,40,.8);border:1px solid rgba(99,102,241,.14);border-radius:11px;padding:.9rem;text-align:center;}
.stat-num{font-size:2rem;font-weight:900;color:#6366F1;line-height:1.1;}
.stat-num.green{color:#10B981;}.stat-num.amber{color:#F59E0B;}
.stat-lbl{font-size:.64rem;color:#475569;text-transform:uppercase;letter-spacing:.08em;margin-top:.25rem;font-weight:500;}
.card{background:rgba(13,19,40,.8);border:1px solid rgba(99,102,241,.12);border-radius:12px;padding:1.1rem;margin-bottom:.9rem;}
.card-header{display:flex;align-items:center;justify-content:space-between;margin-bottom:.85rem;flex-wrap:wrap;gap:.4rem;}
.card-title{font-size:.84rem;font-weight:700;color:#818CF8;}
.btn-action{padding:.38rem .8rem;border-radius:7px;border:none;cursor:pointer;font-size:.72rem;font-weight:600;font-family:'Inter',sans-serif;transition:.2s;text-decoration:none;display:inline-flex;align-items:center;gap:.3rem;}
.btn-dl{background:rgba(16,185,129,.1);color:#10B981;border:1px solid rgba(16,185,129,.2);}
.btn-dl:hover{background:rgba(16,185,129,.18);}
.btn-refresh{background:rgba(99,102,241,.1);color:#818CF8;border:1px solid rgba(99,102,241,.2);}
.btn-refresh:hover{background:rgba(99,102,241,.18);}
.tbl-wrap{overflow-x:auto;border-radius:8px;border:1px solid rgba(99,102,241,.1);}
table{width:100%;border-collapse:collapse;font-size:.76rem;}
thead{background:rgba(99,102,241,.1);}
thead th{padding:.55rem .8rem;text-align:left;font-weight:700;color:#818CF8;white-space:nowrap;border-bottom:1px solid rgba(99,102,241,.18);font-size:.72rem;letter-spacing:.02em;}
tbody tr{border-bottom:1px solid rgba(255,255,255,.04);transition:background .15s;}
tbody tr:hover{background:rgba(99,102,241,.04);}
tbody td{padding:.48rem .8rem;color:#C8D8F8;}
.role-admin{color:#F59E0B;font-weight:700;font-size:.7rem;background:rgba(245,158,11,.1);padding:.12rem .45rem;border-radius:5px;border:1px solid rgba(245,158,11,.2);}
.role-user{color:#60A5FA;font-weight:600;font-size:.7rem;background:rgba(59,130,246,.1);padding:.12rem .45rem;border-radius:5px;border:1px solid rgba(59,130,246,.2);}
.action-chip{display:inline-block;padding:.12rem .5rem;border-radius:5px;font-size:.68rem;font-weight:600;background:rgba(99,102,241,.12);color:#818CF8;border:1px solid rgba(99,102,241,.2);}
.chip-login{background:rgba(16,185,129,.1);color:#10B981;border-color:rgba(16,185,129,.2);}
.chip-logout{background:rgba(107,114,128,.1);color:#6B7280;border-color:rgba(107,114,128,.18);}
.chip-open{background:rgba(59,130,246,.1);color:#60A5FA;border-color:rgba(59,130,246,.2);}
.no-data{text-align:center;padding:2rem;color:#334155;}
.online-dot{width:7px;height:7px;border-radius:50%;background:#10B981;display:inline-block;box-shadow:0 0 5px #10B981;}
.offline-dot{width:7px;height:7px;border-radius:50%;background:#475569;display:inline-block;}
.analytics-grid{display:grid;grid-template-columns:1fr 1fr;gap:.85rem;margin-bottom:.85rem;}
@media(max-width:700px){.analytics-grid{grid-template-columns:1fr;}}
.mini-bar-row{display:flex;align-items:center;gap:.5rem;margin-bottom:.35rem;font-size:.72rem;}
.mini-bar-label{width:100px;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;color:#C8D8F8;flex-shrink:0;}
.mini-bar-track{flex:1;height:5px;background:rgba(99,102,241,.12);border-radius:3px;overflow:hidden;}
.mini-bar-fill{height:100%;background:linear-gradient(90deg,#4338CA,#6366F1);border-radius:3px;transition:width .6s;}
.mini-bar-val{width:28px;text-align:right;color:#818CF8;font-weight:600;flex-shrink:0;}
.spinner{animation:spin 1s linear infinite;display:inline-block;}
@keyframes spin{to{transform:rotate(360deg);}}




/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="top-bar">
  <span class="top-mini-brand">APSG · Admin</span>
  <div class="top-sep"></div>
  <span class="top-page-label">⚙ Admin Panel</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>
<div class="page-content">
  <div class="admin-hero">
    <div class="admin-hero-icon">⚙</div>
    <div>
      <div class="admin-hero-title">Admin Panel</div>
      <div class="admin-hero-sub">Manage users, view activity logs, and analyse system usage.</div>
    </div>
  </div>
  <div class="tabs">
    <button class="tab active" onclick="showTab('analytics')">📊 Analytics Dashboard</button>
    <button class="tab" onclick="showTab('users')">👥 Users</button>
    <button class="tab" onclick="showTab('activity')">📋 Activity Log</button>
  </div>
  <!-- ── Analytics Panel ── -->
  <div class="panel active" id="tab-analytics">
    <div class="stats-grid" id="analyticsStats">
      <div class="stat-card"><div class="stat-num spinner">⟳</div><div class="stat-lbl">Loading…</div></div>
    </div>
    <div class="analytics-grid">
      <div class="card">
        <div class="card-header"><div class="card-title">👤 User Login Activity</div></div>
        <div id="loginBarChart" class="mini-bar-wrap"><div class="no-data"><span class="spinner">⟳</span></div></div>
      </div>
      <div class="card">
        <div class="card-header"><div class="card-title">🛠 Module Usage</div></div>
        <div id="moduleBarChart" class="mini-bar-wrap"><div class="no-data"><span class="spinner">⟳</span></div></div>
      </div>
    </div>
    <div class="card">
      <div class="card-header">
        <div class="card-title">👥 User Activity Summary</div>
        <div style="display:flex;gap:.5rem">
          <button class="btn-action btn-refresh" onclick="loadAnalytics()">↺ Refresh</button>
          <a href="/api/admin/analytics/download" class="btn-action btn-dl">⬇ Download Excel</a>
        </div>
      </div>
      <div class="tbl-wrap">
        <table>
          <thead><tr><th>#</th><th>Username</th><th>Full Name</th><th>Role</th><th>Total Actions</th><th>First Seen</th><th>Last Active</th><th>Status</th></tr></thead>
          <tbody id="analyticsTbl"><tr><td colspan="8" class="no-data"><span class="spinner">⟳</span> Loading…</td></tr></tbody>
        </table>
      </div>
    </div>
  </div>
  <!-- ── Users Panel ── -->
  <div class="panel" id="tab-users">
    <div class="card">
      <div class="card-header">
        <div class="card-title">Registered Users</div>
        <button class="btn-action btn-refresh" onclick="loadUsers()">↺ Refresh</button>
      </div>
      <div class="tbl-wrap">
        <table>
          <thead><tr><th>#</th><th>Username</th><th>Password</th><th>Full Name</th><th>Role</th><th>Registered</th><th>Actions</th></tr></thead>
          <tbody id="usersTbl"><tr><td colspan="5" class="no-data"><span class="spinner">⟳</span></td></tr></tbody>
        </table>
      </div>
    </div>
  </div>
  <!-- ── Activity Panel ── -->
  <div class="panel" id="tab-activity">
    <div class="card">
      <div class="card-header">
        <div class="card-title">Activity Log (latest 200)</div>
        <div style="display:flex;gap:.4rem;margin-left:auto;">
          <button class="btn-action btn-refresh" onclick="loadActivity('all')" id="actBtnAll" style="background:rgba(99,102,241,.25)">All</button>
          <button class="btn-action btn-refresh" onclick="loadActivity('login')" id="actBtnLogin">🔐 Logins Only</button>
        </div>
        <div style="display:flex;gap:.5rem">
          <button class="btn-action btn-refresh" onclick="loadActivity()">↺ Refresh</button>
          <a href="/api/admin/activity/download" class="btn-action btn-dl">⬇ Download</a>
        </div>
      </div>
      <div class="tbl-wrap">
        <table>
          <thead><tr><th>Timestamp</th><th>Username</th><th>Action</th><th>Detail</th><th>IP</th></tr></thead>
          <tbody id="activityTbl"><tr><td colspan="5" class="no-data">Click tab to load</td></tr></tbody>
        </table>
      </div>
    </div>
  </div>
</div>
<script>
function showTab(t){
  document.querySelectorAll('.tab').forEach((el,i)=>{
    el.classList.toggle('active',['analytics','users','activity'][i]===t);
  });
  document.querySelectorAll('.panel').forEach(p=>p.classList.remove('active'));
  document.getElementById(`tab-${t}`).classList.add('active');
  if(t==='analytics')loadAnalytics();
  else if(t==='users')loadUsers();
  else loadActivity();
}

async function loadAnalytics(){
  try{
    const res=await fetch('/api/admin/analytics');
    const d=await res.json();
    // Stats cards
    const now=new Date();
    const recentCutoff=new Date(now-24*60*60*1000);
    document.getElementById('analyticsStats').innerHTML=`
      <div class="stat-card"><div class="stat-num">${d.total_users}</div><div class="stat-lbl">Total Users</div></div>
      <div class="stat-card"><div class="stat-num green">${d.total_logins}</div><div class="stat-lbl">Total Logins</div></div>
      <div class="stat-card"><div class="stat-num amber">${d.total_actions}</div><div class="stat-lbl">Total Actions</div></div>
      <div class="stat-card"><div class="stat-num">${d.users.filter(u=>u.role==='admin').length}</div><div class="stat-lbl">Admin Users</div></div>
    `;
    // Login bar chart
    if(d.logins&&d.logins.length>0){
      const maxL=Math.max(...d.logins.map(l=>l.cnt));
      document.getElementById('loginBarChart').innerHTML=d.logins.slice(0,10).map(l=>`
        <div class="mini-bar-row">
          <div class="mini-bar-label" title="${esc(l.username)}">${esc(l.username)}</div>
          <div class="mini-bar-track"><div class="mini-bar-fill" style="width:${Math.round(l.cnt/maxL*100)}%"></div></div>
          <div class="mini-bar-val">${l.cnt}</div>
        </div>`).join('');
    }else{document.getElementById('loginBarChart').innerHTML='<div class="no-data">No login data yet</div>';}
    // Module usage chart
    const moduleMap={'OPEN_APP':'Open App','PPT_UPLOAD':'PPT Upload','PPT_GENERATE':'PPT Generate','EXCEL_REJECTION':'Excel Rejection','PHOTO_MERGE_START':'Photo Merge','DAILY_REPORT':'Daily Report','PPT_EXCEL':'PPT Excel','PPT_ZIP':'PPT ZIP'};
    const modCounts={};
    // Count from activity summary
    for(const u of d.activity_summary){
      if(u.actions_used){u.actions_used.split(',').forEach(a=>{
        const key=moduleMap[a.trim()]||a.trim();
        if(a.trim()!=='LOGIN'&&a.trim()!=='LOGOUT'&&a.trim()!=='REGISTER')
          modCounts[key]=(modCounts[key]||0)+1;
      });}
    }
    const modEntries=Object.entries(modCounts).sort((a,b)=>b[1]-a[1]).slice(0,8);
    if(modEntries.length>0){
      const maxM=Math.max(...modEntries.map(e=>e[1]));
      document.getElementById('moduleBarChart').innerHTML=modEntries.map(([k,v])=>`
        <div class="mini-bar-row">
          <div class="mini-bar-label" title="${esc(k)}">${esc(k)}</div>
          <div class="mini-bar-track"><div class="mini-bar-fill" style="width:${Math.round(v/maxM*100)}%;background:linear-gradient(90deg,#065F38,#10B981)"></div></div>
          <div class="mini-bar-val">${v}</div>
        </div>`).join('');
    }else{document.getElementById('moduleBarChart').innerHTML='<div class="no-data">No module usage data yet</div>';}
    // Activity summary table
    const actMap={};
    for(const a of d.activity_summary){actMap[a.username]=a;}
    document.getElementById('analyticsTbl').innerHTML=d.users.length?d.users.map((u,i)=>{
      const act=actMap[u.username]||{};
      const lastTs=act.last_seen?new Date(act.last_seen):null;
      const isRecent=lastTs&&(new Date()-lastTs)<24*60*60*1000;
      const status=isRecent?'<span class="online-dot"></span> Active (24h)':'<span class="offline-dot"></span> Inactive';
      return`<tr>
        <td style="color:#475569">${i+1}</td>
        <td><strong>${esc(u.username)}</strong></td>
        <td>${esc(u.name)}</td>
        <td><span class="${u.role==='admin'?'role-admin':'role-user'}">${u.role}</span></td>
        <td style="color:#818CF8;font-weight:600">${act.total_actions||0}</td>
        <td style="font-size:.72rem;color:#475569">${act.first_seen?act.first_seen.slice(0,16):'—'}</td>
        <td style="font-size:.72rem;color:#475569">${act.last_seen?act.last_seen.slice(0,16):'—'}</td>
        <td style="font-size:.72rem">${status}</td>
      </tr>`;
    }).join(''):'<tr><td colspan="8" class="no-data">No users found</td></tr>';
  }catch(e){console.error(e);}
}

async function loadUsers(){
  const tb=document.getElementById('usersTbl');
  tb.innerHTML='<tr><td colspan="6" class="no-data"><span class="spinner">⟳</span></td></tr>';
  try{
    const res=await fetch('/api/admin/users');const d=await res.json();
    tb.innerHTML=d.length?d.map((u,i)=>`<tr>
      <td style="color:#475569;font-size:.72rem">${i+1}</td>
      <td><strong>${esc(u.username)}</strong></td>
      <td>
        <span style="display:flex;align-items:center;gap:.4rem;">
          <span id="pw-${esc(u.username)}" style="font-family:monospace;font-size:.8rem;color:#C8D8F8;letter-spacing:.08em;">••••••••</span>
          <button onclick="togglePwView('${esc(u.username)}')" title="Show/Hide password"
            style="background:none;border:none;cursor:pointer;font-size:.85rem;padding:.1rem .2rem;color:#818CF8;">
            <span id="pw-eye-${esc(u.username)}">👁</span>
          </button>
        </span>
      </td>
      <td>${esc(u.name)}</td>
      <td><span class="${u.role==='admin'?'role-admin':'role-user'}">${u.role}</span></td>
      <td style="font-size:.7rem;color:#475569">${(u.created_at||'—').slice(0,16)}</td>
      <td>
        <div style="display:flex;gap:.35rem;flex-wrap:wrap">
          <button class="btn-action btn-pw" onclick="openChangePw('${esc(u.username)}','${esc(u.name)}')" title="Change Password">🔑 Change PW</button>
          <button class="btn-action btn-reset" onclick="resetPw('${esc(u.username)}')" title="Generate temporary password">🎲 Reset PW</button>
          ${u.username!=='admin'?`<button class="btn-action btn-del" onclick="deleteUser('${esc(u.username)}','${esc(u.name)}')" title="Delete user">🗑 Delete</button>`:''}
        </div>
      </td>
    </tr>`).join(''):'<tr><td colspan="6" class="no-data">No users found</td></tr>';
  }catch(e){tb.innerHTML=`<tr><td colspan="6" class="no-data">Error: ${e.message}</td></tr>`;}
}

// ── Change Password Modal ──────────────────────────────────────────────────────
const _pwCache={};  // username → plaintext (session cache)

async function togglePwView(username){
  const el=document.getElementById(`pw-${username}`);
  const eyeEl=document.getElementById(`pw-eye-${username}`);
  if(!el)return;
  if(el.dataset.shown==='1'){
    el.textContent='••••••••'; el.dataset.shown='0';
    eyeEl.textContent='👁';
    return;
  }
  // Already fetched?
  if(_pwCache[username]){
    el.textContent=_pwCache[username]; el.dataset.shown='1';
    eyeEl.textContent='🙈';
    return;
  }
  el.textContent='⏳';
  try{
    const res=await fetch(`/api/admin/user/${encodeURIComponent(username)}/view_password`);
    const d=await res.json();
    if(d.ok){
      _pwCache[username]=d.password;
      el.textContent=d.password; el.dataset.shown='1';
      eyeEl.textContent='🙈';
      if(d.source==='generated'){
        showModalResult(`⚠️ A new temp password was generated for @${username}: <strong>${d.password}</strong><br>Share it with the user.`);
      }
    }else{
      el.textContent='Error'; el.dataset.shown='0';
    }
  }catch(e){el.textContent='Error';}
}

function openChangePw(username, name){
  document.getElementById('pwModalUser').textContent=`${name} (@${username})`;
  document.getElementById('pwModalUsername').value=username;
  document.getElementById('pwModalNewPw').value='';
  document.getElementById('pwModalConfirm').value='';
  document.getElementById('pwModalResult').innerHTML='';
  document.getElementById('pwShowToggle').checked=false;
  document.getElementById('pwModalNewPw').type='password';
  document.getElementById('pwModalConfirm').type='password';
  document.getElementById('pwModal').style.display='flex';
}
function closePwModal(){document.getElementById('pwModal').style.display='none';}

async function submitChangePw(){
  const username=document.getElementById('pwModalUsername').value;
  const newPw=document.getElementById('pwModalNewPw').value.trim();
  const confirm=document.getElementById('pwModalConfirm').value.trim();
  const resultEl=document.getElementById('pwModalResult');
  if(!newPw||newPw.length<6){showModalResult('Password must be at least 6 characters.','error');return;}
  if(newPw!==confirm){showModalResult('Passwords do not match.','error');return;}
  try{
    const res=await fetch(`/api/admin/user/${encodeURIComponent(username)}/password`,{
      method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({password:newPw})});
    const d=await res.json();
    if(d.ok){showModalResult(`✅ Password updated successfully for @${username}`,'ok');setTimeout(closePwModal,1800);}
    else{showModalResult(d.error||'Failed','error');}
  }catch(e){showModalResult(e.message,'error');}
}

// ── Reset Password (generates temp password) ──────────────────────────────────
async function resetPw(username){
  if(!confirm(`Generate a new temporary password for @${username}?`))return;
  try{
    const res=await fetch(`/api/admin/user/${encodeURIComponent(username)}/view_password`);
    const d=await res.json();
    if(d.ok){
      document.getElementById('tempPwUser').textContent=`${d.name} (@${d.username})`;
      document.getElementById('tempPwValue').textContent=d.temp_password;
      document.getElementById('tempPwModal').style.display='flex';
    }else{alert(d.error||'Failed');}
  }catch(e){alert(e.message);}
}
function closeTempPwModal(){document.getElementById('tempPwModal').style.display='none';}
function copyTempPw(){
  const txt=document.getElementById('tempPwValue').textContent;
  navigator.clipboard.writeText(txt).then(()=>{
    const btn=document.getElementById('copyTempBtn');
    btn.textContent='✅ Copied!';setTimeout(()=>btn.textContent='📋 Copy',1500);
  });
}

// ── Delete User ────────────────────────────────────────────────────────────────
async function deleteUser(username,name){
  if(!confirm(`Delete user "${name}" (@${username})?\nThis action cannot be undone.`))return;
  try{
    const res=await fetch(`/api/admin/user/${encodeURIComponent(username)}/delete`,{method:'DELETE'});
    const d=await res.json();
    if(d.ok){loadUsers();loadAnalytics();}
    else{alert(d.error||'Delete failed');}
  }catch(e){alert(e.message);}
}

function showModalResult(msg,type){
  const el=document.getElementById('pwModalResult');
  el.textContent=msg;
  el.className='modal-result '+(type==='ok'?'result-ok':'result-err');
}
function togglePwShow(checked){
  const t=checked?'text':'password';
  document.getElementById('pwModalNewPw').type=t;
  document.getElementById('pwModalConfirm').type=t;
}

async function loadActivity(filter='all'){
  document.getElementById('actBtnAll').style.background=filter==='all'?'rgba(99,102,241,.25)':''; 
  const loginBtn=document.getElementById('actBtnLogin'); if(loginBtn) loginBtn.style.background=filter==='login'?'rgba(16,185,129,.25)':'';
  const tb=document.getElementById('activityTbl');
  tb.innerHTML='<tr><td colspan="5" class="no-data"><span class="spinner">⟳</span></td></tr>';
  try{
    const actUrl=filter==='login'?'/api/admin/activity?limit=200&action=LOGIN':'/api/admin/activity?limit=200';
    const res=await fetch(actUrl);const d=await res.json();
    const chipClass=(a)=>a==='LOGIN'?'chip-login':a==='LOGOUT'?'chip-logout':a.startsWith('OPEN')?'chip-open':'';
    tb.innerHTML=d.length?d.map(r=>`<tr>
      <td style="font-size:.72rem;color:#475569;white-space:nowrap">${r.ts||''}</td>
      <td><strong>${esc(r.username)}</strong></td>
      <td><span class="action-chip ${chipClass(r.action)}">${esc(r.action)}</span></td>
      <td style="font-size:.75rem;color:#6B7280">${esc(r.detail||'')}</td>
      <td style="font-size:.72rem;color:#475569">${r.ip||''}</td>
    </tr>`).join(''):'<tr><td colspan="5" class="no-data">No activity yet</td></tr>';
  }catch(e){tb.innerHTML=`<tr><td colspan="5" class="no-data">Error: ${e.message}</td></tr>`;}
}

function esc(s){return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
loadAnalytics();
</script>

<!-- ── Change Password Modal ── -->
<div class="modal-overlay" id="pwModal" onclick="if(event.target===this)closePwModal()">
  <div class="modal-card">
    <div class="modal-title">🔑 Change Password</div>
    <div class="modal-sub" id="pwModalUser"></div>
    <input type="hidden" id="pwModalUsername">
    <div class="modal-field">
      <label>New Password</label>
      <input type="password" id="pwModalNewPw" placeholder="Enter new password (min 6 chars)">
    </div>
    <div class="modal-field">
      <label>Confirm Password</label>
      <input type="password" id="pwModalConfirm" placeholder="Re-enter new password">
    </div>
    <label class="show-pw-row">
      <input type="checkbox" id="pwShowToggle" onchange="togglePwShow(this.checked)">
      Show passwords
    </label>
    <div class="modal-result" id="pwModalResult"></div>
    <div class="modal-actions">
      <button class="modal-btn modal-btn-cancel" onclick="closePwModal()">✕ Cancel</button>
      <button class="modal-btn modal-btn-primary" onclick="submitChangePw()">✔ Update Password</button>
    </div>
  </div>
</div>

<!-- ── Temp Password Display Modal ── -->
<div class="modal-overlay" id="tempPwModal" onclick="if(event.target===this)closeTempPwModal()">
  <div class="modal-card">
    <div class="modal-title">🎲 Temporary Password Generated</div>
    <div class="modal-sub" id="tempPwUser"></div>
    <div class="temp-pw-box">
      <div class="temp-pw-label">Temporary Password</div>
      <div class="temp-pw-value" id="tempPwValue"></div>
    </div>
    <div class="temp-pw-warn">
      ⚠️ Share this password securely with the user. It will be required to log in.
      The user should change it immediately after logging in.
    </div>
    <div class="modal-actions" style="margin-top:1rem">
      <button class="modal-btn modal-btn-cancel" onclick="closeTempPwModal()">✕ Close</button>
      <button class="modal-btn modal-btn-primary" id="copyTempBtn" onclick="copyTempPw()">📋 Copy Password</button>
    </div>
  </div>
</div>

<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""


# ═══════════════════════════════════════════════════════════════════════════════
#  STARTUP
# ═══════════════════════════════════════════════════════════════════════════════

# Initialize DB (must be after all function definitions)
try:
    init_db()
    print("✓ Database initialized")
except Exception as _init_e:
    print(f"DB init warning: {_init_e}")

if __name__ == "__main__":
    port=int(os.environ.get("PORT",5000))
    app.run(host="0.0.0.0",port=port,debug=False)
